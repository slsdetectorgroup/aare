"""Build docs/cf_cuda_fused.pptx — the algorithm + kernel + hardware half of
docs/cf_cuda_kernel.pptx fused with the opt1→opt7 optimization story.

The ladder is told in three acts, ordered by which bar is tallest:

    ACT I   [f64]  feed the GPU        opt1 opt2 opt3 opt4   (+ route A, rejected)
    ACT II  [f64]  get results back    opt5 opt6             (+ routes B', B'', rejected)
    ACT III [f32]  the kernel          opt7

Act III comes last because it cannot be justified earlier: measured through
collect() the f32 kernel is worth 1.5 % end-to-end, and only once the host is off
the critical path is the same change worth 21 %.

cf_cuda_kernel.pptx is kept solely as the base presentation — it donates the PSI
theme and title slide, and every other slide of it is deleted below.

All hardware numbers are re-measured against the CURRENT kernel, not taken
from the kernel deck (whose implementation details and timings are stale):

    nvcc -arch=sm_89 --ptxas-options=-v          -> registers, spills
    cudaOccupancyMaxActiveBlocksPerMultiprocessor -> blocks/SM, occupancy

    3x3 : 34 regs/thread, 0 spill, 1296 B smem, 6 blocks/SM, 100.0% occupancy
    9x9 : 128 regs/thread, 0 spill, 2304 B smem, 2 blocks/SM,  33.3% occupancy
    32x32 blocks @ 9x9: 0 blocks/SM -- 1024 x 128 regs > 65536 regs/SM

Performance numbers come from docs/ClusterFinderCUDA_benchmark_results.md (quotable
rows only).
"""
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pathlib import Path
from PIL import Image

DOCS = Path(__file__).resolve().parent.parent
FIGS = DOCS / "figures"
BASE = DOCS / "cf_cuda_kernel.pptx"          # PSI theme + title slide
OUT = DOCS / "cf_cuda_fused.pptx"

# ---------------------------------------------------------------- design tokens
BG     = RGBColor(0x0B, 0x10, 0x18)
PANEL  = RGBColor(0x12, 0x1A, 0x28)
CODEBG = RGBColor(0x0E, 0x14, 0x20)
RULE   = RGBColor(0x1E, 0x28, 0x36)
ACCENT = RGBColor(0x1E, 0x90, 0xC2)
AMBER  = RGBColor(0xE8, 0xB2, 0x5C)
PALE   = RGBColor(0xE7, 0xED, 0xF4)
TEXT2  = RGBColor(0xA5, 0xB2, 0xC4)
MUTED  = RGBColor(0x6B, 0x7A, 0x90)
CARD   = RGBColor(0xF4, 0xF6, 0xF9)          # light card for white figures

UI, MONO = "Segoe UI", "Consolas"
W, H = 13.333, 7.5
M = 0.7                       # left margin
COL = 7.9                     # left column width
RAIL_X, RAIL_W = 9.2, 3.5     # right rail

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

prs = Presentation(str(BASE))
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[0]                  # 'Blank Slide' — zero shapes
N_SLIDES = 34


# ------------------------------------------------------------------- helpers
def keep_only_slide(prs, keep=0):
    """Drop every slide but one from the base presentation."""
    lst = prs.slides._sldIdLst
    for i, sldId in enumerate(list(lst)):
        if i != keep:
            prs.part.drop_rel(sldId.get(f"{R}id"))
            lst.remove(sldId)


def set_para_texts(shape, texts):
    """Replace paragraph texts in-place, keeping each paragraph's formatting."""
    for p, txt in zip(shape.text_frame.paragraphs, texts):
        if not p.runs:
            continue
        p.runs[0].text = txt
        for r in p.runs[1:]:
            r.text = ""


def new_slide():
    """A dark slide on the PSI master.

    Two independent guards, because the base template's master carries a PSI
    background picture and logo that must not bleed through:
      1. showMasterSp="0" + a slide-level <p:bg> (correct schema position:
         first child of <p:cSld>), which is what PowerPoint honours;
      2. a full-bleed rectangle as the first shape, which every renderer
         honours regardless of how it treats (1).
    """
    s = prs.slides.add_slide(BLANK)
    s._element.set("showMasterSp", "0")

    cSld = s._element.find(f"{P}cSld")
    bg = etree.Element(f"{P}bg")
    pr = etree.SubElement(bg, f"{P}bgPr")
    fill = etree.SubElement(pr, f"{A}solidFill")
    clr = etree.SubElement(fill, f"{A}srgbClr")
    clr.set("val", "0B1018")
    etree.SubElement(pr, f"{A}effectLst")
    cSld.insert(0, bg)

    rect(s, 0, 0, W, H, BG)
    return s


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, first=False, space_after=0, space_before=0, line=None, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    if line: p.line_spacing = line
    if align: p.alignment = align
    return p


def run(p, text, size=11, color=TEXT2, font=UI, bold=False, italic=False, spc=None):
    r = p.add_run(); r.text = text
    f = r.font
    f.name, f.size, f.bold, f.italic = font, Pt(size), bold, italic
    f.color.rgb = color
    if spc is not None:
        r.font._rPr.set("spc", str(int(spc * 100)))
    return r


# ------------------------------------------------------------------ chrome
def chrome(s, idx, eyebrow, title, title_size=27):
    rect(s, M, 0.60, 0.35, 0.035, ACCENT)
    tf = tb(s, 1.17, 0.50, 10.33, 0.32)
    run(para(tf, True), eyebrow.upper(), 9, MUTED, bold=True, spc=1.6)

    tf = tb(s, M, 0.86, 11.9, 1.0)
    run(para(tf, True, line=1.05), title, title_size, PALE, bold=True)

    span, n = 11.0, N_SLIDES
    pitch = span / n; wseg = pitch * 0.90
    for i in range(n):
        rect(s, M + i * pitch, 7.28, wseg, 0.045, ACCENT if i <= idx - 1 else RULE)
    tf = tb(s, 12.0, 7.14, 0.9, 0.3)
    run(para(tf, True, align=PP_ALIGN.RIGHT), f"{idx} / {n}", 8.5, MUTED)


N_ANNEX = 5          # annex GROUPS, not slides


def annex_chrome(s, grp, eyebrow, title, part=None, nparts=None, title_size=27):
    """Same chrome, amber, on its own progress track.

    The annex is GROUPED, not enumerated: slides that belong together carry the
    same tag (A2 is both CUDA-Graphs slides, A5 all three artefacts), because
    what a reader needs at a glance is which slides form one argument. The badge
    is rendered in the header, not only in the corner.
    """
    tf = tb(s, M, 0.44, 1.0, 0.34)
    run(para(tf, True), f"A{grp}", 15, AMBER, bold=True)
    tf = tb(s, 1.42, 0.50, 10.1, 0.32)
    tag = f"ANNEX · {eyebrow}" + (f" · {part} of {nparts}" if nparts else "")
    run(para(tf, True), tag.upper(), 9, MUTED, bold=True, spc=1.6)
    tf = tb(s, M, 0.86, 11.9, 1.0)
    run(para(tf, True, line=1.05), title, title_size, PALE, bold=True)
    pitch = 11.0 / N_ANNEX
    for i in range(N_ANNEX):
        rect(s, M + i * pitch, 7.28, pitch * 0.90, 0.045,
             AMBER if i <= grp - 1 else RULE)
    tf = tb(s, 11.7, 7.14, 1.2, 0.3)
    foot = f"A{grp}" + (f" · {part}/{nparts}" if nparts else "")
    run(para(tf, True, align=PP_ALIGN.RIGHT), foot, 8.5, MUTED)


def table(s, x, y, w, header, rows, colw, size=9.5, rowh=0.62):
    """Minimal header + zebra table. colw are fractions of w."""
    xs, acc = [], 0.0
    for c in colw:
        xs.append(x + acc * w)
        acc += c
    rect(s, x, y, w, 0.34, PANEL)
    for cx, h in zip(xs, header):
        tf = tb(s, cx + 0.16, y + 0.08, w, 0.24)
        run(para(tf, True), _up(h), 8, MUTED, bold=True, spc=1.2)
    yy = y + 0.38
    for i, row in enumerate(rows):
        if i % 2 == 0:
            rect(s, x, yy, w, rowh, PANEL)
        for j, (cx, cell) in enumerate(zip(xs, row)):
            col = PALE if j == 0 else TEXT2
            tf = tb(s, cx + 0.16, yy + 0.10, colw[j] * w - 0.24, rowh)
            p = para(tf, True, line=1.2)
            for k, part in enumerate(str(cell).split("**")):
                if part:
                    run(p, part, size, AMBER if k % 2 else col, bold=bool(k % 2))
        yy += rowh + 0.04
    return yy


def bullets(s, x, y, w, items, size=11, gap=7):
    tf = tb(s, x, y, w, 0.3)
    for i, it in enumerate(items):
        color, txt = (it if isinstance(it, tuple) else (TEXT2, it))
        p = para(tf, i == 0, space_after=gap, line=1.25)
        run(p, "•  ", size, MUTED)
        for j, part in enumerate(txt.split("**")):
            if part:
                run(p, part, size, PALE if j % 2 else color, bold=bool(j % 2))
    return tf


def code(s, x, y, w, lines, size=8.5, title=None):
    lh = 0.148
    h = 0.24 + len(lines) * lh + (0.22 if title else 0)
    rect(s, x, y, w, h, CODEBG, MSO_SHAPE.ROUNDED_RECTANGLE)
    ty = y + 0.12
    if title:
        tf = tb(s, x + 0.18, ty, w - 0.36, 0.2)
        run(para(tf, True), title, 7.5, MUTED, bold=True, spc=1.2)
        ty += 0.22
    tf = tb(s, x + 0.18, ty, w - 0.36, h - 0.24)
    for i, ln in enumerate(lines):
        p = para(tf, i == 0, line=1.12)
        if ln.strip().startswith(("//", "#")):
            run(p, ln, size, MUTED, MONO)
            continue
        for j, part in enumerate(ln.split("«")):
            for k, seg in enumerate(part.split("»")):
                if not seg: continue
                hi = (j > 0 and k == 0)
                run(p, seg, size, ACCENT if hi else TEXT2, MONO, bold=hi)
    return h


def callout(s, x, y, w, text, h=0.78, color=ACCENT, size=10.5):
    rect(s, x + 0.045, y, w - 0.045, h, PANEL)
    rect(s, x, y, 0.045, h, color)
    tf = tb(s, x + 0.28, y + 0.10, w - 0.5, h - 0.2, MSO_ANCHOR.MIDDLE)
    p = para(tf, True, line=1.2)
    for j, part in enumerate(text.split("**")):
        if part:
            run(p, part, size, PALE if j % 2 else TEXT2, bold=bool(j % 2))


def _up(txt):
    """upper() for labels, but 'µ'.upper() is Greek capital Mu — which renders as
    an 'M' and turns 'µs' into 'MS', i.e. microseconds into milliseconds."""
    return txt.upper().replace("\u039c", "µ")


def rail(s, items, y0=2.0, divider=True):
    if divider:
        rect(s, 8.95, 2.0, 0.012, 4.55, RULE)
    y = y0
    for it in items:
        kind = it[0]
        if kind == "label":
            tf = tb(s, RAIL_X, y, RAIL_W, 0.26)
            run(para(tf, True), _up(it[1]), 8.5, MUTED, bold=True, spc=1.4)
            y += 0.28
        elif kind == "stat":
            _, lab, val, col = it
            tf = tb(s, RAIL_X, y, RAIL_W, 0.24)
            run(para(tf, True), _up(lab), 8.5, MUTED, spc=1.2)
            tf = tb(s, RAIL_X, y + 0.24, RAIL_W, 0.6)
            run(para(tf, True), val, 26, col, bold=True)
            y += 0.98
        elif kind == "row":
            _, lab, val, col = it
            tf = tb(s, RAIL_X, y, RAIL_W, 0.24)
            run(para(tf, True), _up(lab), 8.5, MUTED, spc=1.2)
            tf = tb(s, RAIL_X, y + 0.22, RAIL_W, 0.3)
            run(para(tf, True), val, 13, col, bold=True)
            y += 0.66
        elif kind == "note":
            tf = tb(s, RAIL_X, y, RAIL_W, 0.9)
            run(para(tf, True, line=1.25), it[1], 9, TEXT2)
            y += 0.30 + 0.17 * (len(it[1]) // 42 + 1)
        elif kind == "gap":
            y += it[1]
    return y


def figure(s, name, x, y, w):
    p = FIGS / f"{name}.png"
    iw, ih = Image.open(p).size
    h = w * ih / iw
    s.shapes.add_picture(str(p), In(x), In(y), In(w), In(h))
    return h


def card_figure(s, name, x, y, w, pad=0.10):
    """A light-background figure (imported, not re-rendered) on a light card."""
    p = FIGS / f"{name}.png"
    iw, ih = Image.open(p).size
    h = w * ih / iw
    rect(s, x - pad, y - pad, w + 2 * pad, h + 2 * pad, CARD,
         MSO_SHAPE.ROUNDED_RECTANGLE)
    s.shapes.add_picture(str(p), In(x), In(y), In(w), In(h))
    return h + 2 * pad


def notes(s, text):
    """Speaker notes. Detail that belongs in the talk, not on the slide.

    python-pptx creates the notes slide on first access, so this is safe to call
    on any slide. Used to relieve slides that carry a figure: the mechanism goes
    on the screen, the API detail goes here.
    """
    s.notes_slide.notes_text_frame.text = text


def caption(s, x, y, w, text, size=9):
    tf = tb(s, x, y, w, 0.3)
    run(para(tf, True, line=1.25), text, size, MUTED)


def flow(s, x, y, w, steps, h=0.62):
    """Numbered left-to-right step strip."""
    n = len(steps); gap = 0.30
    bw = (w - gap * (n - 1)) / n
    for i, t in enumerate(steps):
        bx = x + i * (bw + gap)
        rect(s, bx, y, bw, h, PANEL)
        rect(s, bx, y, 0.03, h, ACCENT)
        tf = tb(s, bx + 0.20, y + 0.05, bw - 0.32, h - 0.10, MSO_ANCHOR.MIDDLE)
        p = para(tf, True, line=1.1)
        run(p, f"{i + 1}  ", 9, ACCENT, bold=True, font=MONO)
        run(p, t, 9.5, PALE)
        if i < n - 1:
            tf = tb(s, bx + bw, y + 0.05, gap, h - 0.10, MSO_ANCHOR.MIDDLE)
            run(para(tf, True, align=PP_ALIGN.CENTER), "›", 15, MUTED, bold=True)


def statstrip(s, x, y, w, items, h=0.80):
    n = len(items); gap = 0.22
    bw = (w - gap * (n - 1)) / n
    for i, (lab, val) in enumerate(items):
        bx = x + i * (bw + gap)
        rect(s, bx, y, bw, h, PANEL)
        tf = tb(s, bx + 0.20, y + 0.12, bw - 0.4, 0.22)
        run(para(tf, True), lab.upper(), 8, MUTED, bold=True, spc=1.2)
        tf = tb(s, bx + 0.20, y + 0.37, bw - 0.4, 0.34)
        run(para(tf, True), val, 15, PALE, bold=True)


# ------------------------------------------------------------ interstitial
_NARROW, _WIDE = set("ijlt.,;:'!|()[]"), set("mwMW")


def _em(txt):
    """Width of `txt` in ems of Segoe UI Bold, near enough to wrap by.

    A character count is not good enough: 'now the tallest bar' and
    'measured, and in what' are the same length and differ by 8 % in width, which
    is exactly the margin that decides whether a title takes two lines or three.
    Weights are calibrated against a LibreOffice render of two title lines and
    reproduce both to within 1 %.
    """
    return sum(1.08 if c in _WIDE else
               0.37 if c in _NARROW else
               0.34 if c == " " else 0.68 for c in txt)


def _fit(title, w, sizes=(36, 32, 28, 24), lines=2):
    """Largest size at which `title` wraps into at most `lines` lines across `w`.

    python-pptx cannot measure text and PowerPoint's autofit does not apply until
    a render, so a title that grows one line silently walks over the rule beneath
    it. No tolerance is granted — a title that only just fits is one font
    substitution away from not fitting on someone else's machine.
    """
    for size in sizes:
        cap = w / (size / 72)
        n, cur = 1, 0.0
        for word in title.split():
            need = _em(word) + (0.34 if cur else 0)
            if cur + need > cap and cur:
                n, cur = n + 1, _em(word)
            else:
                cur += need
        if n <= lines:
            return size
    return sizes[-1]


def section(kicker, title, thesis, items, rng, col=ACCENT, carry=None,
            annex=False):
    """An unnumbered beat between sections: where we are, and what is coming.

    Deliberately sparse — it exists to buy 10–15 s of stage setting, so it has
    to be readable at a glance and finished before the audience starts reading
    ahead. It carries no slide number and takes no tick of its own on the
    progress track: slides 3–34 keep the numbers they have, so the annex's
    cross-references ("expands slide 26") stay true. What it lights up instead
    is the *range* the section covers, which is the thing the audience wants.
    """
    s = new_slide()
    rect(s, 0, 0, 0.16, H, col)

    tf = tb(s, M + 0.3, 1.52, 5.8, 0.3)
    run(para(tf, True), kicker.upper(), 9.5, MUTED, bold=True, spc=1.8)
    tf = tb(s, M + 0.3, 1.90, 6.0, 1.35)
    run(para(tf, True, line=1.03), title, _fit(title, 6.0), PALE, bold=True)
    rect(s, M + 0.3, 3.42, 1.5, 0.03, col)
    tf = tb(s, M + 0.3, 3.70, 5.7, 1.6)
    run(para(tf, True, line=1.4), thesis, 13.5, TEXT2)

    if carry:
        lab, val, sub = carry
        rect(s, M + 0.3, 5.55, 5.7, 1.12, PANEL)
        rect(s, M + 0.3, 5.55, 0.035, 1.12, col)
        tf = tb(s, M + 0.60, 5.72, 5.2, 0.24)
        run(para(tf, True), _up(lab), 8.5, MUTED, bold=True, spc=1.4)
        tf = tb(s, M + 0.60, 5.94, 5.2, 0.4)
        run(para(tf, True), val, 24, col, bold=True)
        tf = tb(s, M + 0.60, 6.40, 5.2, 0.24)
        run(para(tf, True), sub, 9.5, MUTED)

    rect(s, 7.15, 1.95, 0.012, 4.4, RULE)
    step = 0.42 if len(items) > 7 else 0.46 if len(items) > 5 else 0.54
    y = 1.95 + (4.4 - len(items) * step) / 2
    tf = tb(s, 7.45, y - 0.42, 5.0, 0.26)
    run(para(tf, True), "COMING UP", 8.5, MUTED, bold=True, spc=1.6)
    for num, txt in items:
        tf = tb(s, 7.45, y, 0.7, 0.3)
        run(para(tf, True), str(num), 12, col, bold=True, font=MONO)
        tf = tb(s, 8.15, y, 4.5, 0.3)
        p = para(tf, True)
        for j, part in enumerate(txt.split("**")):
            if part:
                run(p, part, 12, col if j % 2 else PALE, bold=bool(j % 2))
        y += step

    # The annex divider sits on the annex's own track: the main arc is finished
    # behind it, so lighting main-track segments would misreport where we are.
    n_track = N_ANNEX if annex else N_SLIDES
    pitch = 11.0 / n_track
    for i in range(n_track):
        n = i + 1
        # the section ahead in its own colour, what is already behind us dimmed,
        # the rest dark — so the divider agrees with the chrome on either side.
        c = col if annex or rng[0] <= n <= rng[1] else (
            MUTED if n < rng[0] else RULE)
        rect(s, M + i * pitch, 7.28, pitch * 0.90, 0.045, c)
    return s


# =========================================================== 1 · PSI TITLE
keep_only_slide(prs, 0)
title_slide = prs.slides[0]
by_name = {sh.name: sh for sh in title_slide.shapes}
set_para_texts(by_name["TextShape 1"], ["The CUDA ClusterFinder"])
set_para_texts(by_name["CustomShape 4"],
               ["Kernel design · hardware limits · seven optimization steps"])
set_para_texts(by_name["CustomShape 2"], ["Khalil Daniel Ferjaoui"])
set_para_texts(by_name["CustomShape 3"],
               ["Paul Scherrer Institut · aare", "August 2026"])

# =========================================================== 2 · HERO
s = new_slide()
rect(s, 0, 0, 0.16, H, ACCENT)
tf = tb(s, M + 0.3, 0.85, 11, 0.3)
run(para(tf, True), "AARE · HYBRID PIXEL DETECTORS · CUDA CLUSTERFINDER",
    9.5, MUTED, bold=True, spc=1.8)

tf = tb(s, M + 0.3, 1.28, 11.4, 1.45)
run(para(tf, True, line=1.02),
    "The kernel was never the bottleneck — feeding it was", 42, PALE, bold=True)
tf = tb(s, M + 0.3, 2.80, 11.4, 0.45)
run(para(tf, True, line=1.05),
    "One kernel, one thread per pixel, and seven steps to keep it fed", 22, ACCENT)

tf = tb(s, M + 0.3, 3.40, 11.4, 0.9)
run(para(tf, True, line=1.3),
    "The stencil was fast almost immediately: at 3×3 the kernel needs 5.5 µs per "
    "frame while getting that frame across PCIe costs 16.6 — 13.2 even "
    "uncontended. The frame is gated by the wire, not the arithmetic. Six of the "
    "seven steps get data in, get results back, and measure honestly.",
    12.5, TEXT2)

tf = tb(s, M + 0.3, 4.40, 11.4, 0.26)
run(para(tf, True),
    "WHERE THIS ENDS UP — THE SHIPPED f32 BUILD AFTER ALL SEVEN STEPS · "
    "ENGINE TIMES [f32 · s4] · RECONCILED IN A1",
    8.5, AMBER, bold=True, spc=1.4)

stats = [("×9.1", "VS 24-THREAD CPU", ACCENT), ("61,312", "FRAMES / SECOND", PALE),
         ("16.3 µs", "PER FRAME, END TO END", PALE), ("6 / 23 M", "CLUSTER MISMATCH VS CPU TWIN", AMBER)]
for i, (v, l, c) in enumerate(stats):
    x = M + 0.3 + i * 2.85
    rect(s, x, 4.78, 0.035, 0.95, c)
    tf = tb(s, x + 0.22, 4.78, 2.5, 0.55)
    run(para(tf, True), v, 30, c, bold=True)
    tf = tb(s, x + 0.22, 5.41, 2.5, 0.3)
    run(para(tf, True), l, 8.5, MUTED, spc=1.2)

rect(s, M + 0.3, 6.05, 11.0, 0.012, RULE)
tf = tb(s, M + 0.3, 6.25, 11.4, 0.6)
run(para(tf, True, line=1.35),
    "RTX 4090 (Ada, sm_89) · PCIe 4.0 ×16 · Mönch 400×400 uint16 · 3×3 clusters · "
    "100 000 frames · Cu fluorescence, MAX IV", 10, MUTED)

# ------------------------------------------------------- divider · context
section("Context · what the code does",
        "What the kernel does, and what limits it",
        "No optimizations yet, only what the hardware has to do.",
        [("3–4", "The algorithm"),
         ("5–6", "The two machines"),
         ("7–8", "The CUDA kernel"),
         ("9–10", "The hardware limit")],
        rng=(3, 10))

# =========================================================== 3 · THE PHYSICS
s = new_slide()
chrome(s, 3, "The algorithm · what it is for",
       "A photon is not a pixel — it is a cluster")
bullets(s, M, 1.90, 7.9, [
    "Charge from one absorbed photon **spreads over neighbouring pixels**. "
    "Summing that 3×3 patch recovers the photon energy; a single pixel does not.",
    "The histogram of those cluster energies **is** the measurement: peak position "
    "and width give the detector's gain and **energy resolution**.",
], size=10.5)
figure(s, "fig_frame", M, 2.98, 7.45)
rail(s, [("label", "Why it matters")], y0=1.95)
h = card_figure(s, "img_spectra", RAIL_X, 2.32, RAIL_W)
caption(s, RAIL_X, 2.32 + h + 0.10, RAIL_W,
        "Cluster-energy spectra from an energy scan, against allpix² simulation. "
        "Each peak is one beam energy; its width is the resolution being measured.",
        size=8.5)
rail(s, [
    ("row", "Pixels per frame", "400 × 400 = 160 000", TEXT2),
    ("row", "Peak pixels = photons / frame", "~2 330   ·   1.5 %", ACCENT),
    ("row", "Pixels above 5σ (2.4 per photon)", "~5 700   ·   3.6 %", ACCENT),
], y0=5.10, divider=False)
caption(s, M, 6.68, 7.9,
        "Real MOENCH data, Cu fluorescence, MAX IV beamtime. One cluster is emitted "
        "per local maximum, so 2 330 counts photons, not lit pixels; the recorded "
        "3×3 windows cover 12.7 % of the frame.")

# =========================================================== 4 · PER FRAME
s = new_slide()
chrome(s, 4, "The algorithm · per frame", "Per pixel: subtract, threshold, update the pedestal")
bullets(s, M, 1.95, COL, [
    "Per pixel: subtract a **running pedestal** (mean ± rms), keep pixels above "
    "**nσ · rms**, cut a 3×3 cluster around each local maximum.",
    "400×400 = 160 k pixels, **312.5 kB per frame**; Cu data yields ~2 330 clusters "
    "per frame at 3×3.",
    "The pedestal is **updated by every pixel that sees no photon** — about 80 % of "
    "them — every frame, so the arithmetic and the data movement are coupled.",
])
code(s, M, 3.55, COL, [
    "// the whole algorithm, per pixel",
    "v   = frame[i] - pedestal_mean[i]",
    "rms = pedestal rms at i",
    "m   = max(v) over the 3x3 window at i",
    "if (m > «nSigma» * rms)                  // a photon is within reach",
    "    if (v == m) -> emit cluster        //   ... and I am its peak",
    "    else        -> «nothing»             //   ... I am in its shadow",
    "else            -> update pedestal     // I saw nothing",
], title="THE WHOLE ALGORITHM · THREE OUTCOMES, NOT TWO")
callout(s, M, 5.55, COL,
        "**Thesis of this talk:** the compute was fast almost immediately. "
        "Six of the seven steps are about feeding it.")
rail(s, [
    ("label", "The shape of the work"),
    ("gap", 0.15),
    ("stat", "Work items per frame", "160 000", PALE),
    ("row", "Operations on each", "~5, identical", TEXT2),
    ("gap", 0.12),
    ("row", "Communication between them", "none", ACCENT),
    ("row", "Order they may run in", "any", ACCENT),
    ("gap", 0.22),
    ("note", "What a pixel does never depends on what its neighbours decided, only "
             "on what they measured. That one property is what the next two slides "
             "point two very different machines at."),
])

# ==================================================== 5 · THE CPU
# The two machine slides. They exist because the audience is asked, from slide 7
# on, to accept "one thread per pixel" and "occupancy" without ever having been
# shown what a thread costs on either machine. Both diagrams are the CS149 ones
# (credited in the captions): redrawing them would lose the shared visual
# grammar — orange fetch/decode, yellow ALU, blue execution context — which is
# the entire reason the pair reads at a glance.
s = new_slide()
chrome(s, 5, "The machine we are starting from",
       "CPU: latency-oriented, built to finish one thread fast")
figure(s, "img_cpu_core", M, 1.98, 7.3)
callout(s, M, 5.58, COL,
        "Count the boxes: **6 fetch/decode**, out-of-order instruction selection, two "
        "levels of private cache, all of it to keep **two** instruction streams fed. "
        "The ALUs are the small part.", h=0.86, size=10)
caption(s, M, 6.54, COL,
        "One core, schematically. Intel Skylake is shown. The Zen 4 core in this machine "
        "differs in detail (4 FP pipes rather than 3; AVX-512 double-pumped on 256-bit "
        "datapaths) but not in kind: ~6-wide front end, 4 scalar ALUs, 2 SMT contexts, "
        "private L1 + L2. Load/store units not drawn. After Stanford CS149, Fall 2025.",
        size=8)
rail(s, [
    ("label", "pc-moench-04 · AMD Ryzen 9 7950X"),
    ("gap", 0.10),
    ("stat", "Cores × SMT", "16 × 2", PALE),
    ("gap", 0.04),
    ("stat", "CPU, 1 thread", "1.75 ms", MUTED),
    ("stat", "CPU MT, 24 threads", "148 µs", PALE),
    ("gap", 0.06),
    ("row", "That is the bar", "6 762 frames / s", ACCENT),
    ("gap", 0.16),
    ("note", "Thread count is swept per cluster size, not assumed: "
             "24 at 3×3, 32 at 9×9."),
])
notes(s, """The CPU slide. The point is not that CPUs are bad -- it is what the
silicon is SPENT on.

Six fetch/decode units, an out-of-order instruction selector, branch prediction
(not even drawn), L1 + L2 private cache: all of that machinery exists to find
independent work INSIDE one instruction stream, and to hide memory latency
behind a cache. It is the right design when you have a few threads that must
each finish fast.

Our problem has the opposite shape: 160 000 work items that are already
independent. We do not need a machine to FIND the parallelism -- it is handed to
us. Every transistor spent looking for it is a transistor not doing arithmetic.

On the thread sweep: 6 762 FPS is the best of a measured sweep (perf/cpu_threads.py),
not a default. 48 threads on 16 cores is 24 % SLOWER than 24 -- worth saying out
loud, because an oversubscribed baseline is the easiest way to inflate a GPU
speedup without lying about anything.""")

# ==================================================== 6 · THE GPU
s = new_slide()
chrome(s, 6, "The machine we are moving to",
       "GPU: throughput-oriented, the whole frame at once")
figure(s, "img_gpu_die", M, 2.10, 2.75)
caption(s, M, 4.92, 2.75,
        "AD102 · 144 blocks, 128 enabled on this card. One SM boxed.", size=7.5)
figure(s, "img_gpu_sm", M + 3.05, 1.92, 4.75)
# The colour key is a separate crop: in the source it spans the full slide width
# while the diagram spans 60 % of it, so one rectangle cannot hold both.
figure(s, "img_gpu_legend", M + 3.05, 5.26, 4.75)
callout(s, M, 6.02, COL,
        "Same grammar, inverted proportions: **4 fetch/decode** for **64 warp "
        "contexts** and a wall of lanes. Nothing reorders instructions: when a warp "
        "stalls on memory, the selector just **runs a different one**.", h=0.86,
        size=10)
caption(s, M, 6.90, COL,
        "One SM: a V100 is shown; this card's is the same idea (128 FP32 lanes, 48 warp "
        "slots, 100 kB shared memory). Note how few of the units are FP64, and remember "
        "it at opt7. Diagrams after Stanford CS149, Fall 2025.", size=8)
rail(s, [
    ("label", "NVIDIA GeForce RTX 4090"),
    ("gap", 0.12),
    ("stat", "FP32 lanes", "16 384", PALE),
    ("row", "128 SMs × 128 lanes", "vs 32 CPU streams", TEXT2),
    ("gap", 0.14),
    ("stat", "Resident thread slots", "196 608", ACCENT),
    ("row", "The frame needs", "160 000  →  all at once", AMBER),
    ("gap", 0.20),
    ("note", "Device memory runs at 1 008 GB/s (384-bit, 21 Gbps); PCIe delivers 25. "
             "Two copy engines, so H2D and D2H run at the same time."),
])
notes(s, """The GPU slide, and the number to land.

128 SMs x 1 536 threads = 196 608 thread slots that can be RESIDENT at the same
time. Our frame is 160 000 pixels. The entire frame fits in the machine at once,
one thread per pixel, with room left over -- which is why slide 7's "one thread
per pixel" is not a figure of speech, and why slide 9 can talk about occupancy
as a real quantity (3x3 achieves 100 %: 1 536 threads resident per SM).

Latency hiding, in one sentence: the CPU hides memory latency with a cache and
out-of-order execution; the GPU hides it by having 48 other warps ready to run.
That is why there is no reorder buffer on this diagram and no branch predictor.

The trap to pre-empt: 16 384 lanes vs 32 streams is a factor of 512, and we
measure x9.1. Say so before someone else does. We are not lane-limited; we are
BANDWIDTH-limited -- 1 TB/s of device memory, and 25 GB/s of PCIe to reach it.
That gap between 512 and 9 IS the talk.""")

# =========================================================== 7 · THE KERNEL
s = new_slide()
chrome(s, 7, "The CUDA kernel · execution model", "One thread per pixel")
flow(s, M, 1.90, 11.9, ["load tile + halo", "__syncthreads", "stencil reduction",
                        "classify", "append or update pedestal"])
bullets(s, M, 2.85, 7.5, [
    "A **16×16 block = 256 threads** covers 256 pixels; the grid tiles the "
    "whole 400×400 frame. Cluster geometry is a **template parameter**, so the "
    "stencil is fully unrolled at compile time.",
    "The output is **sparse**: only detections touch global memory, through one "
    "atomic bump of a per-frame counter. The **decision work is dense**: every "
    "pixel is tested, independently and identically.",
    "That is exactly the shape a GPU wants: regular, independent, repeated "
    "160 000 times per frame.",
], size=10.5)
code(s, M, 4.35, 7.5, [
    "block = dim3(BLOCK_X, BLOCK_Y);            // 16 x 16",
    "grid  = dim3((ncols + BLOCK_X - 1)/BLOCK_X,",
    "             (nrows + BLOCK_Y - 1)/BLOCK_Y);",
    "device::find_clusters_in_single_frame<<<grid, block, «shmem», stream>>>(",
    "    d_frame, d_pd_mean, d_pd_sum, d_pd_sum2, d_pd_off, n_pd_samples,",
    "    nSigma, nrows, ncols, d_clusters, d_cluster_count, max_clusters);",
], size=8, title="LAUNCH CONFIGURATION · ClusterFinderCUDA.hpp")
caption(s, M, 5.80, 7.5,
        "400×400 pixels → a 25×25 grid of 16×16 blocks = 625 blocks per frame, "
        "handed to 128 SMs. Nothing about the launch depends on how many clusters "
        "the frame happens to contain, which is what makes the work uniform.")
code(s, 8.5, 2.85, 4.1, [
    "// ClusterFinder.hpp — the serial CPU",
    "if (max > nSigma * rms) {",
    "    if (value < max)",
    "        continue;  // Not max go to the",
    "                   // next pixel — but",
    "                   // also no pedestal",
    "                   // update",
    "} else {",
    "    pedestal.«push_fast»(iy, ix, ...);",
    "}",
], size=8, title="THREE OUTCOMES, NOT TWO")
callout(s, 8.5, 4.80, 4.1,
        "That comment is **verbatim from the CPU source**. A pixel in a photon's "
        "shadow is **neither recorded nor fed back** — and the CUDA kernel "
        "reproduces it exactly.",
        h=1.15, size=10)
callout(s, 8.5, 6.10, 4.1,
        "**~80 % of threads update the pedestal**: every pixel with no photon in "
        "its window. ~1.5 % are peaks, ~18 % are shadow. The update, not the "
        "cluster write, dominates.",
        h=1.1, size=10, color=AMBER)

# =========================================================== 6 · TILING
s = new_slide()
chrome(s, 8, "The CUDA kernel · shared memory", "Load the tile once, reuse it nine times")
bullets(s, M, 1.92, 12.0, [
    "Neighbouring threads need **overlapping** 3×3 windows. Without shared memory "
    "each pixel would be fetched from global memory up to nine times.",
    "Each block stages a tile of (16 + 2r) × (16 + 2r) **pedestal-subtracted** "
    "values, the halo is the price of the stencil, and it is loaded cooperatively "
    "by the threads on the block edges.",
], size=10.5)
figure(s, "fig_tile", M, 3.05, 7.5)
code(s, 8.4, 3.05, 4.3, [
    "extern __shared__ unsigned char smem[];",
    "auto *sh = (COMPUTE_TYPE*)smem;",
    "auto stride = blockDim.x + 2*col_radius;",
    "auto tid = (threadIdx.y + row_radius)*stride",
    "         + (threadIdx.x + col_radius);",
    "// pedestal subtraction fused into the load",
    "sh[tid] = d_frame[gid] - «d_pd_mean»[gid];",
], size=8, title="clusterfinder_kernel.cuh")
callout(s, 8.4, 4.60, 4.3,
        "Only **odd** cluster sizes are supported (3×3, 5×5, 7×7, 9×9), so that "
        "the centre pixel is unique and local-maximum suppression is well defined.",
        h=1.15, size=10)
caption(s, 8.4, 5.95, 4.3,
        "The tile is stored in COMPUTE_TYPE (float), not in the pedestal type: "
        "1.3 KB for 3×3, 2.3 KB for 9×9, against 100 KB of shared memory per SM. "
        "Even the old double-precision tile only reached 4.5 KB.")
caption(s, M, 6.70, 7.5,
        "Halo cost falls with block size: 56 % of the tile at 8×8, 27 % at 16×16, "
        "13 % at 32×32, which is the first half of the block-size argument. "
        "The second half is registers, next slide.")

# =========================================================== 7 · OCCUPANCY
s = new_slide()
chrome(s, 9, "Hardware · occupancy", "3×3 reaches 100 % occupancy, 9×9 only 33 %")
statstrip(s, M, 1.82, 11.9, [
    ("block size", "16 × 16"),
    ("threads / block", "256"),
    ("max threads / SM · 48 warps", "1 536"),
    ("registers / SM", "65 536"),
    ("shared mem / block", "1.3 – 2.3 KB"),
])
figure(s, "fig_occupancy", 0.9, 2.78, 11.5)
callout(s, M, 6.12, 11.9,
        "**16×16 is the balance point**, enough threads to amortise the halo, few "
        "enough that 6 blocks still fit per SM. At 3×3 that is 48 of 48 warps "
        "resident: **fully occupied**. Shared memory never binds; **registers do**.",
        h=0.66, size=10.5)
caption(s, M, 6.84, 11.9,
        "Shipped f32 build, measured not estimated: cuobjdump -res-usage for registers "
        "and spills, cudaOccupancyMaxActiveBlocksPerMultiprocessor for blocks/SM. "
        "RTX 4090 (sm_89). Reproduce: python/tests/perf/kernel_resources.py")

# =========================================================== 8 · REGISTERS
s = new_slide()
chrome(s, 10, "Hardware · the real limiter", "Occupancy is a means, not the goal")
bullets(s, M, 1.92, COL, [
    "The limiter is **register pressure**: every thread keeps a private "
    "clusterData[CSX × CSY] staging array, so demand grows with the **square** of "
    "the cluster size. 3×3 costs **38 registers**, 9×9 costs **128**.",
    "**Zero spills either way.** ptxas would rather give up occupancy than spill to "
    "local memory, and on this kernel that is the right trade.",
    "Low occupancy is not automatically a problem: at 9×9 one kernel nearly fills "
    "the machine on its own.",
], size=10)
figure(s, "fig_regpressure", M, 3.24, 7.9)
callout(s, M, 5.00, COL,
        "**Occupancy is an output.** At 9×9 the register file is full at two blocks, "
        "so two thirds of the thread slots are stranded: nothing else can be resident "
        "no matter how much work is queued.", h=0.80, size=10)
code(s, M, 5.92, COL, [
    "«cuobjdump -res-usage» build/aare/_aare_cuda*.so | c++filt",
    "  3x3: REG:«38»  STACK:0 LOCAL:0     # STACK/LOCAL 0 = no spills",
    "  9x9: REG:«128» STACK:0 LOCAL:0",
    "python python/tests/perf/«kernel_resources.py»   # blocks/SM, runtime check",
], size=7.5, title="HOW THESE NUMBERS WERE OBTAINED · READ FROM THE BUILT .SO, NO REBUILD")
rail(s, [
    ("label", "Register budget · sm_89 · f32"),
    ("gap", 0.15),
    ("stat", "3×3 · occupancy", "100 %", ACCENT),
    ("stat", "9×9 · occupancy", "33 %", AMBER),
    ("gap", 0.05),
    ("row", "Spills, either case", "0 bytes", TEXT2),
    ("row", "3×3 on the f64 build", "47 regs → 83 %", AMBER),
    ("gap", 0.20),
    ("note", "Which is why the 9×9 kernel is the one worth optimising, see opt7."),
])
notes(s, """The arithmetic behind the bars, and the one caveat.

One SM holds 65 536 registers and 1 536 thread slots. A 16x16 block is 256
threads, so a block costs regs_per_thread x 256 registers. At 3x3 that is 9 728,
six blocks fit, and the THREAD SLOTS run out first at 100 %. At 9x9 it is
32 768, so two blocks exactly fill the REGISTER FILE and strand two thirds of
the slots. Same kernel, same block size, opposite binding resource.

The caveat on "low occupancy is fine": it is build-dependent. On the shipped f32
build the 9x9 kernel is long enough that four streams overlap by only 1.02x, so
the extra streams buy nothing in kernel concurrency. On the f64 build the kernel
is 3.5x longer and DOES overlap, 1.32x, and there four streams genuinely lower
the floor. Quote the overlap factor, never "occupancy is fine" on its own.

Build dependence of the register count itself: the f64 pedestal costs 3x3 nine
extra registers (47), which loses a block per SM and drops it to 83 %. 9x9 is
unmoved at 128 -- there the limiter is the clusterData[9][9] staging array, not
the pedestal, which is why opt7 helps 3x3's occupancy and not 9x9's.""")

# =========================================================== 10 · THE LADDER
s = new_slide()
chrome(s, 11, "Roadmap",
       "Three acts, ordered by which bar is tallest")
rows = [
    ("act", "ACT I  ·  feed the GPU", "the host cannot submit work fast enough", "at 3×3  [f64]", ACCENT),
    ("opt1", "First CUDA port", "1 stream, one launch per frame", "×2.34", ACCENT),
    ("opt2", "Streams + batching", "4 streams, 2 000-frame batches", "×3.66", ACCENT),
    ("opt3", "Pipeline rework", "sync barriers removed", "×4.32", ACCENT),
    ("opt4", "Pinned memory", "DMA-speed host transfers", "×5.69", ACCENT),
    ("act", "ACT II  ·  get the results back", "the host copy is now the tallest bar", "at 3×3  [f64]", PALE),
    ("opt5", "Host↔GPU overlap", "chunked submit / collect", "×7.45", PALE),
    ("opt6", "Zero-copy collection", "read in place, never copy", "×8.65", PALE),
    ("act", "ACT III  ·  the kernel", "the kernel is the tallest bar — at 9×9", "at 9×9  [f32]", AMBER),
    ("opt7", "FP32 pedestal + variance rewrite", "the only kernel change in the deck",
     "kernel −41%", AMBER),
]
y = 1.70
for i, (tag, name, sub, gain, col) in enumerate(rows):
    if tag == "act":
        rect(s, M, y + 0.30, 11.9, 0.016, col)
        tf = tb(s, M + 0.02, y, 8.0, 0.28)
        run(para(tf, True), name.upper(), 11, col, bold=True, spc=1.4)
        tf = tb(s, M + 5.2, y + 0.03, 6.7, 0.26)
        run(para(tf, True, align=PP_ALIGN.RIGHT), f"{sub}   {gain}", 9, MUTED)
        y += 0.40
        continue
    rect(s, M, y, 11.9, 0.48, PANEL)
    rect(s, M, y, 0.035, 0.48, col)
    tf = tb(s, M + 0.28, y + 0.09, 1.0, 0.32)
    run(para(tf, True), tag, 12.5, col, bold=True, font=MONO)
    tf = tb(s, M + 1.45, y + 0.04, 5.0, 0.28)
    run(para(tf, True), name, 11.5, PALE, bold=True)
    tf = tb(s, M + 1.45, y + 0.26, 5.6, 0.26)
    run(para(tf, True), sub, 9, MUTED)
    tf = tb(s, 9.4, y + 0.09, 3.1, 0.35)
    run(para(tf, True, align=PP_ALIGN.RIGHT), gain, 12.5, col, bold=True)
    y += 0.52
caption(s, M, 6.62, 11.9,
        "Speedups are 3×3 vs the best CPU configuration, 24 threads. Two routes were measured and "
        "rejected, two result transports inside Act II, shown next to the step they "
        "lost to, and CUDA Graphs, which is in the annex. The rule that predicts "
        "the wins predicts the failures too.")

# --------------------------------------------------------- divider · ACT I
section("Act I of III · feed the GPU",
        "The host cannot submit work fast enough",
        "The kernel was fast almost immediately. This act is about the host, and it "
        "is told at 3×3, where the wire is the floor.",
        [(12, "opt1 · first port"),
         (13, "opt2 · streams + batching"),
         (14, "opt3 · no barriers"),
         (15, "opt4 · pinned memory")],
        rng=(12, 15),
        carry=("Starting from", "6 762 FPS",
               "24-thread CPU · 14.8 s for 100 000 frames"))

# =========================================================== 11 · OPT1
s = new_slide()
chrome(s, 12, "Act I · opt1 · the first CUDA port",
       "The first port runs at 26 % of the GPU's floor")
bullets(s, M, 1.95, COL, [
    "Shared-memory tiling with **halo loading** for any cluster size; pedestal "
    "subtraction fused into the tile load.",
    "Cluster geometry is a **compile-time template parameter** → the 3×3 stencil "
    "is fully unrolled.",
    "One cudaMemcpy in, one kernel, one cudaMemcpy out; **the host blocks "
    "on every frame**.",
])
code(s, M, 3.18, COL, [
    "// one frame at a time, the host waits at every step",
    "cudaMemcpy(d_frame, h_frame, bytes, cudaMemcpyHostToDevice);",
    "find_clusters_in_single_frame<ClusterType, FRAME_TYPE>",
    "    <<<grid, block, shmem>>>(d_frame, d_pd_mean, ...);",
    "cudaMemcpy(h_out, d_out, out_bytes, cudaMemcpyDeviceToHost);",
], title="ClusterFinderCUDAOpt2.hpp · find_clusters()")
callout(s, M, 4.50, COL,
        "**PCIe is full-duplex**: H2D, kernel and D2H run on independent engines and "
        "overlap, so the floor is **max(H2D, kernel, D2H)**, never the sum. Each term "
        "is that engine's **busy time per frame at 4 streams**, the union of its "
        "intervals, not how long one operation takes. At 3×3 **[f64 · s4]**: "
        "max(**16.17**, 15.17, 7.69) = **16.2 µs → 61 859 FPS**. That max is a "
        "**profiled estimate**, so the deck's floor is the **lower of it and the best "
        "sustained rate** — the same 16.17 here; **A1** shows where the two part at 9×9.",
        h=1.30, size=10.5)
figure(s, "fig_opt1_timeline", M, 5.85, COL)
rail(s, [
    ("label", "opt1 · 3×3 · 100 k frames · f64"),
    ("gap", 0.10),
    ("stat", "Throughput", "15 807 FPS", PALE),
    ("stat", "vs 24-thread CPU", "×2.34", ACCENT),
    ("gap", 0.05),
    ("row", "Per frame", "63.3 µs", TEXT2),
    ("row", "The GPU floor", "61 859 FPS", ACCENT),
    ("row", "% of floor", "26 %", AMBER),
])

# =========================================================== 12 · OPT2
s = new_slide()
chrome(s, 13, "Act I · opt2 · streams and batching",
       "Four streams and 2 000-frame batches: ×1.56")
bullets(s, M, 1.95, COL, [
    "A **stream** is an ordered queue of GPU work. Work in **different** streams may "
    "overlap, so a copy can run while another stream computes.",
    "Each stream gets its own **StreamContext**: device frame buffer, output buffer "
    "and pedestal. Frames are handed out **round-robin**.",
    "The host now submits **2 000 frames per call** instead of one.",
])
code(s, M, 3.30, COL, [
    "struct StreamContext {",
    "    cudaStream_t   stream;",
    "    FRAME_TYPE    *d_frame;      ClusterType *d_clusters;",
    "    PEDESTAL_TYPE *d_pd_mean, *d_pd_sum, *d_pd_sum2;",
    "};",
    "auto &sc = v_sc[frame_idx % «n_streams»];   // round-robin",
], title="ClusterFinderCUDA.hpp · per-stream state")
figure(s, "fig_opt2_timeline", M, 4.77, COL)
callout(s, M, 6.48, COL,
        "**Scaffolding, not yet the payoff.** The streams exist, but the host still "
        "synchronises after every round: see opt3.", h=0.60)
rail(s, [
    ("label", "opt2 · 3×3 · 4 streams · batch 2 000"),
    ("gap", 0.10),
    ("stat", "Throughput", "24 726 FPS", PALE),
    ("stat", "vs 24-thread CPU", "×3.66", ACCENT),
    ("gap", 0.05),
    ("row", "Per frame", "40.4 µs", TEXT2),
    ("row", "Step gain over opt1", "×1.56", ACCENT),
    ("row", "% of floor · 61 859 FPS", "40 %  (was 26)", AMBER),
])

# =========================================================== 13 · OPT3
s = new_slide()
chrome(s, 14, "Act I · opt3 · remove the sync barriers",
       "One sync per batch, not one per round: ×1.18")
bullets(s, M, 1.95, 7.4, [
    "opt2 synchronised **all streams after every round** of n_streams frames. "
    "The GPU drained to empty each time.",
    "opt3 submits every frame's H2D → kernel → D2H **asynchronously**, then "
    "synchronises **once at the end of the batch**.",
], size=10.5)
figure(s, "fig_streams", M, 3.05, 6.55)
code(s, 8.35, 1.95, 4.25, [
    "// opt2:  barrier after every round",
    "for (round) {",
    "   submit(n_streams frames);",
    "   «cudaDeviceSynchronize»();",
    "}",
    "",
    "// opt3:  submit everything, sync once",
    "for (frame : batch) {",
    "   cudaMemcpyAsync(..., sc.stream);",
    "   kernel<<<..., sc.stream>>>(...);",
    "   cudaMemcpyAsync(..., sc.stream);",
    "}",
    "for (sc : streams)",
    "   «cudaStreamSynchronize»(sc.stream);",
], size=8, title="THE ONE-LINE IDEA")
callout(s, 8.35, 5.05, 4.25,
        "**29 188 FPS · ×4.32**\n34.3 µs/frame · 47 % of floor (was 40)", h=0.86, size=11)
caption(s, 8.35, 6.15, 4.25,
        "Each lane is one stream. Removing the barrier lets a stream start its next "
        "frame while its neighbours are still copying. The three panels are scheduled, "
        "not sketched: H2D and D2H are one FIFO engine each, so a stream waits for the "
        "copy engine, never for another stream's copy to finish overlapping it.")

# =========================================================== 14 · OPT4
s = new_slide()
chrome(s, 15, "Act I · opt4 · pinned (page-locked) memory",
       "Pinning the input buys DMA-speed H2D: ×1.32")
bullets(s, M, 1.95, 12.0, [
    "Normal host memory is **pageable**: the OS may move or swap it. A DMA engine "
    "cannot safely read that, so the driver first copies your data into a **hidden "
    "pinned staging buffer**. Every transfer is copied twice.",
    "**Pinning** locks the pages in physical RAM. The GPU's DMA engine then reads "
    "host memory **directly**, no staging copy, and the transfer can be truly asynchronous.",
], size=10.5)
figure(s, "fig_pinning", M, 3.15, 7.6)
code(s, 8.5, 3.15, 4.1, [
    "// pin the whole dataset once",
    "«cudaHostRegister»(ptr, bytes,",
    "                 cudaHostRegisterDefault);",
    "",
    "// ... run the whole campaign ...",
    "",
    "«cudaHostUnregister»(ptr);",
], size=8, title="ClusterFinderCUDA.hpp")
callout(s, 8.5, 4.72, 4.1,
        "**38 486 FPS · ×5.69**\n26.0 µs/frame · 62 % of floor, the largest step in Act I",
        h=0.86, size=11)
caption(s, 8.5, 5.82, 4.1,
        "Measured H2D [s1, uncontended]: one 400×400 uint16 frame (312.5 KiB = "
        "320 000 B) in 13.2 µs = 24.2 GB/s, 77 % of PCIe 4.0 ×16 theoretical, i.e. "
        "true DMA speed. In the shipped pipeline it reads 16.6 [s4], +26 % of "
        "H2D↔D2H contention (A1). Pageable staging runs ~15 GB/s.")
callout(s, M, 6.45, 7.6,
        "**The rule, first sighting: ×1.32 at 3×3 but only ×1.02 at 9×9.** Pinning "
        "attacks H2D, the tallest bar at 3×3, the shortest at 9×9.",
        h=0.72, size=10, color=AMBER)

# -------------------------------------------------------- divider · ACT II
section("Act II of III · get the results back",
        "The host copy is now the tallest bar",
        "Frames go in at DMA speed. The results still come back slowly. Still 3×3 — "
        "but 9×9 is where this act pays most.",
        [(16, "opt5 · host↔GPU overlap"),
         (17, "opt6 · zero-copy"),
         (18, "two rejected routes")],
        rng=(16, 18), col=PALE,
        carry=("Arriving at", "38 486 FPS",
               "opt4 · 26.0 µs per frame · 62 % of the GPU floor"))

# =========================================================== 15 · OPT5
s = new_slide()
chrome(s, 16, "Act II · opt5 · host↔GPU overlap",
       "Overlapping host and GPU hides min(host, GPU): ×1.31")
bullets(s, M, 1.92, COL, [
    "opt3 overlapped H2D ∥ kernel ∥ D2H **across streams, inside one batch**, but "
    "never the **host** with the GPU: find_clusters_batched synchronised, then built "
    "thousands of ClusterVectors with the GPU idle.",
    "opt5 keeps **one batch in flight while materialising the previous one**: chunk "
    "i+1 is submitted before chunk i is collected.",
])
figure(s, "fig_overlap", M - 0.15, 3.28, COL + 0.30)
callout(s, M, 5.86, COL,
        "The time hidden is **min(GPU, host)**, so this pays most when the two are "
        "comparable, and cannot rescue a host term that is simply larger.")
notes(s, """opt5 — host<->GPU overlap. Code and chunk sizing are on annex A3.

    tok = cf.submit_batch(data[a0:b0], first_frame=a0)
    for a, b in bounds[1:]:
        nxt = cf.submit_batch(data[a:b], first_frame=a)   # GPU starts N+1 ...
        results.extend(cf.collect(tok))                   # ... host unpacks N
        tok = nxt
    results.extend(cf.collect(tok))                       # drain

1. This is now INTERNAL to find_clusters_batched(), so every caller gets it for
   free. submit_batch/collect stay public for anyone who wants the token by hand.

2. The chunk size is rounded to a multiple of n_streams, because the device
   pedestal is per-stream: an uneven chunk would advance the four stream
   pedestals by different amounts and the finders would stop being comparable.

Why the gain differs by cluster size: the saving is min(GPU, host) per chunk, so
it is largest when the two terms are comparable. At 3x3 they nearly are (GPU
16.2 us, host ~9.8) and opt5 is worth x1.31. At 9x9 the host term is roughly
twice the GPU term, so overlap hides only the smaller one and opt5 is worth
x1.20 - which is exactly the diagnosis that motivates opt6: you cannot overlap
your way out of a host term that is simply larger. Report SS8.2.""")
rail(s, [
    ("label", "opt5 · 3×3 and 9×9 · no CUDA work at all"),
    ("gap", 0.10),
    ("stat", "3×3 throughput", "50 410 FPS", PALE),
    ("row", "per frame · step · of floor", "19.8 µs · ×1.31 · 81 %", ACCENT),
    ("gap", 0.14),
    ("stat", "9×9 throughput", "15 063 FPS", PALE),
    ("row", "per frame · step · of floor", "66.4 µs · ×1.20 · 45 %", AMBER),
    ("gap", 0.15),
    ("note", "For clusters that must outlive the finder, this is the endpoint at "
             "3×3: opt6 lends, it does not give."),
])

# =========================================================== 19 · OPT6
s = new_slide()
chrome(s, 17, "Act II · opt6 · zero-copy collection",
       "Read the results in place: ×2.21 at 9×9")
bullets(s, M, 1.92, COL, [
    "The D2H lands in a **pinned host buffer**. collect() then allocates one "
    "ClusterVector per frame and memcpys into it; at 9×9 that is **467 kB per frame, "
    "~9.3 GB per run**, single-threaded.",
    "collect_view() returns a **BatchView**: strided numpy views straight onto the "
    "pinned buffer. It withholds ownership past the chunk, **not access**, every "
    "cluster's payload and coordinates are readable.",
], size=10.5)
figure(s, "fig_resultpath", M - 0.15, 3.05, COL + 0.30)
callout(s, M, 5.86, COL,
        "The win is **max(0, host copy − GPU floor)**: at 3×3 the 8 µs copy hides under "
        "a 16.2 µs floor and opt5 had already absorbed most of it; at 9×9 the 40 µs copy "
        "is **larger than the floor** and cannot hide at any overlap.", h=0.80, size=10)
caption(s, M, 6.78, COL,
        "The two bars are the competing costs, not the two steps; the step times are "
        "on the right.", size=8.5)
rail(s, [
    ("label", "opt6 · 3×3 and 9×9 · collect_view()"),
    ("gap", 0.10),
    ("stat", "3×3 throughput", "58 495 FPS", PALE),
    ("row", "per frame · step · of floor", "17.1 µs · ×1.16 · 95 %", ACCENT),
    ("gap", 0.14),
    ("stat", "9×9 throughput", "33 323 FPS", PALE),
    ("row", "per frame · step · of floor", "30.0 µs · ×2.21 · 100 %", AMBER),
    ("gap", 0.12),
    ("note", "opt5 → opt6: 19.8 → 17.1 and 66.4 → 30.0 µs, bit-identical, "
             "0.2 % spread, zero warm faults. FLOOR = lower of the s4 engine max "
             "and the best sustained rate: at 3×3 the 16.17 µs max sets it; at 9×9 "
             "f64 the run BEATS the 32.66 µs max, so the 30.0 sustained sets it [A1]."),
])

# =========================================================== 20 · ACT II rejected
s = new_slide()
chrome(s, 18, "Act II · two rejected routes",
       "The copy is allocation-bound, not bandwidth-bound")
rows = [
    ("B′", "One allocation per chunk", "collect_packed()",
     "Removes the per-frame malloc but keeps the copy, and the copy is ~80% of the "
     "cost. Worse, the replacement allocation is 1.17 GB, far above glibc's mmap "
     "threshold, so it is mmap'd and munmap'd every chunk: 606 566 faults, ~21 µs/frame.",
     "69.3 µs  ·  deleted from the API"),
    ("B″", "Parallel materialisation", "8-thread copy pool",
     "Each worker gets its own glibc arena, which destroys the cross-run heap reuse "
     "that makes the single-threaded path cheap. Faults went 9 700 → 2 270 000; "
     "MALLOC_ARENA_MAX=1 collapsed them back to 138 k, which is the proof.",
     "+6% at best, −33% when results are freed promptly"),
]
y = 2.00
for tag, name, how, body, verdict in rows:
    rect(s, M, y, 11.9, 2.05, PANEL)
    rect(s, M, y, 0.035, 2.05, AMBER)
    tf = tb(s, M + 0.30, y + 0.22, 1.0, 0.4)
    run(para(tf, True), tag, 20, AMBER, bold=True, font=MONO)
    tf = tb(s, M + 1.30, y + 0.20, 6.0, 0.3)
    run(para(tf, True), name, 14, PALE, bold=True)
    tf = tb(s, M + 1.30, y + 0.52, 6.0, 0.3)
    run(para(tf, True), how, 9.5, MUTED, font=MONO)
    tf = tb(s, M + 1.30, y + 0.90, 10.2, 1.0)
    run(para(tf, True, line=1.3), body, 10, TEXT2)
    tf = tb(s, M + 1.30, y + 1.68, 10.2, 0.3)
    run(para(tf, True), verdict, 10.5, AMBER, bold=True)
    y += 2.25
callout(s, M, 6.38, 11.9,
        "**Copying faster does not help when the cost is the OS populating pages.** "
        "The only winning move is not to allocate — which is exactly what opt6 does. "
        "materialize_slot() is deliberately single-threaded and carries a comment "
        "saying so, to stop the experiment being repeated.", h=0.80, size=10.5)

# ------------------------------------------------------- divider · ACT III
section("Act III of III · the kernel",
        "Only now is the kernel the tallest bar",
        "The story moves to 9×9. At 3×3 the wire was always the floor and the kernel "
        "never mattered; at 9×9 it is the tallest bar, which is why this act exists.",
        [(19, "opt7 · FP32 pedestal"),
         (20, "the correctness trap"),
         (21, "the variance rewrite"),
         (22, "why it comes last")],
        rng=(19, 22), col=AMBER,
        carry=("Arriving at", "58 495 FPS",
               "opt6 · 3×3, and 33 323 FPS at 9×9, where this act pays"))

# =========================================================== 18 · OPT7 why
s = new_slide()
chrome(s, 19, "Act III · opt7 · FP32 device pedestal",
       "FP32 halves pedestal traffic: −41 % kernel time")
bullets(s, M, 1.95, 7.5, [
    "**~80 % of pixels** take the **pedestal-update** branch: it reads off, sum, sum² "
    "and writes back sum, sum², mean. All four pedestal arrays are DEVICE_PED_TYPE, "
    "so one typedef halves every one of those six accesses — **48 bytes per updating "
    "pixel in FP64**, **24 bytes in FP32**.",
    "The kernel is **bandwidth-bound**, so halving that traffic nearly halves the time.",
    "Second effect: on a GeForce part, **FP64 arithmetic runs at 1/64 of FP32**. "
    "The pedestal update was paying that tax on every pixel.",
    "Third, quietly: the narrower accumulators free **9 registers at 3×3**, 47 → 38, "
    "which buys back a block per SM and takes occupancy **83 % → 100 %**.",
], size=10.5)
figure(s, "fig_f32_kernel", M, 4.05, 7.5)
code(s, 8.5, 1.95, 4.1, [
    "// clusterfinder_kernel.cuh",
    "using COMPUTE_TYPE    = float;",
    "using DEVICE_PED_TYPE = «float»;",
    "//            was: double",
], size=8.5, title="ONE TYPEDEF")
callout(s, 8.5, 3.05, 4.1,
        "Kernel, 9×9 **[s1 · cap 1700 · exclusive]**\n**39.86 µs → 23.70 µs  (−40.5%)**\n"
        "Shipped **[s4]**: **30.0 → 25.1 µs** end to end.",
        h=1.32, size=10.5)
caption(s, 8.5, 4.48, 4.1,
        "Both panels, both builds, same git rev, 20 000 frames. The deck's rule, and "
        "the reason there are two: quote s1 for how long an engine takes, s4 for which "
        "one sets the floor. A1 has the full grid.")
callout(s, 8.5, 5.24, 4.1,
        "Naive FP32 is **wrong** (next two slides), and even when correct it only pays "
        "**because Act II came first**.", h=0.80, size=10, color=AMBER)
callout(s, 8.5, 6.20, 4.1,
        "At 3×3 the same typedef is worth **−70.6 %** (14.72 → 4.32 µs [s1]) yet buys "
        "only **4.6 %** end to end: there the kernel was never the tallest bar.",
        h=0.88, size=10)

# =========================================================== 19 · OPT6 trap
s = new_slide()
chrome(s, 20, "Act III · opt7 · the correctness trap",
       "The naive FP32 pedestal loses the variance to cancellation")
bullets(s, M, 1.95, 7.9, [
    "The running variance was computed as **var = E[X²] − mean²**. With a pedestal "
    "mean of ~4 655 ADU, both operands are ≈ 2.17 × 10⁷ while the answer is ≈ 2 000.",
    "In FP32 each operand sits on a grid of **2 ADU²**, so the answer inherits "
    "**±3 — an absolute error that does not shrink**. For a pixel whose true "
    "variance is 9 that is a third of it; below rms ≈ 2 it goes negative and the "
    "**rms clamps to 0**.",
], size=10.5)
figure(s, "fig_cancellation", M, 3.25, 7.5)
rail(s, [
    ("label", "What it looked like"),
    ("gap", 0.15),
    ("stat", "Extra clusters", "+28.06%", AMBER),
    ("gap", 0.05),
    ("note", "Quiet pixels got rms → 0, so their threshold became 0 and they fired "
             "on every single frame, producing a large unphysical high-energy tail "
             "in the spectrum."),
    ("gap", 0.35),
    ("row", "Affected pixels", "~1–2% of the sensor", TEXT2),
    ("row", "Written up in", "docs/pedestal_precision_…", MUTED),
])

# =========================================================== 20 · OPT6 fix
s = new_slide()
chrome(s, 21, "Act III · opt7 · the variance rewrite",
       "Accumulate what is small, not what is large")
bullets(s, M, 1.95, COL, [
    "Freeze a per-pixel baseline **X₀ = round(mean)** once, at the end of pedestal "
    "training, and never move it again.",
    "Accumulate the **centred** value Y = X − X₀ instead of X, and report the mean "
    "as **X₀ + sum/n** — so nothing downstream changes.",
])
figure(s, "fig_variance_rewrite", M, 2.95, COL)
code(s, M, 4.47, COL, [
    "// before: both terms ~2.17e7, answer ~2000",
    "var = sum2/n - mean*mean;",
    "",
    "// after: centred on a frozen per-pixel offset X0",
    "DEVICE_PED_TYPE resid  = mean - «d_pd_off»[i];      // ~O(1)",
    "DEVICE_PED_TYPE var_px = sum2[i]/n - resid*resid;  // no cancellation",
], title="clusterfinder_kernel.cuh")
callout(s, M, 5.94, COL, 
        "Result: the 100% FP32 build now matches the FP64 build to "
        "**3 × 10⁻⁷**, 70 clusters out of 233 million.")
rail(s, [
    ("label", "Why it works"),
    ("gap", 0.15),
    ("note", "Precision is relative. Floats resolve small numbers finely and large "
             "numbers coarsely, so never let a small answer be the difference of "
             "two large numbers."),
    ("gap", 0.5),
    ("row", "f32 vs f64 counts", "3 × 10⁻⁷", ACCENT),
    ("row", "vs CPU", "0.0039%", ACCENT),
    ("gap", 0.3),
    ("note", "X₀ must never be updated: the accumulators are defined relative to it."),
])

# =========================================================== 21 · OPT6 when
s = new_slide()
chrome(s, 22, "Act III · why this act comes last",
       "The saving never grew — the frame around it shrank")
figure(s, "fig_f32_absolute", M, 1.95, 11.9)
callout(s, M, 4.86, 5.85,
        "**The same typedef saves −4.63 µs at opt4 and −4.87 µs at opt6** — the two "
        "readings whose fault counts match. What changes is the denominator: the frame "
        "falls 79.8 → 30.0 µs, so an identical saving reads **−5.8 %, then −16.2 %**. "
        "Act II did not make the kernel win bigger; it made it **separable**.",
        h=1.28, size=10.5)
callout(s, 6.75, 4.86, 5.85,
        "**And the act ends by handing the floor away.** At **s4**, the config that "
        "ships, the f64 arm is kernel-bound: **32.66** against a **25.25 µs** D2H. The "
        "typedef puts the kernel at **23.94** — **below a D2H that never moved**. "
        "The constraint is now the result path, not the arithmetic. [engines: 19/34]",
        h=1.28, size=10.5, color=AMBER)
caption(s, M, 6.30, 11.9,
        "9×9 · 20 000 frames · 4 streams · cap 1 700 · warm · both arms at the same git "
        "rev. opt3 is excluded: its two arms sat in different allocator states, so that "
        "comparison does not report the typedef at all. † opt5's arms differ the same "
        "way; its reading agrees with the other two but is not independently "
        "attributable. Both are worked through in annex A4. The f32 bar at opt6 is "
        "opt7, the shipped build; the f32 bars at opt4 and opt5 are the same typedef "
        "applied at earlier steps — configurations that exist only to make this "
        "comparison controlled.")

# ------------------------------------------------------- divider · results
section("Results · what came out of it",
        "The whole ladder, and how to use it",
        "Both cluster sizes end to end, and the audit behind the numbers.",
        [("23–24", "Results, both cluster sizes"),
         ("25", "Where the time went"),
         ("26–27", "What the numbers survived")],
        rng=(23, 27), col=PALE,
        carry=("Arriving at", "58 495 FPS",
               "opt6 · everything after this is the ladder seen whole"))

# =========================================================== 22 · RESULTS
s = new_slide()
chrome(s, 23, "Results · 3×3", "×9.1 at 3×3, sitting on the H2D floor")
figure(s, "fig_arc", 1.95, 1.88, 9.4)
callout(s, M, 5.80, 5.85,
        "**×9.1 over 24 CPU threads**, 14.8 s → 1.63 s for 100 000 frames, "
        "and **at the H2D floor**.", h=0.8)
callout(s, 6.75, 5.80, 5.85,
        "Every step is **monotonic**, and correctness is held constant **throughout**, "
        "0.004 % against the CPU baseline, and **exact** against the CPU twin "
        "that isolates the port (slides 28–31).", h=0.8, color=AMBER)
caption(s, M, 6.62, 11.9,
        "3×3 clusters · nσ = 5 · 100 000 frames · batch 2 000 · 4 streams · 5 reps · "
        "warm = best of reps 1–4 (collect() does not converge, it oscillates between "
        "allocator states) · each step in its own process · CPU baseline = "
        "ClusterFinderMT at its best thread count, 24 here, first pass only.")

# =========================================================== 23 · RESULTS 9x9
s = new_slide()
chrome(s, 24, "Results · 9×9", "×26.5 at 9×9, and opt7 hands the floor to D2H")
figure(s, "fig_arc_9x9", 1.95, 1.88, 9.4)
callout(s, M, 5.80, 5.85,
        "**×26.5 over 32 CPU threads**; the kernel is the tallest bar for the whole "
        "f64 arm, and opt7's −40 % drops it **below D2H**.", h=0.8)
callout(s, 6.75, 5.80, 5.85,
        "opt4 buys **×1.03** here and **×1.32** at 3×3. Same code, opposite regimes, "
        "the rule, twice.", h=0.8, color=AMBER)
caption(s, M, 6.62, 11.9,
        "9×9 · cap 1 700 (lossless; 1 500 truncated 0.0095 % of clusters) · 20 000 frames · "
        "opt1/opt2 are 3×3-only · CPU baseline ClusterFinderMT at 32 threads. The floor is the "
        "lower of the nsys estimate and the best sustained rate; both arms are sustained-bound, "
        "corroborated to 8.8 % (f64) and 0.4 % (f32). Why the arms differ: 544.5 kB of D2H "
        "costs 25.2 µs on both [s4], hidden under the 32.7 µs f64 kernel, not under "
     "the 23.9 µs f32 one.")

# =========================================================== 24 · WHERE TIME GOES
s = new_slide()
chrome(s, 25, "Where the time actually went",
       "The host bar dies first, then the floor itself drops")
figure(s, "fig_overhead", 1.95, 1.95, 9.4)
callout(s, M, 5.30, 5.85,
        "**Acts I and II never touch the arithmetic.** The GPU floor is a flat "
        "16.2 µs at 3×3 / 30.0 µs at 9×9; what collapses is everything stacked on it.",
        h=0.86, size=10.5)
callout(s, 6.75, 5.30, 5.85,
        "**Act III is the only step that lowers the floor itself**, and it could not "
        "have been seen until the stack above it was gone.", h=0.86, size=10.5,
        color=AMBER)
caption(s, M, 6.45, 11.9,
        "Blue/white/amber = the GPU floor for that act's build — the LOWER of the s4 "
        "engine max (max of H2D, kernel, D2H; PCIe is full duplex, so never the sum) "
        "and the best sustained rate. At 3×3 the engine max sets it, 16.17 µs; at 9×9 "
        "f64 the sustained rate does, 30.01 against a 32.66 µs max. Grey = everything "
        "the host adds on top. At 9×9 the host contributes +50 µs at opt3 and nothing "
        "at opt6.")

# =========================================================== 24 · MEASUREMENT AUDIT
s = new_slide()
chrome(s, 26, "Behind the numbers · what they had to survive",
       "Three ways a GPU benchmark lies")
items = [
    ("First-touch page faults", AMBER,
     "Each run materialises ~10 GB of clusters. The first pass faults in ~2.6 M "
     "pages at 0.7 µs each, up to 4 s of pure OS work inside the timer.",
     "Fix: re-run until getrusage() minor faults plateau (< 200 k)."),
    ("CUDA-event kernel timing", AMBER,
     "avg_kernel_time_ms() measures elapsed time on a stream, including waiting "
     "for other streams. Under 8-stream load it over-reads by up to 3.5×.",
     "Fix: Nsight Systems per-instance times; 1 stream for exclusive numbers."),
    ("The profiler itself", AMBER,
     "Under nsys, wall time per frame inflates ~4× from API tracing.",
     "Fix: GPU op times from nsys, wall times from unprofiled runs."),
]
x = M
for title, col, body, fix in items:
    rect(s, x, 2.0, 3.83, 3.15, PANEL)
    rect(s, x, 2.0, 3.83, 0.035, col)
    tf = tb(s, x + 0.26, 2.28, 3.3, 0.6)
    run(para(tf, True, line=1.15), title, 13, PALE, bold=True)
    tf = tb(s, x + 0.26, 3.02, 3.3, 1.5)
    run(para(tf, True, line=1.3), body, 10, TEXT2)
    tf = tb(s, x + 0.26, 4.42, 3.3, 0.65)
    run(para(tf, True, line=1.25), fix, 9.5, ACCENT)
    x += 4.03
code(s, M, 5.4, 11.9, [
    "# every timed cell in the benchmark notebook is bracketed with:",
    "mf0 = resource.getrusage(resource.RUSAGE_SELF).ru_minflt",
    "...   t = time.perf_counter() - t0   ...",
    "print(f'minor faults: {mf1-mf0:,}')   # quote the run where this plateaus",
], title="THE FAULT PROTOCOL · python/tests/ClusterFinderCUDA_perf.ipynb")
callout(s, M, 6.52, 11.2,
        "Validated: **wall = steady-state + faults × 0.68 µs** reproduced a 6.110 s "
        "run to within **1 ms**. Kernel time stayed constant throughout; the GPU was never the variable.",
        h=0.70, size=9.5)

# =========================================================== 25 · FIRST RUN
s = new_slide()
chrome(s, 27, "Behind the numbers · what a user actually gets",
       "A first run loses a third of its throughput to page faults")
figure(s, "fig_first_run", 1.37, 1.70, 10.6)
callout(s, M, 5.98, 5.85,
        "**Everything that materialises clusters loses a third of its throughput "
        "on the first run**, +7 to +20 µs per frame, depending on how much of it "
        "the GPU can hide behind its own work.", h=0.68, size=9.5)
callout(s, 6.75, 5.98, 5.85,
        "**Only the two ends escape, for opposite reasons.** opt1 never grows the "
        "heap: it discards each frame. opt6 never needs one, and reaches **98 % "
        "of its peak on a cold process**.", h=0.68, size=9.5, color=AMBER)
caption(s, M, 6.76, 11.9,
        "Single pass, one process, every ClusterVector retained, "
        "python/tests/ClusterFinderCUDA_perf.ipynb, f32, the same 100 000 frames. Not the "
        "campaign's \"cold\" rep, which discards results and so never grows the heap. Repeat "
        "the run and the amber bars climb onto the blue ones; opt1 and opt6 never move, "
        "because neither ever paid.", size=8.5)

# --------------------------------------------------- divider · VALIDATION
section("Validation · does it find the same photons",
        "Every number so far assumed the answers are identical",
        "Whether the CUDA finder returns the same clusters as the CPU.",
        [("28–29", "The fair comparison"),
         ("30–31", "The residual, dissected"),
         ("32–33", "For users"),
         ("34", "What is next")],
        rng=(28, 34), col=PALE,
        carry=("Established", "×9.1 and ×26.5",
               "on the hardware floor at both cluster sizes, if the physics holds"))

# ============================================ 26 · PEDESTAL UPDATE TIMING
s = new_slide()
chrome(s, 28, "Validation · why a CPU twin was needed",
       "CPU and CUDA update the pedestal at different moments")
figure(s, "fig_pedtiming", M - 0.15, 1.90, 12.2)
callout(s, M, 5.30, 5.85,
        "The serial CPU finder updates the pedestal **as the raster scan reaches each "
        "pixel**. A CUDA thread cannot: 160 000 of them read the pedestal at once, so "
        "the update is **applied at the frame boundary**.", h=1.02, size=10.5)
callout(s, 6.75, 5.30, 5.85,
        "So a straight CPU↔CUDA comparison moves **two** things at once. "
        "**ClusterFinderFrozen** is the serial finder with only the update moved to "
        "the frame end: same arithmetic, same gates, same scan.", h=1.02, size=10.5,
        color=AMBER)
caption(s, M, 6.52, 11.9,
        "Frozen is a diagnostic twin, not a product: it exists so the next two slides "
        "can attribute each disagreement to exactly one cause. cpu vs frozen = update "
        "timing; frozen vs cuda = the port. Every finder on these slides is trained on "
        "the same 1 000 pedestal frames and run over the same 10 000 data frames.")
notes(s, """Why this slide is here.

The obvious experiment - run the CPU finder and the CUDA finder over the same
frames and count the differences - cannot answer the question anyone actually
asks, which is "is the port correct?" It moves two variables at once:

  1. WHEN the pedestal is updated. The serial CPU finder updates per pixel,
     during the raster scan, so a pixel late in the frame is judged against a
     pedestal that already contains this frame's earlier pixels. That is
     scan-order dependent by construction.
  2. WHAT arithmetic runs. float vs double, a different local-max expression,
     a different rounding point.

ClusterFinderFrozen holds (2) fixed and changes only (1): it is the serial
finder, same decision logic line for line, with the update deferred to the frame
boundary. That is the CUDA update model, on the CPU.

So the three-way comparison factorises:
    cpu    vs frozen  -> update timing alone      (19 of 23.2 M)
    frozen vs cuda    -> the port alone           (6 of 23.2 M)

and the second is the number that answers the question. It is also why the
headline mismatch figure in this deck is 6 / 23 M and not 25 / 23 M: the larger
number is dominated by an effect that has nothing to do with CUDA.

Frozen ships in the library as a diagnostic, not as the recommended finder.""")

# =========================================================== 26 · CORRECTNESS
s = new_slide()
chrome(s, 29, "Validation · isolating one variable at a time",
       "CUDA and its CPU twin agree exactly: 0 in 23 million")
bullets(s, M, 1.86, 12.0, [
    "**ClusterFinderFrozen** makes byte-for-byte the same decisions as ClusterFinder "
    "and differs in exactly one thing: **when** the pedestal is updated. Frozen per "
    "frame, pushed at frame end. That is the CUDA model, so comparing against it "
    "isolates everything else the port changes.",
], size=10.5)
table(s, M, 2.52, 12.0,
      ["comparison", "the one thing that differs", "A-only / B-only",
       "% of clusters"],
      [["serial CPU  vs  frozen CPU", "update timing alone: **a CPU-only effect**",
        "8 / 11", "0.000082 %"],
       ["serial CPU  vs  CUDA  [f64 ped]", "the same thing, **and nothing else**",
        "8 / 11", "0.000082 %"],
       ["frozen CPU  vs  CUDA  [f64 ped]", "**nothing**", "**0 / 0**", "**0 %**"],
       ["frozen CPU  vs  CUDA  [f32 ped]", "the float32 pedestal EMA drifting: "
        "see next slide", "0 / 6", "0.000026 %"]],
      colw=[0.28, 0.40, 0.17, 0.15], size=9.5, rowh=0.60)
callout(s, M, 5.76, 5.85,
        "**Identical, frame by frame.** 23 244 605 clusters, not one disagreement, "
        "and **cpu vs cuda** equals **cpu vs frozen** exactly, so the port adds nothing "
        "of its own. Payloads: **99.9941 % bit-identical**, worst case 1 ADU on one "
        "pixel.", h=0.74, size=9.5)
callout(s, 6.75, 5.76, 5.85,
        "**The 0.004 % headline is the baseline disagreeing with itself.** "
        "ClusterFinderMT builds **48 ClusterFinders, each with its own Pedestal**, "
        "after 100 k frames each has ~2 083 updates, not 100 000. It measures the "
        "baseline, not the port.", h=0.74, size=9.5, color=AMBER)
caption(s, M, 6.58, 11.9,
        "3×3 · 10 000 frames · 23.2 M clusters · same pedestal, same frames · exact "
        "centre-set difference at tol = 0 · python/tests/validation_tiers.py. [f64 ped] "
        "and [f32 ped] are the same source built with DEVICE_PED_TYPE double / float; "
        "COMPUTE_TYPE is float in both, so the stencil arithmetic is identical across "
        "the two rows.")

# ================================================ 28 · THE MISMATCH, SEEN
s = new_slide()
chrome(s, 30, "Validation · the disagreement, seen",
       "The whole disagreement is one duplicate centre")
figure(s, "fig_mismatch147", 1.97, 1.78, 9.4)
callout(s, M, 5.82, 5.85,
        "Only the **3×3 footprints of each finder's own centres** are drawn, "
        "everything else is blank, so the panels differ **exactly** where the finders "
        "do. Values are pedestal-subtracted ADU.", h=0.94, size=10.5)
callout(s, 6.75, 5.82, 5.85,
        "cuda's patch is **one row taller**: a second centre directly below the one "
        "both found, so the two 3×3 windows overlap. **The charge is already "
        "counted**: a duplicate, not a new photon.", h=0.94, size=10.5, color=AMBER)
caption(s, M, 6.92, 11.9,
        "Frame 147, the strongest of the six residuals · shipping f32 build · every "
        "other centre in the patch agrees, including the ordinary photon at bottom right.")
notes(s, """The masked view, and how to read it.

Each panel masks every pixel that is not inside the 3x3 footprint of one of THAT
finder's own cluster centres. So a pixel is visible on the left only if frozen
claimed it, and visible on the right only if cuda did. Where the panels look
identical, the finders agreed pixel for pixel.

The red dots are cluster centres. The amber ring is the one centre cuda keeps
that frozen does not: (202, 8), directly below (202, 7) which both finders keep.
Because the two 3x3 windows overlap, cuda's patch is four rows tall where
frozen's is three - that visible extra row IS the disagreement.

Note what is NOT happening: there is no cluster in one panel that is absent from
the other's neighbourhood entirely. Nothing was missed, and nothing was invented
out of noise. The two finders are arguing about which of two adjacent pixels owns
a photon they both found.

Same construction as helper.plot_masked_mismatch() in the notebook, so this is
reproducible from ClusterFinderFrozen_vs_CUDA.ipynb directly.""")

# ========================================================== 27 · RESIDUALS
s = new_slide()
chrome(s, 31, "Validation · the six residuals, dissected",
       "float32 cannot tell these two pixels apart")
code(s, M, 1.92, 6.35, [
    "frame 147    centre (x=202, y=8)      3×3 window",
    "",
    "raw window (ADU)     pedestal-subtracted, 1 decimal",
    "[[4646 5282 4703]    [[ 45.3 «638.4» -12.1]   frozen and cuda",
    " [4857 5318 4950]     [ 43.7 «638.4»  -7.5]   print the SAME",
    " [4763 4640 4858]]    [  1.1  70.7  136.4]]   window",
    "",
    "the two contenders, at full precision:",
    "                          rival (dy=-1)        centre",
    "frozen  [f64 ped]      638.383019956   638.382773664",
    "cuda    [f32 ped]      638.382812500   «638.382812500»",
    "",
    "gate:  accept if  centre >= max(window)",
    "   frozen   638.382773664 >= 638.383019956  ->  reject",
    "   cuda     638.382812500 >= 638.382812500  ->  «ACCEPT»",
], size=8, title="THE ONLY TEST THAT FLIPS, AND WHY")
callout(s, M, 4.62, 6.35,
        "Separation under the f64 pedestal: **0.000246 ADU**. One float32 ULP at "
        "4 679.6 ADU is **0.000488**. The two pixels are **half a ULP apart**; in "
        "float32 they are *the same number*, and the gate accepts on a tie.",
        h=1.04, size=10)
figure(s, "fig_spectra_valid", 7.30, 1.92, 5.40)
callout(s, M, 5.82, 11.9,
        "**With the f64 pedestal there are none at all.** In the shipping f32 build "
        "each of the six sits **one pixel from a cluster both finders found**, a "
        "duplicate neighbour, never a spurious photon and never a missed one.",
        h=0.76, size=10.5, color=AMBER)
caption(s, M, 6.74, 11.9,
        "The same frame 147 as the previous slide, recomputed under each finder's "
        "decision-time pedestal. Both finders run the same gate (CPU value == max, CUDA "
        "!(val < max)), so the tie separates them, not the expression.")
notes(s, """Frame 147, centre (x=202, y=8). The full numbers.

                              centre               rival      centre - rival
  frozen (f64 host ped)  638.382773664       638.383019956        -0.000246291
  cuda   (f32 dev  ped)  638.382812500       638.382812500         0.000000000

  pedestal mean, centre:  f64 4679.617226336   f32 4679.617187500
  pedestal mean, rival :  f64 4643.616980044   f32 4643.617187500

  1 float32 ULP at 4679.6 ADU = 0.000488 ADU
  the f64 separation between the two pixels = 0.50 ULP

Read that last line slowly: the two pixels differ by HALF a float32 ULP. There is
no float32 number between them, so once the pedestal is stored as float32 they
round to the identical bit pattern. The local-max gate is `value >= max`, which
accepts on equality - so CUDA keeps the centre and frozen, which can still see
the rival is 2.5e-4 ADU higher, rejects it.

This is not drift in the usual sense. It is not that the f32 EMA wandered away
from the f64 one over thousands of updates - look at the pedestal columns, they
agree to 4e-5 and 2e-4 ADU. It is that the QUESTION being asked ("which of these
two pixels is larger?") has an answer that float32 cannot represent.

Consequences worth stating out loud if asked:
  - It is one-directional. cuda-only 6, frozen-only 0. A tie can only ever ADD a
    centre, never remove one, because >= accepts.
  - All six are at Chebyshev distance 1 from a cluster both finders found: a
    duplicate neighbour, never a spurious photon.
  - With DEVICE_PED_TYPE = double the residual is 0 / 0 at both 3x3 and 9x9.
  - Fixing it in f32 would mean a strict > in the gate, which changes the CPU's
    documented behaviour on genuine ties. Not worth 6 clusters in 23 million.""")

# =========================================================== 28 · API 1
s = new_slide()
chrome(s, 32, "For users · Python API", "The fast path in eight lines")
code(s, M, 1.95, 7.6, [
    "from aare import File, ClusterFinderCUDA",
    "",
    "cf = ClusterFinderCUDA(image_size=(400, 400), cluster_size=(3, 3),",
    "                       n_sigma=5, «n_streams»=4,",
    "                       «max_clusters_per_frame»=3000)",
    "",
    "for _ in range(1000):                    # 1. train the pedestal",
    "    cf.push_pedestal_frame(pd.read_frame())",
    "",
    "data = f.read_n(100_000)                 # 2. one contiguous array",
    "cf.«register_input_buffer»(data)           # 3. pin it once",
    "",
    "for s in range(0, N, 2000):              # 4. batch through it",
    "    clusters = cf.«find_clusters_batched»(data[s:s+2000], first_frame=s)",
    "",
    "cf.unregister_input_buffer()             # 5. release the pages",
], size=9, title="THE RECOMMENDED PATTERN")
bullets(s, 8.6, 2.0, 4.1, [
    "find_clusters_batched returns **one ClusterVector per frame**, in order, and "
    "does opt5's chunked overlap internally, so you get it for free.",
    "register_input_buffer is what turns opt3 into opt4: **one call**.",
], size=10)
code(s, 8.6, 3.55, 4.1, [
    "# opt6: never copy",
    "for v in cf.«find_cluster_views_",
    "          batched_iter»(data, 2000):",
    "    hist.fill(v.sums())",
    "    # views die with the chunk",
], size=8, title="IF YOU REDUCE AS YOU GO")
callout(s, 8.6, 5.25, 4.1,
        "**×1.16 at 3×3, ×2.21 at 9×9.** The views expose every cluster; they only "
        "withhold ownership past the chunk.", h=0.95, size=10)
callout(s, 8.6, 6.35, 4.1,
        "Pin **once**, outside the loop. Slices of a registered array inherit the "
        "pinning.", h=0.72, size=10, color=AMBER)

# =========================================================== 27 · API 2
s = new_slide()
chrome(s, 33, "For users · choosing the knobs", "Four knobs, and the one that silently truncates")
hdr = [("Parameter", 1.05), ("What it does", 3.6), ("Guidance", 5.2)]
y = 2.0
rect(s, M, y, 11.9, 0.4, PANEL)
for lab, dx in hdr:
    tf = tb(s, M + dx - 0.85 if dx > 1.05 else M + 0.28, y + 0.09, 5.0, 0.3)
    run(para(tf, True), lab.upper(), 9, MUTED, bold=True, spc=1.3)
y += 0.44
params = [
    ("n_streams", "How many frames may be in flight at once — an upper bound, not a "
     "count: the copy engines cap the real number below it.",
     "4 at both sizes. 8 buys no kernel concurrency at 9×9 (+1% instance time) and "
     "inflates the CUDA-event timer 3.5×."),
    ("max_clusters_per_frame", "Fixed size of the per-frame D2H transfer.",
     "Must exceed the real maximum or clusters are silently dropped, and it sets the "
     "D2H bar directly. Measured at 9×9: the maximum is 1 633, and the lossless "
     "cap of 1 700 already makes D2H the bottleneck on the f32 build."),
    ("batch size", "Frames per find_clusters_batched call.",
     "2 000 amortises launch overhead without a large pinned footprint."),
    ("cluster_size", "Compile-time stencil geometry.",
     "3×3 and 9×9 are registered; 9×9 moves the bottleneck off H2D and onto the "
     "kernel on f64, and onto D2H once opt7 shortens it."),
    ("register_input_buffer", "Page-locks the host array for DMA.",
     "Always, if the data is already in RAM. Check the pinning budget first."),
]
for i, (p_, what, guide) in enumerate(params):
    if i % 2 == 0:
        rect(s, M, y, 11.9, 0.82, PANEL)
    tf = tb(s, M + 0.28, y + 0.14, 2.6, 0.5)
    run(para(tf, True, line=1.1), p_, 9.5, ACCENT, font=MONO, bold=True)
    tf = tb(s, M + 3.0, y + 0.14, 2.9, 0.6)
    run(para(tf, True, line=1.2), what, 9.5, PALE)
    tf = tb(s, M + 6.15, y + 0.14, 5.4, 0.6)
    run(para(tf, True, line=1.2), guide, 9.5, TEXT2)
    y += 0.80
callout(s, M, 6.55, 11.2,
        "The single most common mistake: leaving **max_clusters_per_frame** too low. "
        "It does not error; it truncates, and every frame quietly returns the same count.",
        h=0.66, size=10, color=AMBER)

# =========================================================== 28 · NEXT
s = new_slide()
chrome(s, 34, "Where this leaves us",
       "The bottleneck has walked from the host, to the GPU, to the wire")
cards = [
    ("DONE", ACCENT, "×9.1 at 3×3, ×26.5 at 9×9",
     "16.3 and 25.1 µs/frame end to end, both on their hardware floor. Against the CPU "
     "twin that isolates the port the decisions are identical; the shipping f32 "
     "pedestal adds 6 duplicates in 23 million."),
    ("DONE", ACCENT, "FP32 pedestal, safely",
     "−40.5% kernel, and correct, because the variance is now accumulated on a frozen "
     "per-pixel offset instead of a raw second moment."),
    ("NEXT", AMBER, "3×3: transfer granularity",
     "The 16.31 µs sustained sits 3.16 µs above the uncontended 13.15 µs H2D rate "
     "[f32: s4 vs s1]: "
     "2 000 separate 320 kB descriptors, plus 26% of H2D↔D2H contention."),
    ("NEXT", AMBER, "9×9: the result path, not the kernel",
     "At a lossless cap D2H already binds [f32 · s4]: 25.24 µs against a 23.94 µs "
     "kernel. More "
     "kernel work buys nothing until the D2H slot stops being cap-sized."),
]
for i, (tag, col, title, body) in enumerate(cards):
    cx = M + (i % 2) * 6.05
    cy = 2.05 + (i // 2) * 2.35
    rect(s, cx, cy, 5.85, 2.05, PANEL)
    rect(s, cx, cy, 5.85, 0.035, col)
    tf = tb(s, cx + 0.3, cy + 0.26, 1.4, 0.26)
    run(para(tf, True), tag, 8.5, col, bold=True, spc=1.5)
    tf = tb(s, cx + 0.3, cy + 0.60, 5.2, 0.4)
    run(para(tf, True, line=1.1), title, 14, PALE, bold=True)
    tf = tb(s, cx + 0.3, cy + 1.12, 5.2, 0.85)
    run(para(tf, True, line=1.3), body, 10, TEXT2)
callout(s, M, 6.58, 11.2,
        "Full numbers, methodology and reproduction steps: **docs/ClusterFinderCUDA_benchmark_results.md** · "
        "notebook **python/tests/ClusterFinderCUDA_perf.ipynb** · measurement campaign **python/tests/perf/**",
        h=0.66, size=9.5)

# --------------------------------------------------------- divider · ANNEX
section("Annex · the evidence behind the claims",
        "Kept back so the arc stays readable",
        "The measurement artefacts, the rejected route, and the opt5 code.",
        [("A1", "Every engine number, reconciled"),
         ("A2", "CUDA Graphs, and why they lost"),
         ("A3", "opt5 · the overlap code"),
         ("A4", "The fault model, tested"),
         ("A5", "The three benchmark artefacts")],
        rng=(1, N_ANNEX), col=AMBER, annex=True,
        carry=("Everything so far", "34 slides",
               "the arc is finished; what follows answers questions"))

# ===========================================================================
# ANNEX — the measurement detail behind slides 26–27, and the rejected route A
# ===========================================================================

# ---- A1 · THE CONVENTION -------------------------------------------------
s = new_slide()
annex_chrome(s, 1, "measurement convention · expands slide 12",
             "Uncontended, or as the pipeline runs it")
bullets(s, M, 1.90, 12.0, [
    "Every engine time in this deck is tagged **[build · s1|s4]**. **s1** is one "
    "stream with nothing else running: what an engine does **on its own**, which is "
    "the right number for a capability claim and for the headroom that remains. "
    "**s4** is the shipped four-stream pipeline: each engine's **busy time per "
    "frame**, the union of its intervals — the only number that can set a floor.",
], size=10.5)
table(s, M, 2.78, 5.8,
      ["3×3 · µs/frame", "s1 f64", "s1 f32", "s4 f64", "s4 f32"],
      [["H2D", "13.14", "13.15", "16.17", "16.63"],
       ["kernel", "14.72", "4.32", "15.17", "5.53"],
       ["D2H", "5.31", "5.27", "7.69", "7.57"],
       ["engine max [s4]", "—", "—", "16.17", "16.63"],
       ["FLOOR = lower of max, sustained", "—", "—", "**16.17**", "**16.31**"]],
      colw=[0.28, 0.18, 0.18, 0.18, 0.18], size=9, rowh=0.44)
table(s, 7.0, 2.78, 5.8,
      ["9×9 · cap 1700", "s1 f64", "s1 f32", "s4 f64", "s4 f32"],
      [["H2D", "13.20", "13.22", "20.77", "20.54"],
       ["kernel", "39.86", "23.70", "32.66", "23.94"],
       ["D2H", "21.97", "21.95", "25.25", "25.24"],
       ["engine max [s4]", "—", "—", "32.66", "25.24"],
       ["FLOOR = lower of max, sustained", "—", "—", "**30.01**", "**25.14**"]],
      colw=[0.28, 0.18, 0.18, 0.18, 0.18], size=9, rowh=0.44)
callout(s, M, 5.66, 12.0,
        "**Two traps this table closes.** (1) s4 is engine *occupancy*, not duration — "
        "the 9×9 kernel row **falls** 39.86 → 32.66 while every transfer rises, and "
        "nothing that measures a duration falls under load. (2) **the FLOOR is not the "
        "engine max**: the max is profiled and runs 2–8 % high, so the floor is whichever is "
        "lower, it or the best unprofiled sustained rate — which is why opt6 reports "
        "**30.01 µs against a 32.66 max**.", h=0.92, size=10.5)
caption(s, M, 6.66, 12.0,
        "THE D2H SHIFT IS AN s4 PHENOMENON: at s1 the kernel binds in BOTH arms "
        "(39.86 and 23.70 against a 21.97 / 21.95 D2H); only under four-stream "
        "contention does D2H climb to 25.24 and overtake the 23.94 f32 kernel. Any "
        "claim about which engine binds must be read from the s4 columns. "
        "NEITHER DIRECTION IS FASTER THAN THE OTHER. At s1, H2D moves 320 000 B in "
        "13.15 µs = 24.3 GB/s; D2H moves 120 004 B in 5.27 µs = 22.8 GB/s; at 9×9 D2H "
        "moves 557 604 B in 21.95 µs = 25.4 GB/s. All sit at 72–81 % of PCIe 4.0 ×16's "
        "31.5 GB/s, the smallest transfer paying the most fixed cost per byte — so a "
        "taller bar always means MORE BYTES, never a slower wire. That is why cap 1 700 "
        "hands the 9×9 floor to D2H. Source: probes.csv in perf/results/"
        "2026-08-18_{f64,f32}/ (3×3) and 2026-08-20_{f64,f32}_capAB/ (9×9). The floor is the "
        "LOWER of the s4 estimate and the best unprofiled sustained rate — 16.31 vs "
        "16.63 at 3×3 f32, 25.14 vs 25.24 at 9×9 f32 — so a probe roofline is never a "
        "hard denominator.")

# ---- A2 · ROUTE A · IDEA -------------------------------------------------
s = new_slide()
annex_chrome(s, 2, "rejected route · CUDA graphs · expands slide 11",
             "CUDA Graphs, a sound idea that the next act overtook", part=1, nparts=2)
bullets(s, M, 1.95, 12.0, [
    "Every cudaMemcpyAsync / kernel launch costs the **CPU** a few microseconds of "
    "driver work, per frame and per operation. After opt4 that looked like the budget.",
    "A **CUDA Graph** captures the whole dependency DAG once. Replaying it is a "
    "**single** cudaGraphLaunch: the driver already knows every node and edge.",
], size=10.5)
figure(s, "fig_graphs", M, 3.15, 7.6)
code(s, 8.5, 3.15, 4.1, [
    "// record once, at setup",
    "cudaStreamBeginCapture(sc.stream, ...);",
    "   submit_h2d_kernel_d2h(sc);",
    "cudaStreamEndCapture(sc.stream, &sc.graph);",
    "«cudaGraphInstantiate»(&sc.graphExec, ...);",
    "",
    "// per batch: one call",
    "«cudaGraphLaunch»(sc.graphExec, sc.stream);",
], size=8, title="ClusterFinderCUDA_graph.hpp")
callout(s, 8.5, 4.88, 4.1,
        "**REJECTED**\n3×3: 39 752 FPS, inside noise of opt4.\n9×9: **11 072 FPS, 12 % slower**.",
        h=1.10, size=10.5, color=AMBER)
caption(s, M, 6.62, 12.0,
        "Its 3×3 edge was never established: the graph finder recorded no CUDA events "
        "while the stream finder did, and that instrumentation tax (2.8 µs) is larger "
        "than the gap (0.8 µs). More decisively, it never received the chunked pipeline "
        "of opt5, so it is competing on ~2 µs of launch cost against a 24 µs floor. "
        "Launch overhead stops binding one step later; the technique aimed at it can no "
        "longer pay.")

# ---- A3 · ROUTE A · BUDGET -----------------------------------------------
s = new_slide()
annex_chrome(s, 2, "rejected route · CUDA graphs",
             "What a CUDA Graph actually saves, in microseconds", part=2, nparts=2)
bullets(s, M, 1.90, COL, [
    "The stream path issues **four runtime calls per frame**, a memset to clear the "
    "cluster counter, H2D, the launch, D2H. A graph replaces all four with **one** "
    "cudaGraphLaunch, so the ceiling on what it can save is **¾ of the submission cost**.",
], size=10.5)
table(s, M, 2.72, COL,
      ["call", "per frame", "host cost", "µs/frame"],
      [["cudaMemcpyAsync", "2", "1.98 µs", "3.97"],
       ["cudaLaunchKernel", "1", "2.13 µs", "2.13"],
       ["cudaMemsetAsync", "1", "1.57 µs", "1.57"],
       ["**submission total**", "**4**", "", "**7.67**"]],
      colw=[0.40, 0.18, 0.22, 0.20], size=9, rowh=0.44)
code(s, M, 5.42, COL, [
    "7.67 us/frame x 3/4  = 5.75 us/frame   eliminated, AS MEASURED (under nsys)",
    "5.75 / 4 (see A5)    ~ 1.4 us/frame    eliminated, unprofiled estimate",
], size=8, title="THE ARITHMETIC")
callout(s, M, 6.28, COL,
        "**~1.4 µs against a 16.17 µs floor = 8.7 %**, real while the host is the "
        "critical path, and **worth nothing after opt5**, which hides host work under "
        "the GPU entirely.", h=0.80, size=10, color=AMBER)
rail(s, [
    ("label", "route A · measured verdict"),
    ("gap", 0.10),
    ("stat", "3×3 vs opt4", "×1.03", TEXT2),
    ("stat", "9×9 vs opt4", "×0.88", AMBER),
    ("gap", 0.05),
    ("row", "vs opt5 at 9×9", "36 % behind", AMBER),
    ("gap", 0.20),
    ("note", "Read from CUPTI_..._RUNTIME, the host table, so this is the one "
             "order-of-magnitude number in the deck, ±50 %. It is enough to show "
             "graphs cannot pay against a 16 µs floor, and not enough to quote to "
             "two figures."),
])


# ---- A4 · THE OPT5 CODE --------------------------------------------------
s = new_slide()
annex_chrome(s, 3, "opt5 · the overlap code · expands slide 16",
             "The overlap, in six lines, and why you never write them")
code(s, M, 1.86, 7.15, [
    "tok = cf.«submit_batch»(data[a0:b0], first_frame=a0)",
    "for a, b in bounds[1:]:",
    "    nxt = cf.«submit_batch»(data[a:b], first_frame=a)   # GPU starts N+1 …",
    "    results.extend(cf.«collect»(tok))                   # … host unpacks N",
    "    tok = nxt",
    "results.extend(cf.«collect»(tok))                       # drain the last one",
], size=8.5, title="THE WHOLE OF OPT5")
bullets(s, M, 3.66, 7.15, [
    "**You do not write this.** It is inside find_clusters_batched(), which chunks "
    "the batch and runs the loop for you, opt5 arrived as a **speedup, not an API "
    "change**, and every existing caller got it without editing a line.",
    "submit_batch() and collect() stay public for anyone who wants the token by "
    "hand, streaming from a detector, interleaving other work between chunks.",
], size=10)
rail(s, [
    ("label", "chunk sizing · the two constraints"),
    ("gap", 0.12),
    ("row", "multiple of", "n_streams", ACCENT),
    ("row", "capped at", "MAX_SLOT_BYTES", ACCENT),
    ("gap", 0.16),
    ("note", "The chunk MUST be a multiple of n_streams. The device pedestal is "
             "per-stream and advances once per frame the stream sees, so an uneven "
             "chunk leaves the four pedestals at different ages, the finder would "
             "stop being reproducible, and two runs of the same data would not agree."),
    ("gap", 0.14),
    ("note", "It is also capped so the two pinned output slots stay bounded: the "
             "slot is chunk × (4 + cap × sizeof(Cluster)), which at 9×9 and cap 1 700 "
             "is 544.5 kB per frame. chunk_size_for(n) applies both rules; pass the "
             "result to reserve_output_slots() to pre-pay the page-locking outside "
             "any timed region."),
])
callout(s, M, 5.92, 7.15,
        "Two chunks in flight is enough. A third adds pinned memory and no overlap: "
        "the host is already busy for the whole time the GPU is.", h=0.72, size=10)
caption(s, M, 6.80, 12.0,
        "The saving is min(GPU, host) per chunk, so opt5 pays most when the two terms "
        "are comparable, ×1.31 at 3×3, and least when one dominates, ×1.20 at 9×9, "
        "where the host term is roughly twice the GPU term. That gap is the diagnosis "
        "that motivates opt6 (report §8.2).")

# ---- A5 · THE FAULT MODEL ------------------------------------------------
s = new_slide()
annex_chrome(s, 4, "the fault model · expands slide 22",
             "The fault model, tested against every step")
bullets(s, M, 1.90, 12.0, [
    "Slide 26 fits **0.68 µs per first-touch fault**. If that rate is right it should "
    "predict the f64→f32 gap at every step from the fault counts alone — and where it "
    "does, the comparison is measuring **the allocator**, not the typedef.",
], size=10.5)
table(s, M, 2.78, 12.0,
      ["step", "f64 warm (faults)", "f32 warm (faults)", "Δ wall",
       "Δ faults", "predicted", "verdict"],
      [["opt3", "82.44 µs  (128 k)", "95.66 µs  (521 k)", "**+13.22**", "+393 k",
        "**+13.37**", "**the allocator**"],
       ["opt4", "79.83 µs  (128 k)", "75.20 µs  (127 k)", "−4.63", "−0.5 k",
        "−0.02", "clean"],
       ["opt5", "66.39 µs  (152 k)", "61.85 µs  (10 k)", "−4.54", "−141 k",
        "**−4.81**", "**not separable**"],
       ["opt6", "30.01 µs  (0)", "25.14 µs  (0)", "−4.87", "0", "0.00", "clean"]],
      colw=[0.08, 0.18, 0.18, 0.10, 0.10, 0.13, 0.23], size=8.5, rowh=0.50)
callout(s, M, 5.42, 12.0,
        "**opt3 is the whole argument in one row.** A 393 k fault gap predicts +13.37 µs; "
        "+13.22 was observed — agreement to **1 %**. The −40 % kernel is in there "
        "somewhere, invisible under 13 µs of the OS zeroing pages. Only opt4 and opt6, "
        "where the fault term is ~0, report the typedef at all.", h=0.90, size=10.5)
caption(s, M, 6.52, 12.0,
        "ladder_9x9.csv in results/2026-08-20_{f64,f32}_cap1700/, warm = best of reps "
        "1–4, faults are that rep's own getrusage minor-fault count. Predicted = Δfaults "
        "× 0.68 µs ÷ 20 000 frames. This table replaces an earlier figure that quoted "
        "opt3's +16 % as a measurement of the result path; it is a measurement of two "
        "allocator states, and the model above is how that was established.")

# ---- A6 · FAULTS ---------------------------------------------------------
s = new_slide()
annex_chrome(s, 5, "benchmark artefacts · expands slide 26",
             "First-touch page faults: two sources, one counter", part=1, nparts=3)
bullets(s, M, 1.90, 12.0, [
    "A page exists in the process's address space but has no physical frame yet. "
    "On first touch the kernel finds one, **zeroes it** (mandatory), and maps it. "
    "No disk I/O: ru_majflt stays 0 all campaign. At 4 kB/page, **1 GB = 262 144 faults**.",
], size=10.5)
table(s, M, 2.62, 12.0,
      ["", "(a) result heap", "(b) pinned D2H slots"],
      [["allocator", "malloc → mmap, one ClusterVector per frame",
        "cudaMallocHost, in submit_batch"],
       ["cost / page", "**0.7 µs**", "**1.0 µs**: same fault + pin + DMA map"],
       ["recurs?", "**yes**, every alloc/free cycle above the mmap threshold",
        "**no**: once per buffer, for its lifetime"],
       ["removed by", "**collect_view()**: it allocates nothing",
        "nothing; reserve_output_slots() only moves it out of the timer"]],
      colw=[0.14, 0.43, 0.43], size=9)
callout(s, M, 5.62, 5.85,
        "**The two are exactly additive.** Reserving subtracts precisely the pre-pin "
        "count from run 0 and changes nothing else.", h=0.72, size=10)
code(s, 6.75, 5.55, 5.85, [
    "3x3:   572 292 - 455 129 = 117 163   vs 117 192 pre-pin",
    "9x9: 2 759 037 - 2 278 567 = 480 470   vs 480 474",
    "closed form: 2 slots x 2000 x 120 004 B / 4 kB = 117 191",
], size=7.5, title="NOT A CORRELATION, AN IDENTITY")
callout(s, M, 6.42, 11.9,
        "**At 9×9 the heap never plateaus.** ~9.3 GB per pass is above glibc's mmap "
        "threshold, so it is munmap'd and re-faulted every pass: ~292 k faults ≈ "
        "**10 µs/frame, permanently**. No number of re-runs removes it, which is an "
        "independent argument for opt6.", h=0.72, size=10, color=AMBER)

# ---- A7 · EVENTS ---------------------------------------------------------
s = new_slide()
annex_chrome(s, 5, "benchmark artefacts",
             "CUDA events measure the stream, not the kernel", part=2, nparts=3)
bullets(s, M, 1.90, COL, [
    "avg_kernel_time_ms() brackets the launch with **cudaEventRecord on the "
    "kernel's own stream**. What it returns is elapsed time on **that stream's "
    "timeline**, which includes time spent **queued behind other streams**.",
    "So it is honest at 1 stream and inflates under saturation: up to **3.5×** at "
    "8 streams. The tell is that the derived *PCIe + overhead* = wall/N − kernel_ms "
    "**goes negative**: kernels overlap, so wall/frame < kernel/frame.",
    "It is also not free: **~3.6 µs/frame**, i.e. 10–15 % of end-to-end throughput "
    "at 3×3, to produce a number that is unusable exactly when it matters.",
], size=10.5)
code(s, M, 4.75, COL, [
    "if (m_time_kernels)                       // OFF by default, all 3 finders",
    "    cudaEventRecord(start[slot][i], sc.stream);",
    "find_clusters_in_single_frame<<<grid, block, shmem, «sc.stream»>>>(...);",
    "if (m_time_kernels)",
    "    cudaEventRecord(stop[slot][i], sc.stream);   // <- queue-wait lands here",
], size=8, title="ClusterFinderCUDA.hpp")
callout(s, M, 6.10, COL,
        "The flag exists for **comparability**, not preference: with events on for one "
        "finder and off for another, the step between them absorbs the tax.",
        h=0.72, size=10)
rail(s, [
    ("label", "9×9 kernel · f64 · nsys"),
    ("gap", 0.10),
    ("stat", "s1 · duration", "39.86 µs", ACCENT),
    ("stat", "s4 · occupancy", "32.66 µs", AMBER),
    ("gap", 0.05),
    ("row", "overlap factor", "1.32×", TEXT2),
    ("gap", 0.20),
    ("note", "Same kernel. The s4 column is the union of kernel intervals per frame, "
             "not how long one kernel takes. Quote s1 for duration; s4 feeds the floor, "
             "subject to the sustained-rate rule on slide 12. Full grid: A1."),
])

# ---- A8 · NSYS -----------------------------------------------------------
s = new_slide()
annex_chrome(s, 5, "benchmark artefacts",
             "Where nsys is sound, and where it is not", part=3, nparts=3)
bullets(s, M, 1.90, 12.0, [
    "Tracing a **host** call means running a callback on entry and on exit, "
    "**inside the interval being measured**. GPU work is different: the hardware "
    "stamps its own start/end and the host reads those records **afterwards**, so "
    "nothing is injected into the execution path.",
], size=10.5)
table(s, M, 2.60, 12.0,
      ["measurement", "sqlite table", "9×9 [f64 · s1 · cap 1500]", "verdict"],
      [["cudaLaunchKernel: the host call", "CUPTI_..._RUNTIME", "1.85 µs",
        "**inflated ~4×**"],
       ["the kernel executing", "CUPTI_..._KERNEL", "39.93 µs", "sound to ~2 %"],
       ["cudaMemcpyAsync: the host call", "CUPTI_..._RUNTIME", "1.65 µs",
        "**inflated ~4×**"],
       ["the H2D / D2H transfer", "CUPTI_..._MEMCPY", "13.25 / 19.44 µs",
        "sound to ~2 %"]],
      colw=[0.36, 0.24, 0.22, 0.18], size=9, rowh=0.52)
callout(s, M, 5.30, 11.9,
        "**Same cudaMemcpyAsync, two numbers: 1.65 µs to ask for the copy, 13.25 µs "
        "for the copy to happen.** If the profiler must be present at the moment to "
        "measure it, it distorts it; if the hardware records it anyway and the "
        "profiler reads it later, it does not.", h=0.86, size=10.5)
caption(s, M, 6.40, 11.9,
        "Every headline number in this deck (kernel times, transfer times, duty "
        "cycles, overlap and every engine floor) comes from the _KERNEL and _MEMCPY "
        "tables (see gpu_span.py). ⚠ CAP 1500 — the ONLY cap-1500 numbers left in the "
        "deck, kept because the RUNTIME column exists only in the trace and the host "
        "call cost does not depend on the cap. So the "
        "19.44 µs D2H here is NOT the shipped 9×9 bar — at cap 1700 that engine reads "
        "21.95 [s1] and 25.24 [s4], which is what binds the f32 build. The one figure "
        "read from _RUNTIME is the graph "
        "launch budget in A2, and it is quoted to one significant figure "
        "for exactly this reason. Proof the GPU side is sound: opt7 sustains 25.14 µs "
        "unprofiled against a 25.24 µs estimate measured under the profiler, 0.4 % "
        "apart, which a 4× distortion could not survive.")

prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
