"""Rebuild docs/ClusterFinderCUDA_optimizations.pptx — opt1..opt6, in the deck's
own design language (extracted from the original file)."""
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pathlib import Path
from PIL import Image

FIGS = Path(__file__).parent / "figs"
OUT = Path("/home/ferjao_k/aare/docs/ClusterFinderCUDA_optimizations.pptx")

# ---------------------------------------------------------------- design tokens
BG     = RGBColor(0x0B, 0x10, 0x18)
PANEL  = RGBColor(0x12, 0x1A, 0x28)
CODEBG = RGBColor(0x0E, 0x14, 0x20)
RULE   = RGBColor(0x1E, 0x28, 0x36)
ACCENT = RGBColor(0x1E, 0x90, 0xC2)
AMBER  = RGBColor(0xE8, 0xB2, 0x5C)
PALE   = RGBColor(0xE7, 0xED, 0xF4)
TEXT2  = RGBColor(0xA5, 0xB2, 0xC4)
MUTED  = RGBColor(0x6B, 0x7A, 0x90)

UI, MONO = "Segoe UI", "Consolas"
W, H = 13.333, 7.5
M = 0.7                       # left margin
COL = 7.9                     # left column width
RAIL_X, RAIL_W = 9.2, 3.5     # right rail

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]
N_SLIDES = 19


# ------------------------------------------------------------------- helpers
def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = etree.SubElement(s._element, "{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
    pr = etree.SubElement(bg, "{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr")
    fill = etree.SubElement(pr, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    clr = etree.SubElement(fill, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    clr.set("val", "0B1018")
    etree.SubElement(pr, "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
    s._element.insert(0, bg)
    return s


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, first=False, space_after=0, space_before=0, line=None, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    if line: p.line_spacing = line
    if align: p.alignment = align
    return p


def run(p, text, size=11, color=TEXT2, font=UI, bold=False, italic=False, spc=None):
    r = p.add_run(); r.text = text
    f = r.font
    f.name, f.size, f.bold, f.italic = font, Pt(size), bold, italic
    f.color.rgb = color
    if spc is not None:
        r.font._rPr.set("spc", str(int(spc * 100)))
    return r


# ------------------------------------------------------------------ chrome
def chrome(s, idx, eyebrow, title, title_size=27):
    rect(s, M, 0.60, 0.35, 0.035, ACCENT)
    tf = tb(s, 1.17, 0.50, 10.33, 0.32)
    run(para(tf, True), eyebrow.upper(), 9, MUTED, bold=True, spc=1.6)

    tf = tb(s, M, 0.86, 11.9, 1.0)
    run(para(tf, True, line=1.05), title, title_size, PALE, bold=True)

    # progress bar
    span, n = 11.0, N_SLIDES
    pitch = span / n; wseg = pitch * 0.90
    for i in range(n):
        rect(s, M + i * pitch, 7.28, wseg, 0.045, ACCENT if i <= idx - 1 else RULE)
    tf = tb(s, 12.0, 7.14, 0.9, 0.3)
    run(para(tf, True, align=PP_ALIGN.RIGHT), f"{idx} / {n}", 8.5, MUTED)


def bullets(s, x, y, w, items, size=11, gap=7):
    tf = tb(s, x, y, w, 0.3)
    for i, it in enumerate(items):
        color, txt = (it if isinstance(it, tuple) else (TEXT2, it))
        p = para(tf, i == 0, space_after=gap, line=1.25)
        run(p, "•  ", size, MUTED)
        # inline emphasis with **...**
        for j, part in enumerate(txt.split("**")):
            if part:
                run(p, part, size, PALE if j % 2 else color, bold=bool(j % 2))
    return tf


def code(s, x, y, w, lines, size=8.5, title=None):
    lh = 0.148
    h = 0.24 + len(lines) * lh + (0.22 if title else 0)
    rect(s, x, y, w, h, CODEBG, MSO_SHAPE.ROUNDED_RECTANGLE)
    ty = y + 0.12
    if title:
        tf = tb(s, x + 0.18, ty, w - 0.36, 0.2)
        run(para(tf, True), title, 7.5, MUTED, bold=True, spc=1.2)
        ty += 0.22
    tf = tb(s, x + 0.18, ty, w - 0.36, h - 0.24)
    for i, ln in enumerate(lines):
        p = para(tf, i == 0, line=1.12)
        if ln.strip().startswith(("//", "#")):
            run(p, ln, size, MUTED, MONO)
            continue
        for j, part in enumerate(ln.split("«")):
            for k, seg in enumerate(part.split("»")):
                if not seg: continue
                hi = (j > 0 and k == 0)
                run(p, seg, size, ACCENT if hi else TEXT2, MONO, bold=hi)
    return h


def callout(s, x, y, w, text, h=0.78, color=ACCENT, size=10.5):
    rect(s, x + 0.045, y, w - 0.045, h, PANEL)
    rect(s, x, y, 0.045, h, color)
    tf = tb(s, x + 0.28, y + 0.10, w - 0.5, h - 0.2, MSO_ANCHOR.MIDDLE)
    p = para(tf, True, line=1.2)
    for j, part in enumerate(text.split("**")):
        if part:
            run(p, part, size, PALE if j % 2 else TEXT2, bold=bool(j % 2))


def rail(s, items, y0=2.0, divider=True):
    if divider:
        rect(s, 8.95, 2.0, 0.012, 4.55, RULE)
    y = y0
    for it in items:
        kind = it[0]
        if kind == "label":
            tf = tb(s, RAIL_X, y, RAIL_W, 0.26)
            run(para(tf, True), it[1].upper(), 8.5, MUTED, bold=True, spc=1.4)
            y += 0.28
        elif kind == "stat":
            _, lab, val, col = it
            tf = tb(s, RAIL_X, y, RAIL_W, 0.24)
            run(para(tf, True), lab.upper(), 8.5, MUTED, spc=1.2)
            tf = tb(s, RAIL_X, y + 0.24, RAIL_W, 0.6)
            run(para(tf, True), val, 26, col, bold=True)
            y += 0.98
        elif kind == "row":
            _, lab, val, col = it
            tf = tb(s, RAIL_X, y, RAIL_W, 0.24)
            run(para(tf, True), lab.upper(), 8.5, MUTED, spc=1.2)
            tf = tb(s, RAIL_X, y + 0.22, RAIL_W, 0.3)
            run(para(tf, True), val, 13, col, bold=True)
            y += 0.66
        elif kind == "note":
            tf = tb(s, RAIL_X, y, RAIL_W, 0.9)
            run(para(tf, True, line=1.25), it[1], 9, TEXT2)
            y += 0.30 + 0.17 * (len(it[1]) // 42 + 1)
        elif kind == "gap":
            y += it[1]
    return y


def figure(s, name, x, y, w):
    p = FIGS / f"{name}.png"
    iw, ih = Image.open(p).size
    h = w * ih / iw
    s.shapes.add_picture(str(p), In(x), In(y), In(w), In(h))
    return h


def caption(s, x, y, w, text, size=9):
    tf = tb(s, x, y, w, 0.3)
    run(para(tf, True, line=1.25), text, size, MUTED)


# =========================================================== 1 · TITLE
s = new_slide()
rect(s, 0, 0, 0.16, H, ACCENT)
tf = tb(s, M + 0.3, 0.85, 11, 0.3)
run(para(tf, True), "AARE · PSI HYBRID PIXEL DETECTORS · CUDA CLUSTERFINDER",
    9.5, MUTED, bold=True, spc=1.8)

tf = tb(s, M + 0.3, 1.35, 11.4, 1.7)
run(para(tf, True, line=1.02), "Feeding the GPU", 46, PALE, bold=True)
tf = tb(s, M + 0.3, 2.30, 11.4, 0.8)
run(para(tf, True, line=1.05), "Six optimization steps of the CUDA ClusterFinder",
    22, ACCENT)

tf = tb(s, M + 0.3, 3.25, 9.6, 0.8)
run(para(tf, True, line=1.3),
    "Five of the six steps never touch the arithmetic. They are about keeping "
    "a 24 µs kernel supplied with data — and about learning to measure honestly.",
    12.5, TEXT2)

stats = [("×8.29", "VS 48-THREAD CPU", ACCENT), ("39,472", "FRAMES / SECOND", PALE),
         ("25.3 µs", "PER FRAME, END TO END", PALE), ("0.004%", "CLUSTER-COUNT DRIFT", AMBER)]
for i, (v, l, c) in enumerate(stats):
    x = M + 0.3 + i * 2.85
    rect(s, x, 4.35, 0.035, 0.95, c)
    tf = tb(s, x + 0.22, 4.35, 2.5, 0.55)
    run(para(tf, True), v, 30, c, bold=True)
    tf = tb(s, x + 0.22, 4.98, 2.5, 0.3)
    run(para(tf, True), l, 8.5, MUTED, spc=1.2)

rect(s, M + 0.3, 6.05, 11.0, 0.012, RULE)
tf = tb(s, M + 0.3, 6.25, 11.4, 0.6)
run(para(tf, True, line=1.35),
    "RTX 4090 (Ada, sm_89) · PCIe 4.0 ×16 · Mönch 400×400 uint16 · 3×3 clusters · "
    "100 000 frames · Cu fluorescence, MAX IV", 10, MUTED)
tf = tb(s, M + 0.3, 6.62, 11.4, 0.3)
run(para(tf, True), "Khalil Ferjaoui · Paul Scherrer Institut", 10, TEXT2)

# =========================================================== 2 · PROBLEM
s = new_slide()
chrome(s, 2, "The problem & the baseline", "What has to happen to every frame")
bullets(s, M, 1.95, COL, [
    "Per pixel: subtract a **running pedestal** (mean ± rms), keep pixels above "
    "**nσ · rms**, cut a 3×3 cluster around each local maximum.",
    "400×400 = 160 k pixels, **312.5 kB per frame**; Cu data yields ~2 330 clusters "
    "per frame at 3×3.",
    "The pedestal is **updated by every non-photon pixel**, every frame — so the "
    "arithmetic and the data movement are coupled.",
])
code(s, M, 3.55, COL, [
    "// the whole algorithm, per pixel",
    "v      = frame[i] - pedestal_mean[i]",
    "rms    = sqrt(pedestal_sum2[i]/n - pedestal_mean[i]^2)",
    "if (v > «nSigma» * rms  &&  v == max(3x3 window))   ->  emit cluster",
    "else                                                 ->  update pedestal",
], title="THE KERNEL IN FIVE LINES")
callout(s, M, 5.55, COL,
        "**Thesis of this talk:** the compute was fast almost immediately. "
        "Five of six steps are about feeding it.")
rail(s, [
    ("label", "Baseline · same data, same threshold"),
    ("gap", 0.15),
    ("stat", "CPU, 1 thread", "1.75 ms", MUTED),
    ("stat", "CPU MT, 48 threads", "210 µs", PALE),
    ("gap", 0.1),
    ("row", "That is the bar", "4 761 frames / s", TEXT2),
    ("gap", 0.25),
    ("note", "Every CUDA number in this deck is measured against the 48-thread "
             "CPU on the same 100 000 frames."),
])

# =========================================================== 3 · THE LADDER
s = new_slide()
chrome(s, 3, "Roadmap", "Two acts: feed the GPU, then speed up the kernel")
rows = [
    ("opt1", "First CUDA port", "1 stream, one launch per frame", "×3.14", ACCENT),
    ("opt2", "Streams + batching", "4 streams, 2 000-frame batches", "×4.86", ACCENT),
    ("opt3", "Pipeline rework", "sync barriers removed", "×5.58", ACCENT),
    ("opt4", "Pinned memory", "DMA-speed host transfers", "×7.73", ACCENT),
    ("opt5", "CUDA Graphs", "one launch replaces six", "×8.29", ACCENT),
    ("opt6", "FP32 pedestal + variance rewrite", "the first kernel change", "kernel −40%", AMBER),
]
y = 2.05
for i, (tag, name, sub, gain, col) in enumerate(rows):
    rect(s, M, y, 11.9, 0.72, PANEL if i % 2 == 0 else BG)
    rect(s, M, y, 0.035, 0.72, col)
    tf = tb(s, M + 0.28, y + 0.13, 1.0, 0.4)
    run(para(tf, True), tag, 15, col, bold=True, font=MONO)
    tf = tb(s, M + 1.45, y + 0.10, 5.0, 0.3)
    run(para(tf, True), name, 13, PALE, bold=True)
    tf = tb(s, M + 1.45, y + 0.38, 5.6, 0.3)
    run(para(tf, True), sub, 10, MUTED)
    tf = tb(s, 9.4, y + 0.18, 3.1, 0.4)
    run(para(tf, True, align=PP_ALIGN.RIGHT), gain, 15, col, bold=True)
    y += 0.78
rect(s, M, 2.05, 0.012, 4.68, RULE)
caption(s, M, 6.72, 11.9,
        "opt1–opt5 change only how work is scheduled and moved — the arithmetic is "
        "byte-identical. opt6 is the first step that changes the kernel itself.")

# =========================================================== 4 · METHODOLOGY
s = new_slide()
chrome(s, 4, "Before any number is believed", "Three ways a GPU benchmark lies")
items = [
    ("First-touch page faults", AMBER,
     "Each run materialises ~10 GB of clusters. The first pass faults in ~2.6 M "
     "pages at 0.7 µs each — up to 4 s of pure OS work inside the timer.",
     "Fix: re-run until getrusage() minor faults plateau (< 200 k)."),
    ("CUDA-event kernel timing", AMBER,
     "avg_kernel_time_ms() measures elapsed time on a stream — including waiting "
     "for other streams. Under 8-stream load it over-reads by up to 3.5×.",
     "Fix: Nsight Systems per-instance times; 1 stream for exclusive numbers."),
    ("The profiler itself", AMBER,
     "Under nsys, wall time per frame inflates ~4× from API tracing.",
     "Fix: GPU op times from nsys, wall times from unprofiled runs."),
]
x = M
for title, col, body, fix in items:
    rect(s, x, 2.0, 3.83, 3.15, PANEL)
    rect(s, x, 2.0, 3.83, 0.035, col)
    tf = tb(s, x + 0.26, 2.28, 3.3, 0.6)
    run(para(tf, True, line=1.15), title, 13, PALE, bold=True)
    tf = tb(s, x + 0.26, 3.02, 3.3, 1.5)
    run(para(tf, True, line=1.3), body, 10, TEXT2)
    tf = tb(s, x + 0.26, 4.42, 3.3, 0.65)
    run(para(tf, True, line=1.25), fix, 9.5, ACCENT)
    x += 4.03
code(s, M, 5.4, 11.9, [
    "# every timed cell in the benchmark notebook is bracketed with:",
    "mf0 = resource.getrusage(resource.RUSAGE_SELF).ru_minflt",
    "...   t = time.perf_counter() - t0   ...",
    "print(f'minor faults: {mf1-mf0:,}')   # quote the run where this plateaus",
], title="THE FAULT PROTOCOL — python/tests/ClusterFinderCUDA_perf.ipynb")
callout(s, M, 6.52, 11.2,
        "Validated: **wall = steady-state + faults × 0.68 µs** reproduced a 6.110 s "
        "run to within **1 ms**. Kernel time stayed constant throughout — the GPU was never the variable.",
        h=0.70, size=9.5)

# =========================================================== 5 · OPT1
s = new_slide()
chrome(s, 5, "opt1 · the first CUDA port", "One frame, one stream, fully synchronous")
bullets(s, M, 1.95, COL, [
    "Shared-memory tiling with **halo loading** for any cluster size; pedestal "
    "subtraction fused into the tile load.",
    "Cluster geometry is a **compile-time template parameter** → the 3×3 stencil "
    "is fully unrolled.",
    "One cudaMemcpy in, one kernel, one cudaMemcpy out — **the host blocks "
    "on every frame**.",
])
code(s, M, 3.62, COL, [
    "// one frame at a time — the host waits at every step",
    "cudaMemcpy(d_frame, h_frame, bytes, cudaMemcpyHostToDevice);",
    "find_clusters_in_single_frame<ClusterType, FRAME_TYPE>",
    "    <<<grid, block, shmem>>>(d_frame, d_pd_mean, ...);",
    "cudaMemcpy(h_out, d_out, out_bytes, cudaMemcpyDeviceToHost);",
], title="ClusterFinderCUDAOpt2.hpp · find_clusters()")
callout(s, M, 5.62, COL,
        "The stencil was **already fast**: 23 µs of kernel inside a 67 µs frame. "
        "The other 44 µs is the host standing still.")
rail(s, [
    ("label", "opt1 · 3×3 · 100 k frames"),
    ("gap", 0.15),
    ("stat", "End to end", "66.8 µs", PALE),
    ("stat", "vs 48-thread CPU", "×3.14", ACCENT),
    ("gap", 0.05),
    ("row", "Kernel (GPU)", "23 µs", TEXT2),
    ("row", "Host + PCIe", "44 µs", AMBER),
    ("gap", 0.2),
    ("note", "Two thirds of the frame is spent not computing."),
])

# =========================================================== 6 · OPT2
s = new_slide()
chrome(s, 6, "opt2 · streams and batching", "A CUDA stream is a queue the GPU can overlap")
bullets(s, M, 1.95, COL, [
    "A **stream** is an ordered queue of GPU work. Work in **different** streams may "
    "overlap — so a copy can run while another stream computes.",
    "Each stream gets its own **StreamContext**: device frame buffer, output buffer "
    "and pedestal. Frames are handed out **round-robin**.",
    "The host now submits **2 000 frames per call** instead of one.",
])
code(s, M, 3.62, COL, [
    "struct StreamContext {",
    "    cudaStream_t   stream;",
    "    FRAME_TYPE    *d_frame;      ClusterType *d_clusters;",
    "    PEDESTAL_TYPE *d_pd_mean, *d_pd_sum, *d_pd_sum2;",
    "};",
    "auto &sc = v_sc[frame_idx % «n_streams»];   // round-robin",
], title="ClusterFinderCUDA.hpp · per-stream state")
callout(s, M, 5.72, COL,
        "**Scaffolding, not yet the payoff.** The streams exist, but the host still "
        "synchronises after every round — see opt3.")
rail(s, [
    ("label", "opt2 · 4 streams · batch 2 000"),
    ("gap", 0.15),
    ("stat", "End to end", "43.2 µs", PALE),
    ("stat", "vs CPU", "×4.86", ACCENT),
    ("gap", 0.05),
    ("row", "Step gain over opt1", "×1.55", ACCENT),
    ("row", "Host + PCIe", "21 µs  (was 44)", AMBER),
])

# =========================================================== 7 · OPT3
s = new_slide()
chrome(s, 7, "opt3 · remove the sync barriers", "Stop draining the GPU between rounds")
bullets(s, M, 1.95, 7.4, [
    "opt2 synchronised **all streams after every round** of n_streams frames. "
    "The GPU drained to empty each time.",
    "opt3 submits every frame's H2D → kernel → D2H **asynchronously**, then "
    "synchronises **once at the end of the batch**.",
], size=10.5)
figure(s, "fig_streams", M, 3.05, 7.55)
code(s, 8.35, 1.95, 4.25, [
    "// opt2:  barrier after every round",
    "for (round) {",
    "   submit(n_streams frames);",
    "   «cudaDeviceSynchronize»();",
    "}",
    "",
    "// opt3:  submit everything, sync once",
    "for (frame : batch) {",
    "   cudaMemcpyAsync(..., sc.stream);",
    "   kernel<<<..., sc.stream>>>(...);",
    "   cudaMemcpyAsync(..., sc.stream);",
    "}",
    "for (sc : streams)",
    "   «cudaStreamSynchronize»(sc.stream);",
], size=8, title="THE ONE-LINE IDEA")
callout(s, 8.35, 5.05, 4.25,
        "**37.6 µs/frame · ×5.58**\nHost overhead 21 → **14 µs**", h=0.86, size=11)
caption(s, 8.35, 6.15, 4.25,
        "Each lane is one stream. Removing the barrier lets a stream start its next "
        "frame while its neighbours are still copying.")

# =========================================================== 8 · OPT4
s = new_slide()
chrome(s, 8, "opt4 · pinned (page-locked) memory", "What pinning is, and why the GPU cares")
bullets(s, M, 1.95, 12.0, [
    "Normal host memory is **pageable** — the OS may move or swap it. A DMA engine "
    "cannot safely read that, so the driver first copies your data into a **hidden "
    "pinned staging buffer**. Every transfer is copied twice.",
    "**Pinning** locks the pages in physical RAM. The GPU's DMA engine then reads "
    "host memory **directly** — no staging copy, and the transfer can be truly asynchronous.",
], size=10.5)
figure(s, "fig_pinning", M, 3.15, 7.6)
code(s, 8.5, 3.15, 4.1, [
    "// pin the whole dataset once",
    "«cudaHostRegister»(ptr, bytes,",
    "                 cudaHostRegisterDefault);",
    "",
    "// ... run the whole campaign ...",
    "",
    "«cudaHostUnregister»(ptr);",
], size=8, title="ClusterFinderCUDA.hpp")
callout(s, 8.5, 4.72, 4.1,
        "**27.2 µs/frame · ×7.73**\nHost overhead 14 → **3 µs**", h=0.86, size=11)
caption(s, 8.5, 5.82, 4.1,
        "Measured H2D: one 400×400 uint16 frame (312.5 KiB = 320 000 B) in 13.2 µs "
        "= 24.2 GB/s — 77% of PCIe 4.0 ×16 theoretical, i.e. true DMA speed. "
        "Pageable staging runs ~15 GB/s.")
caption(s, M, 6.62, 7.6,
        "Caveat: pinned memory is a finite system resource — it cannot be swapped. "
        "aare exposes a budget helper (print_pinning_budget) before you pin 31 GB.")

# =========================================================== 9 · OPT5
s = new_slide()
chrome(s, 9, "opt5 · CUDA Graphs", "Record the pipeline once, replay it with one call")
bullets(s, M, 1.95, 12.0, [
    "Every cudaMemcpyAsync / kernel launch costs the **CPU** a few microseconds of "
    "driver work — per frame, per operation. At 39 k frames/s that is the budget.",
    "A **CUDA Graph** captures the whole dependency DAG once. Replaying it is a "
    "**single** cudaGraphLaunch — the driver already knows every node and edge.",
], size=10.5)
figure(s, "fig_graphs", M, 3.15, 7.6)
code(s, 8.5, 3.15, 4.1, [
    "// record once, at setup",
    "cudaStreamBeginCapture(sc.stream, ...);",
    "   submit_h2d_kernel_d2h(sc);",
    "cudaStreamEndCapture(sc.stream, &sc.graph);",
    "«cudaGraphInstantiate»(&sc.graphExec, ...);",
    "",
    "// per batch — one call",
    "«cudaGraphLaunch»(sc.graphExec, sc.stream);",
], size=8, title="ClusterFinderCUDA_graph.hpp")
callout(s, 8.5, 4.88, 4.1,
        "**25.3 µs/frame · ×8.29**\nBest end-to-end result", h=0.86, size=11)
caption(s, M, 6.62, 12.0,
        "Worth 7% here because opt4 already removed the transfer cost — what remains "
        "is CPU launch overhead, which is exactly what graphs eliminate. "
        "Trade-off: shapes are frozen at record time, so the batch geometry must be fixed.")

# =========================================================== 10 · OPT6 why
s = new_slide()
chrome(s, 10, "opt6 · FP32 device pedestal", "The first change to the kernel itself")
bullets(s, M, 1.95, 7.5, [
    "~99.9% of pixels take the **pedestal-update** branch, which reads and writes "
    "mean, sum and sum². In FP64 that is **32 bytes per pixel**; in FP32, 16.",
    "The kernel is **bandwidth-bound**, so halving that traffic nearly halves the time.",
    "Second effect: on a GeForce part, **FP64 arithmetic runs at 1/64 of FP32**. "
    "The pedestal update was paying that tax on every pixel.",
], size=10.5)
figure(s, "fig_f32_kernel", M, 4.05, 7.5)
code(s, 8.5, 1.95, 4.1, [
    "// clusterfinder_kernel.cuh",
    "using COMPUTE_TYPE    = float;",
    "using DEVICE_PED_TYPE = «float»;",
    "//            was: double",
], size=8.5, title="ONE TYPEDEF")
callout(s, 8.5, 3.05, 4.1,
        "Kernel, 9×9, measured with Nsight Systems\n**43.0 µs → 25.6 µs  (−40%)**",
        h=0.92, size=11)
caption(s, 8.5, 4.20, 4.1,
        "Exclusive per-instance kernel time, 2 000 instances, 1 stream. "
        "σ = 0.4 µs. Transfers are unchanged, as they must be.")
callout(s, 8.5, 5.30, 4.1,
        "But a faster kernel is **not automatically a faster frame** — and naive FP32 "
        "is **wrong**. Both on the next slides.", h=1.05, size=10, color=AMBER)

# =========================================================== 11 · OPT6 trap
s = new_slide()
chrome(s, 11, "opt6 · the correctness trap", "Why the obvious FP32 pedestal is broken")
bullets(s, M, 1.95, 7.9, [
    "The running variance was computed as **var = E[X²] − mean²**. With a pedestal "
    "mean of ~4 655 ADU, both terms are ≈ 2.17 × 10⁷ while the answer is ≈ 2 000.",
    "In FP32 the spacing between representable numbers at 2.17 × 10⁷ is **2 048** — "
    "larger than the variance itself. This is **catastrophic cancellation**.",
], size=10.5)
figure(s, "fig_cancellation", M, 3.25, 7.5)
rail(s, [
    ("label", "What it looked like"),
    ("gap", 0.15),
    ("stat", "Extra clusters", "+28.06%", AMBER),
    ("gap", 0.05),
    ("note", "Quiet pixels got rms → 0, so their threshold became 0 and they fired "
             "on every single frame — producing a large unphysical high-energy tail "
             "in the spectrum."),
    ("gap", 0.35),
    ("row", "Affected pixels", "~1–2% of the sensor", TEXT2),
    ("row", "Written up in", "docs/pedestal_precision_…", MUTED),
])

# =========================================================== 12 · OPT6 fix
s = new_slide()
chrome(s, 12, "opt6 · the variance rewrite", "Accumulate what is small, not what is large")
bullets(s, M, 1.95, COL, [
    "Freeze a per-pixel baseline **X₀ = round(mean)** once, at the end of pedestal "
    "training, and never move it again.",
    "Accumulate the **centred** value Y = X − X₀ instead of X. Now both sums are "
    "O(rms)-sized — the huge common term is gone before the subtraction.",
    "The reported pedestal mean is still the full value, **X₀ + sum/n**, so nothing "
    "downstream changes.",
])
code(s, M, 4.0, COL, [
    "// before — both terms ~2.17e7, answer ~2000",
    "var = sum2/n - mean*mean;",
    "",
    "// after — centred on a frozen per-pixel offset X0",
    "DEVICE_PED_TYPE resid  = mean - «d_pd_off»[i];      // ~O(1)",
    "DEVICE_PED_TYPE var_px = sum2[i]/n - resid*resid;  // no cancellation",
], title="clusterfinder_kernel.cuh")
callout(s, M, 6.05, COL,
        "Result: the 100% FP32 build now matches the FP64 build to "
        "**3 × 10⁻⁷** — 70 clusters out of 233 million.")
rail(s, [
    ("label", "Why it works"),
    ("gap", 0.15),
    ("note", "Precision is relative. Floats resolve small numbers finely and large "
             "numbers coarsely — so never let a small answer be the difference of "
             "two large numbers."),
    ("gap", 0.5),
    ("row", "f32 vs f64 counts", "3 × 10⁻⁷", ACCENT),
    ("row", "vs CPU", "0.0039%", ACCENT),
    ("gap", 0.3),
    ("note", "X₀ must never be updated — the accumulators are defined relative to it."),
])

# =========================================================== 13 · OPT6 when
s = new_slide()
chrome(s, 13, "opt6 · when does a faster kernel help?", "Only if the kernel was the tallest bar")
figure(s, "fig_bottleneck", M, 2.05, 11.9)
callout(s, M, 5.35, 11.9,
        "**f32 buys exactly the distance between the kernel and the next-tallest bar.** "
        "At 3×3 the kernel already hides inside the transfers, so end-to-end throughput does not move at all. "
        "At 9×9 the kernel is on the critical path, and the same change is worth 8% of the frame.",
        h=0.95, size=11)
caption(s, M, 6.55, 11.9,
        "9×9, 20 000 frames, 8 streams, cap 1 500 (all CPU clusters recorded): "
        "1.540 s → 1.423 s batched, 1.482 s → 1.372 s with graphs. "
        "The FP64 build shows the classic kernel-bound signature — per-frame wall time shorter than per-frame kernel time, because kernels from different streams queue.")

# =========================================================== 14 · RESULTS
s = new_slide()
chrome(s, 14, "Results", "The whole ladder, one dataset, one baseline")
figure(s, "fig_arc", M, 1.95, 11.9)
callout(s, M, 5.55, 5.85,
        "**×8.29 over 48 CPU threads** — 21.0 s → 2.53 s for 100 000 frames.", h=0.8)
callout(s, 6.75, 5.55, 5.85,
        "Every step is **monotonic**, and correctness is held constant at **0.004%** throughout.",
        h=0.8, color=AMBER)
caption(s, M, 6.55, 11.9,
        "3×3 clusters · nσ = 5 · 100 000 frames · batch 2 000 · 4 streams · warm run "
        "(page faults plateaued) · CPU baseline = ClusterFinderMT with 48 threads.")

# =========================================================== 15 · WHERE TIME GOES
s = new_slide()
chrome(s, 15, "Where the time actually went", "The kernel never changed — the overhead collapsed")
figure(s, "fig_overhead", M, 2.05, 5.9)
bullets(s, 7.0, 2.15, 5.6, [
    "The GPU kernel is a **flat ~23 µs** across opt1–opt5. Not one of those steps "
    "made the arithmetic faster.",
    "What changed is everything around it: **44 → 21 → 14 → 3 µs** of host and PCIe "
    "time per frame.",
    "By opt4 the pipeline is **fed almost perfectly** — which is precisely why opt5 "
    "(launch overhead) is the only lever left, and why opt6 shows nothing at 3×3.",
], size=10.5)
callout(s, 7.0, 5.15, 5.6,
        "The bottleneck moved from **the host**, to **PCIe**, and finally — only for "
        "large cluster windows — to **the kernel**.", h=0.95, size=10.5)
caption(s, M, 6.4, 11.9,
        "Per-frame GPU operation profile at 9×9 (Nsight Systems): kernel 43 µs · "
        "D2H 19.4 µs · H2D 13.2 µs. At 3×3 the kernel is 13–24 µs against a ~25 µs "
        "transfer-and-host floor — which is the entire story of this deck in three numbers.")

# =========================================================== 16 · CORRECTNESS
s = new_slide()
chrome(s, 16, "Validation", "Same physics out of every variant")
figure(s, "fig_correctness", M, 2.0, 7.6)
bullets(s, 8.6, 2.05, 4.1, [
    "233 million clusters over 100 000 frames.",
    "All CUDA variants agree with the CPU to **0.004%**.",
    "The residual is **not** precision: it is the CUDA finder updating the pedestal "
    "**once per frame** vs the CPU's per-pixel update.",
], size=10)
code(s, M, 4.85, 7.6, [
    "CPU  (ClusterFinderMT)   233 085 343     2330.85 / frame     reference",
    "opt1 .. opt5  (f64 ped)  233 093 484     2330.93 / frame     0.0035 %",
    "opt6          (f32 ped)  233 094 465     2330.94 / frame     0.0039 %",
], size=8.5, title="CLUSTER COUNTS · 3×3 · 100 000 FRAMES")
callout(s, 8.6, 4.85, 4.1,
        "A local-maximum gate had to be **back-ported** into the opt1/opt2 snapshot so "
        "the whole ladder is compared at identical correctness.", h=1.0, size=10, color=AMBER)
caption(s, M, 6.35, 11.9,
        "Cross-checks: ClusterFinderFrozen (a CPU finder with the CUDA pedestal-update "
        "timing) isolates that residual; energy spectra overlay within statistics.")

# =========================================================== 17 · API 1
s = new_slide()
chrome(s, 17, "For users · Python API", "The fast path in eight lines")
code(s, M, 1.95, 7.6, [
    "from aare import File, ClusterFinderCUDA",
    "",
    "cf = ClusterFinderCUDA(image_size=(400, 400), cluster_size=(3, 3),",
    "                       n_sigma=5, «n_streams»=4,",
    "                       «max_clusters_per_frame»=3000)",
    "",
    "for _ in range(1000):                    # 1. train the pedestal",
    "    cf.push_pedestal_frame(pd.read_frame())",
    "",
    "data = f.read_n(100_000)                 # 2. one contiguous array",
    "cf.«register_input_buffer»(data)           # 3. pin it once",
    "",
    "for s in range(0, N, 2000):              # 4. batch through it",
    "    clusters = cf.«find_clusters_batched»(data[s:s+2000], first_frame=s)",
    "",
    "cf.unregister_input_buffer()             # 5. release the pages",
], size=9, title="THE RECOMMENDED PATTERN")
bullets(s, 8.6, 2.0, 4.1, [
    "find_clusters_batched returns **one ClusterVector per frame**, in order.",
    "register_input_buffer is what turns opt3 into opt4 — **one call**.",
    "Swap in ClusterFinderCUDAGraph for opt5; the API is identical.",
], size=10)
callout(s, 8.6, 4.55, 4.1,
        "GIL is released for both find_clusters and find_clusters_batched, so "
        "reading the next file can overlap with the GPU.", h=1.05, size=10)
callout(s, 8.6, 5.80, 4.1,
        "Pin **once**, outside the loop. Slices of a registered array inherit the pinning.",
        h=0.85, size=10, color=AMBER)

# =========================================================== 18 · API 2
s = new_slide()
chrome(s, 18, "For users · choosing the knobs", "What to set, and what it costs you")
hdr = [("Parameter", 1.05), ("What it does", 3.6), ("Guidance", 5.2)]
y = 2.0
rect(s, M, y, 11.9, 0.4, PANEL)
for lab, dx in hdr:
    tf = tb(s, M + dx - 0.85 if dx > 1.05 else M + 0.28, y + 0.09, 5.0, 0.3)
    run(para(tf, True), lab.upper(), 9, MUTED, bold=True, spc=1.3)
y += 0.44
params = [
    ("n_streams", "How many frames are in flight at once.",
     "4 is right for 3×3. Larger windows saturate the GPU — 8 helps at 9×9."),
    ("max_clusters_per_frame", "Fixed size of the per-frame D2H transfer.",
     "Must exceed the real maximum or clusters are silently dropped. Too high wastes PCIe."),
    ("batch size", "Frames per find_clusters_batched call.",
     "2 000 amortises launch overhead without a large pinned footprint."),
    ("cluster_size", "Compile-time stencil geometry.",
     "3×3 and 9×9 are registered; 9×9 shifts the bottleneck onto the kernel."),
    ("register_input_buffer", "Page-locks the host array for DMA.",
     "Always, if the data is already in RAM. Check the pinning budget first."),
]
for i, (p_, what, guide) in enumerate(params):
    if i % 2 == 0:
        rect(s, M, y, 11.9, 0.82, PANEL)
    tf = tb(s, M + 0.28, y + 0.14, 2.6, 0.5)
    run(para(tf, True, line=1.1), p_, 9.5, ACCENT, font=MONO, bold=True)
    tf = tb(s, M + 3.0, y + 0.14, 2.9, 0.6)
    run(para(tf, True, line=1.2), what, 9.5, PALE)
    tf = tb(s, M + 6.15, y + 0.14, 5.4, 0.6)
    run(para(tf, True, line=1.2), guide, 9.5, TEXT2)
    y += 0.80
callout(s, M, 6.55, 11.2,
        "The single most common mistake: leaving **max_clusters_per_frame** too low. "
        "It does not error — it truncates, and every frame quietly returns the same count.",
        h=0.66, size=10, color=AMBER)

# =========================================================== 19 · NEXT
s = new_slide()
chrome(s, 19, "Where this leaves us", "The bottleneck has moved — twice")
cards = [
    ("DONE", ACCENT, "×8.29 over 48 CPU threads",
     "25.3 µs/frame end to end, at 0.004% cluster agreement. Five pipeline steps "
     "and one kernel step."),
    ("DONE", ACCENT, "FP32 pedestal, safely",
     "−40% kernel, and correct — because the variance is now accumulated on a frozen "
     "per-pixel offset instead of a raw second moment."),
    ("NEXT", AMBER, "Attack the transfers, not the kernel",
     "At 3×3 the kernel is already hidden. The remaining per-frame cost is PCIe: "
     "492 kB of clusters out, 312 kB of frame in."),
    ("NEXT", AMBER, "Keep results on the device",
     "On-GPU reduction, eta/interpolation on device, or compressed cluster formats — "
     "so the D2H bar stops setting the floor."),
]
x, y = M, 2.05
for i, (tag, col, title, body) in enumerate(cards):
    cx = M + (i % 2) * 6.05
    cy = 2.05 + (i // 2) * 2.35
    rect(s, cx, cy, 5.85, 2.05, PANEL)
    rect(s, cx, cy, 5.85, 0.035, col)
    tf = tb(s, cx + 0.3, cy + 0.26, 1.4, 0.26)
    run(para(tf, True), tag, 8.5, col, bold=True, spc=1.5)
    tf = tb(s, cx + 0.3, cy + 0.60, 5.2, 0.4)
    run(para(tf, True, line=1.1), title, 14, PALE, bold=True)
    tf = tb(s, cx + 0.3, cy + 1.12, 5.2, 0.85)
    run(para(tf, True, line=1.3), body, 10, TEXT2)
callout(s, M, 6.58, 11.2,
        "Full numbers, methodology and reproduction steps: **docs/benchmark_opt1_opt6_results.md** · "
        "notebook **python/tests/ClusterFinderCUDA_perf.ipynb** · profiler probe **python/tests/nsys_kernel_probe.py**",
        h=0.66, size=9.5)

prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
