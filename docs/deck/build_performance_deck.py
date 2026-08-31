"""Build docs/cf_cuda_performance.pptx — the algorithm + kernel + hardware half of
docs/cf_cuda_kernel.pptx fused with the opt1→opt7 optimization story.

The ladder is told in three acts, ordered by which bar is tallest:

    ACT I   [f64]  feed the GPU        opt1 opt2 opt3 opt4   (+ route A, rejected)
    ACT II  [f64]  get results back    opt5 opt6             (+ routes B', B'', rejected)
    ACT III [f32]  the kernel          opt7

Act III comes last because it cannot be justified earlier: measured through
collect() the f32 kernel is worth 1.5 % end-to-end, and only once the host is off
the critical path is the same change worth 21 %.

cf_cuda_kernel.pptx is kept solely as the base presentation — it donates the PSI
theme and title slide, and every other slide of it is deleted below.

All hardware numbers are re-measured against the CURRENT kernel, not taken
from the kernel deck (whose implementation details and timings are stale):

    nvcc -arch=sm_89 --ptxas-options=-v          -> registers, spills
    cudaOccupancyMaxActiveBlocksPerMultiprocessor -> blocks/SM, occupancy

    3x3 : 34 regs/thread, 0 spill, 1296 B smem, 6 blocks/SM, 100.0% occupancy
    9x9 : 128 regs/thread, 0 spill, 2304 B smem, 2 blocks/SM,  33.3% occupancy
    32x32 blocks @ 9x9: 0 blocks/SM -- 1024 x 128 regs > 65536 regs/SM

Performance numbers come from docs/ClusterFinderCUDA_benchmark_results.md (quotable
rows only).
"""
import re

from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pathlib import Path
from PIL import Image

DOCS = Path(__file__).resolve().parent.parent
FIGS = DOCS / "figures"
BASE = DOCS / "cf_cuda_kernel.pptx"          # PSI theme + title slide
OUT = DOCS / "cf_cuda_performance.pptx"

# ---------------------------------------------------------------- design tokens
BG     = RGBColor(0x0B, 0x10, 0x18)
PANEL  = RGBColor(0x12, 0x1A, 0x28)
# Code panels are framed, not filled, to set them apart. The fill can only ever
# be a hair lighter than the slide (1.16:1 at #17202E) before the accent-coloured
# tokens inside start losing contrast against it -- and a projector's black level
# crushes every dark tone together anyway. The 1 pt edge reaches 2.18:1 against
# the background without touching the ink's ground, so the boundary is drawn by
# the line and the fill only has to say "a different surface".
CODEBG = RGBColor(0x17, 0x20, 0x2E)
CODEEDGE = RGBColor(0x3A, 0x4C, 0x66)
# MUTED was tuned against the slide background; on the lighter code fill it drops
# to 3.75:1, and it carries the comments and the panel title. This puts them back
# above MUTED's original 4.22 -- inside code panels only.
CODEDIM = RGBColor(0x7C, 0x8A, 0x9E)
RULE   = RGBColor(0x1E, 0x28, 0x36)
ACCENT = RGBColor(0x1E, 0x90, 0xC2)
AMBER  = RGBColor(0xE8, 0xB2, 0x5C)
PALE   = RGBColor(0xE7, 0xED, 0xF4)
TEXT2  = RGBColor(0xA5, 0xB2, 0xC4)
MUTED  = RGBColor(0x6B, 0x7A, 0x90)
CARD   = RGBColor(0xF4, 0xF6, 0xF9)          # light card for white figures
RED    = RGBColor(0xE2, 0x54, 0x54)          # pointing; the s1 label (s21/A6);
                                             # the ⟪…⟫ code marker (s33)
GREEN  = RGBColor(0x5C, 0xC8, 0xA0)          # rooflines / floors, as in make_figs

UI, MONO = "Segoe UI", "Consolas"
W, H = 13.333, 7.5
M = 0.50                      # left margin
COL = 8.16                    # left column width
RAIL_X, RAIL_W = 9.28, 3.62   # right rail

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

prs = Presentation(str(BASE))
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[0]                  # 'Blank Slide' — zero shapes
N_SLIDES = 34


# ------------------------------------------------------------------- helpers
def keep_only_slide(prs, keep=0):
    """Drop every slide but one from the base presentation."""
    lst = prs.slides._sldIdLst
    for i, sldId in enumerate(list(lst)):
        if i != keep:
            prs.part.drop_rel(sldId.get(f"{R}id"))
            lst.remove(sldId)


def set_para_texts(shape, texts):
    """Replace paragraph texts in-place, keeping each paragraph's formatting."""
    for p, txt in zip(shape.text_frame.paragraphs, texts):
        if not p.runs:
            continue
        p.runs[0].text = txt
        for r in p.runs[1:]:
            r.text = ""


def new_slide():
    """A dark slide on the PSI master.

    Two independent guards, because the base template's master carries a PSI
    background picture and logo that must not bleed through:
      1. showMasterSp="0" + a slide-level <p:bg> (correct schema position:
         first child of <p:cSld>), which is what PowerPoint honours;
      2. a full-bleed rectangle as the first shape, which every renderer
         honours regardless of how it treats (1).
    """
    s = prs.slides.add_slide(BLANK)
    s._element.set("showMasterSp", "0")

    cSld = s._element.find(f"{P}cSld")
    bg = etree.Element(f"{P}bg")
    pr = etree.SubElement(bg, f"{P}bgPr")
    fill = etree.SubElement(pr, f"{A}solidFill")
    clr = etree.SubElement(fill, f"{A}srgbClr")
    clr.set("val", "0B1018")
    etree.SubElement(pr, f"{A}effectLst")
    cSld.insert(0, bg)

    rect(s, 0, 0, W, H, BG)
    return s


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, first=False, space_after=0, space_before=0, line=None, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    if line: p.line_spacing = line
    if align: p.alignment = align
    return p


# The projection floor, in points, for text set directly in PowerPoint.
#
# The deck is shown in a room where the back row is 6-7 m from the screen. On a
# 13.33 x 7.5 in slide, 9 pt is about 1/60 of the slide height, which is the
# conventional lower bound for supporting detail. It is deliberately the SAME
# number that make_figs.py checks every string inside every PNG against, so the
# deck has one floor rather than two, and nothing is legible on the slide but
# not inside the picture next to it.
#
# It is enforced here, in run(), rather than at the ~400 call sites, because a
# floor applied by hand is a floor that one new caption silently drops through.
MIN_PT = 9.0

# ------------------------------------------------------------ the type scale
# One place, because 186 scattered `size=` arguments cannot be reasoned about.
# Sized for a 13.33 x 7.5 in slide read at 6-7 m: BODY is the argument, and
# everything else is explicitly subordinate to it. Raising BODY costs text --
# 15 pt holds roughly 45 % of the characters 10.5 pt did in the same box -- so
# the scale is a budget, not a preference. Figures are gated separately, in
# make_figs.py, and their floor moves with BODY: type inside a plot that is much
# smaller than the body text reads as an afterthought from the back of a room.
PT_BODY  = 15      # bullets: what the slide is claiming
PT_LEAD  = 13      # callouts: the sentence to remember, already boxed and bold
PT_TABLE = 11.5      # table cells
PT_RAIL  = 13      # rail row values
PT_META  = 11    # captions and rail notes: provenance, deliberately quieter
PT_CODE  = 10.5      # code panels -- read by token, not word by word
PT_API   = 11.5      # the one code panel a user will COPY: slide 32
PT_LABEL = 9.5      # small-caps labels on rails, tables, statstrips
# PT_API is the single deliberate exception to "code is 10.5". Slide 34 is the
# only slide whose code is the deliverable rather than the evidence -- it is the
# eight lines a user types -- and it is also the only code slide with an inch and
# a half of unused height. The longest line, at 72 mono characters, still clears
# the panel's 7.49 in of inner width at this size, with room to spare.


def run(p, text, size=11, color=TEXT2, font=UI, bold=False, italic=False,
        spc=None, sup=False):
    r = p.add_run(); r.text = text
    f = r.font
    f.name, f.size, f.bold, f.italic = font, Pt(max(size, MIN_PT)), bold, italic
    f.color.rgb = color
    if spc is not None:
        r.font._rPr.set("spc", str(int(spc * 100)))
    if sup:
        # python-pptx has no superscript API. `baseline` is the OOXML attribute,
        # in thousandths of a percent; both PowerPoint and LibreOffice honour it.
        r.font._rPr.set("baseline", str(sup if sup is not True else 30000))
    return r


# ------------------------------------------------------------------ chrome
def chrome(s, idx, eyebrow, title, title_size=27, opt=None):
    """Standard slide chrome.

    `opt` marks a rung of the optimisation ladder, in the same corner and the
    same way the annex marks its groups: a slide either IS about an optimisation
    or it is not, and that should be visible before the title is read. Several
    rungs span two slides (opt3, opt5, opt7), so the badge repeats -- it names
    the step, not the slide.
    """
    if opt:
        tf = tb(s, M, 0.44, 1.2, 0.34)
        run(para(tf, True), f"OPT{opt}", 15, ACCENT, bold=True)
        tf = tb(s, 1.45, 0.50, 10.25, 0.32)
    else:
        rect(s, M, 0.60, 0.35, 0.035, ACCENT)
        tf = tb(s, 1.00, 0.50, 10.70, 0.32)
    run(para(tf, True), eyebrow.upper(), 9, MUTED, bold=True, spc=1.6)

    tf = tb(s, M, 0.86, 12.29, 1.0)
    run(para(tf, True, line=1.05), title, title_size, PALE, bold=True)

    span, n = 11.37, N_SLIDES
    pitch = span / n; wseg = pitch * 0.90
    for i in range(n):
        rect(s, M + i * pitch, 7.28, wseg, 0.045, ACCENT if i <= idx - 1 else RULE)
    tf = tb(s, 12.29, 7.14, 0.9, 0.3)
    run(para(tf, True, align=PP_ALIGN.RIGHT), f"{idx} / {n}", 8.5, MUTED)


N_ANNEX = 12         # annex GROUPS, not slides


def annex_chrome(s, grp, eyebrow, title, part=None, nparts=None, title_size=27):
    """Same chrome, amber, on its own progress track.

    The annex is GROUPED, not enumerated: slides that belong together carry the
    same tag (A2 is both CUDA-Graphs slides, A5 all three artefacts), because
    what a reader needs at a glance is which slides form one argument. The badge
    is rendered in the header, not only in the corner.
    """
    tf = tb(s, M, 0.44, 1.0, 0.34)
    run(para(tf, True), f"A{grp}", 15, AMBER, bold=True)
    tf = tb(s, 1.25, 0.50, 10.45, 0.32)
    tag = f"ANNEX · {eyebrow}" + (f" · {part} of {nparts}" if nparts else "")
    run(para(tf, True), tag.upper(), 9, MUTED, bold=True, spc=1.6)
    tf = tb(s, M, 0.86, 12.29, 1.0)
    run(para(tf, True, line=1.05), title, title_size, PALE, bold=True)
    pitch = 11.37 / N_ANNEX
    for i in range(N_ANNEX):
        rect(s, M + i * pitch, 7.28, pitch * 0.90, 0.045,
             AMBER if i <= grp - 1 else RULE)
    tf = tb(s, 11.99, 7.14, 1.2, 0.3)
    foot = f"A{grp}" + (f" · {part}/{nparts}" if nparts else "")
    run(para(tf, True, align=PP_ALIGN.RIGHT), foot, 8.5, MUTED)


def table(s, x, y, w, header, rows, colw, size=PT_TABLE, rowh=0.74):
    """Minimal header + zebra table. colw are fractions of w."""
    xs, acc = [], 0.0
    for c in colw:
        xs.append(x + acc * w)
        acc += c
    rect(s, x, y, w, 0.34, PANEL)
    for cx, h in zip(xs, header):
        tf = tb(s, cx + 0.16, y + 0.08, w, 0.24)
        run(para(tf, True), _up(h), PT_LABEL, MUTED, bold=True, spc=1.2)
    yy = y + 0.38
    for i, row in enumerate(rows):
        if i % 2 == 0:
            rect(s, x, yy, w, rowh, PANEL)
        for j, (cx, cell) in enumerate(zip(xs, row)):
            col = PALE if j == 0 else TEXT2
            tf = tb(s, cx + 0.16, yy + 0.10, colw[j] * w - 0.24, rowh)
            p = para(tf, True, line=1.2)
            for k, part in enumerate(str(cell).split("**")):
                if part:
                    run(p, part, size, AMBER if k % 2 else col, bold=bool(k % 2))
        yy += rowh + 0.04
    return yy


def bullets(s, x, y, w, items, size=PT_BODY, gap=9):
    tf = tb(s, x, y, w, 0.3)
    for i, it in enumerate(items):
        color, txt = (it if isinstance(it, tuple) else (TEXT2, it))
        p = para(tf, i == 0, space_after=gap, line=1.25)
        run(p, "•  ", size, MUTED)
        for j, part in enumerate(txt.split("**")):
            if part:
                run(p, part, size, PALE if j % 2 else color, bold=bool(j % 2))
    return tf


# Three highlights, because a code panel has three kinds of thing worth pointing
# at: «…» is ACCENT, for whatever the slide is arguing about (a knob, a step
# number); ‹…› is PALE bold, reserved for the API surface -- the call the
# audience will actually type; ⟪…⟫ is RED, for the value that DISAGREES with the
# accented one, which so far is only slide 28's f64 rival. They must not share a
# colour or the panel says everything twice.
CODE_MARK = re.compile(r"«([^»]*)»|‹([^›]*)›|⟪([^⟫]*)⟫")


def code(s, x, y, w, lines, size=PT_CODE, title=None):
    # Line height has to follow the font size, and the constant has to match what
    # the RENDERER does, not what python-pptx assumes. Measured off a LibreOffice
    # render at 9 pt Consolas with line=1.12: 0.170 in per line, i.e. 0.0189 in
    # per point. The old 0.0174 under-counted by 8.6 %, which is invisible on a
    # short panel and clips the last line or two on a long one.
    size = max(size, MIN_PT)
    lh = 0.0189 * size
    h = 0.24 + len(lines) * lh + (0.24 if title else 0)
    box = rect(s, x, y, w, h, CODEBG, MSO_SHAPE.ROUNDED_RECTANGLE)
    box.line.fill.solid()             # rect() cleared the line; put one back
    box.line.fill.fore_color.rgb = CODEEDGE
    box.line.width = Pt(1.0)
    box.adjustments[0] = 0.055        # a small radius: a panel, not a pill
    ty = y + 0.12
    if title:
        tf = tb(s, x + 0.18, ty, w - 0.36, 0.2)
        run(para(tf, True), title, PT_LABEL, CODEDIM, bold=True, spc=1.2)
        ty += 0.24
    tf = tb(s, x + 0.18, ty, w - 0.36, h - 0.24)
    for i, ln in enumerate(lines):
        p = para(tf, i == 0, line=1.12)
        if ln.strip().startswith(("//", "#")):
            run(p, ln, size, CODEDIM, MONO)
            continue
        pos = 0
        for m in CODE_MARK.finditer(ln):
            if m.start() > pos:
                run(p, ln[pos:m.start()], size, TEXT2, MONO)
            if m.group(1) is not None:          # «…»  the thing to look at
                run(p, m.group(1), size, ACCENT, MONO, bold=True)
            elif m.group(2) is not None:        # ‹…›  the API surface itself
                run(p, m.group(2), size, PALE, MONO, bold=True)
            else:                               # ⟪…⟫  the value that disagrees
                run(p, m.group(3), size, RED, MONO, bold=True)
            pos = m.end()
        if pos < len(ln):
            run(p, ln[pos:], size, TEXT2, MONO)
    return h


def callout(s, x, y, w, text, h=0.78, color=ACCENT, size=PT_LEAD):
    rect(s, x + 0.045, y, w - 0.045, h, PANEL)
    rect(s, x, y, 0.045, h, color)
    tf = tb(s, x + 0.28, y + 0.10, w - 0.5, h - 0.2, MSO_ANCHOR.MIDDLE)
    p = para(tf, True, line=1.2)
    for j, part in enumerate(text.split("**")):
        if part:
            run(p, part, size, PALE if j % 2 else TEXT2, bold=bool(j % 2))


def _up(txt):
    """upper() for labels, but 'µ'.upper() is Greek capital Mu — which renders as
    an 'M' and turns 'µs' into 'MS', i.e. microseconds into milliseconds."""
    return txt.upper().replace("\u039c", "µ")


def rail(s, items, y0=2.0, divider=True):
    if divider:
        rect(s, 9.03, 2.0, 0.012, 4.55, RULE)
    y = y0
    for it in items:
        kind = it[0]
        if kind == "label":
            tf = tb(s, RAIL_X, y, RAIL_W, 0.26)
            run(para(tf, True), _up(it[1]), PT_LABEL, MUTED, bold=True, spc=1.4)
            y += 0.32
        elif kind == "stat":
            _, lab, val, col = it
            tf = tb(s, RAIL_X, y, RAIL_W, 0.24)
            run(para(tf, True), _up(lab), PT_LABEL, MUTED, spc=1.2)
            tf = tb(s, RAIL_X, y + 0.24, RAIL_W, 0.6)
            run(para(tf, True), val, 26, col, bold=True)
            y += 1.02
        elif kind == "row":
            _, lab, val, col = it
            tf = tb(s, RAIL_X, y, RAIL_W, 0.24)
            run(para(tf, True), _up(lab), PT_LABEL, MUTED, spc=1.2)
            tf = tb(s, RAIL_X, y + 0.24, RAIL_W, 0.32)
            run(para(tf, True), val, PT_RAIL, col, bold=True)
            y += 0.72
        elif kind == "note":
            tf = tb(s, RAIL_X, y, RAIL_W, 0.9)
            run(para(tf, True, line=1.25), it[1], PT_META, TEXT2)
            y += 0.30 + 0.21 * (len(it[1]) // 34 + 1)
        elif kind == "gap":
            y += it[1]
    return y


def figure(s, name, x, y, w):
    p = FIGS / f"{name}.png"
    iw, ih = Image.open(p).size
    h = w * ih / iw
    s.shapes.add_picture(str(p), In(x), In(y), In(w), In(h))
    return h


def card_figure(s, name, x, y, w, pad=0.10):
    """A light-background figure (imported, not re-rendered) on a light card."""
    p = FIGS / f"{name}.png"
    iw, ih = Image.open(p).size
    h = w * ih / iw
    rect(s, x - pad, y - pad, w + 2 * pad, h + 2 * pad, CARD,
         MSO_SHAPE.ROUNDED_RECTANGLE)
    s.shapes.add_picture(str(p), In(x), In(y), In(w), In(h))
    return h + 2 * pad


def frame_rect(s, x, y, w, h, color=RED, wpt=1.75):
    """An outline, not a fill: used to point at one row without recolouring it."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = Pt(wpt)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.06
    return sh


def notes(s, text):
    """Speaker notes. Detail that belongs in the talk, not on the slide.

    python-pptx creates the notes slide on first access, so this is safe to call
    on any slide. Used to relieve slides that carry a figure: the mechanism goes
    on the screen, the API detail goes here.
    """
    s.notes_slide.notes_text_frame.text = text


def caption(s, x, y, w, text, size=PT_META):
    tf = tb(s, x, y, w, 0.3)
    run(para(tf, True, line=1.25), text, size, MUTED)


def flow(s, x, y, w, steps, h=0.62):
    """Numbered left-to-right step strip."""
    n = len(steps); gap = 0.30
    bw = (w - gap * (n - 1)) / n
    for i, t in enumerate(steps):
        bx = x + i * (bw + gap)
        rect(s, bx, y, bw, h, PANEL)
        rect(s, bx, y, 0.03, h, ACCENT)
        tf = tb(s, bx + 0.20, y + 0.05, bw - 0.32, h - 0.10, MSO_ANCHOR.MIDDLE)
        p = para(tf, True, line=1.1)
        run(p, f"{i + 1}  ", PT_LABEL, ACCENT, bold=True, font=MONO)
        run(p, t, 12.5, PALE)
        if i < n - 1:
            tf = tb(s, bx + bw, y + 0.05, gap, h - 0.10, MSO_ANCHOR.MIDDLE)
            run(para(tf, True, align=PP_ALIGN.CENTER), "›", 15, MUTED, bold=True)


def statstrip(s, x, y, w, items, h=0.80):
    """A row of labelled numbers.

    An item may be (label, value) or (label, value, outline_colour). The third
    form rings the box -- same device as frame_rect on the validation table --
    so a strip can say which of its numbers belong together without recolouring
    the type, which would fight the label/value hierarchy the strip is built on.
    """
    n = len(items); gap = 0.22
    bw = (w - gap * (n - 1)) / n
    for i, it in enumerate(items):
        lab, val, ring = (it if len(it) == 3 else (*it, None))
        bx = x + i * (bw + gap)
        rect(s, bx, y, bw, h, PANEL)
        if ring:
            frame_rect(s, bx, y, bw, h, color=ring, wpt=1.5)
        tf = tb(s, bx + 0.20, y + 0.10, bw - 0.4, 0.22)
        run(para(tf, True), lab.upper(), PT_LABEL, MUTED, bold=True, spc=1.2)
        tf = tb(s, bx + 0.20, y + 0.42, bw - 0.4, 0.34)
        run(para(tf, True), val, 17, PALE, bold=True)


# ------------------------------------------------------------ interstitial
_NARROW, _WIDE = set("ijlt.,;:'!|()[]"), set("mwMW")


def _em(txt):
    """Width of `txt` in ems of Segoe UI Bold, near enough to wrap by.

    A character count is not good enough: 'now the tallest bar' and
    'measured, and in what' are the same length and differ by 8 % in width, which
    is exactly the margin that decides whether a title takes two lines or three.
    Weights are calibrated against a LibreOffice render of two title lines and
    reproduce both to within 1 %.
    """
    return sum(1.08 if c in _WIDE else
               0.37 if c in _NARROW else
               0.34 if c == " " else 0.68 for c in txt)


def _fit(title, w, sizes=(36, 32, 28, 24), lines=2):
    """Largest size at which `title` wraps into at most `lines` lines across `w`.

    python-pptx cannot measure text and PowerPoint's autofit does not apply until
    a render, so a title that grows one line silently walks over the rule beneath
    it. No tolerance is granted — a title that only just fits is one font
    substitution away from not fitting on someone else's machine.
    """
    for size in sizes:
        cap = w / (size / 72)
        n, cur = 1, 0.0
        for word in title.split():
            need = _em(word) + (0.34 if cur else 0)
            if cur + need > cap and cur:
                n, cur = n + 1, _em(word)
            else:
                cur += need
        if n <= lines:
            return size
    return sizes[-1]


def section(kicker, title, thesis, items, rng, col=ACCENT, carry=None,
            annex=False, item_pt=15, item_gap=0.20):
    """An unnumbered beat between sections: where we are, and what is coming.

    Deliberately sparse — it exists to buy 10–15 s of stage setting, so it has
    to be readable at a glance and finished before the audience starts reading
    ahead. It carries no slide number and takes no tick of its own on the
    progress track: slides 3–34 keep the numbers they have, so the annex's
    cross-references ("expands slide 24") stay true. What it lights up instead
    is the *range* the section covers, which is the thing the audience wants.
    """
    s = new_slide()
    rect(s, 0, 0, 0.16, H, col)

    tf = tb(s, M + 0.31, 1.52, 5.99, 0.3)
    run(para(tf, True), kicker.upper(), 10.5, MUTED, bold=True, spc=1.8)
    tf = tb(s, M + 0.31, 1.90, 6.2, 1.35)
    run(para(tf, True, line=1.03), title, _fit(title, 6.2), PALE, bold=True)
    rect(s, M + 0.31, 3.42, 1.55, 0.03, col)
    tf = tb(s, M + 0.31, 3.70, 5.89, 1.6)
    run(para(tf, True, line=1.35), thesis, 16, TEXT2)

    if carry:
        lab, val, sub = carry
        rect(s, M + 0.31, 5.55, 5.89, 1.12, PANEL)
        rect(s, M + 0.31, 5.55, 0.035, 1.12, col)
        tf = tb(s, M + 0.62, 5.70, 5.37, 0.24)
        run(para(tf, True), _up(lab), 9.5, MUTED, bold=True, spc=1.4)
        tf = tb(s, M + 0.62, 5.92, 5.37, 0.4)
        run(para(tf, True), val, 24, col, bold=True)
        tf = tb(s, M + 0.62, 6.38, 5.37, 0.26)
        run(para(tf, True), sub, 11, MUTED)

    # The list is the reason the divider exists, so it is set at body size, not
    # at caption size. That costs height, so each entry is measured rather than
    # given a fixed step: a two-line entry gets two lines of room instead of
    # silently sitting on the one below it.
    rect(s, 7.17, 1.95, 0.012, 4.4, RULE)
    # TXT_W 4.58, not 4.85, and the text column starts at 8.32 rather than 8.05:
    # the number is set in MONO, so a range label — "25–26", "28–29", "32–33" —
    # is five characters at 15 pt = 0.63 in and ran to within a hair of the entry
    # it numbers. No entry in the deck wraps at 4.58, so the gutter is free.
    ITEM_PT, TXT_W, TXT_X = item_pt, 4.58, 8.32
    cap = TXT_W / (ITEM_PT / 72)
    nlines = [max(1, -(-_em(t.replace("**", "")) // cap)) for _, t in items]
    step = ITEM_PT * 1.45 / 72
    total = sum(n * step for n in nlines) + item_gap * (len(items) - 1)
    y = 1.95 + (4.4 - total) / 2
    tf = tb(s, 7.47, y - 0.46, 5.17, 0.26)
    run(para(tf, True), "COMING UP", 10, MUTED, bold=True, spc=1.6)
    for (num, txt), nl in zip(items, nlines):
        tf = tb(s, 7.47, y, 0.78, 0.3)
        run(para(tf, True), str(num), ITEM_PT, col, bold=True, font=MONO)
        tf = tb(s, TXT_X, y, TXT_W, 0.3)
        p = para(tf, True, line=1.22)
        for j, part in enumerate(txt.split("**")):
            if part:
                run(p, part, ITEM_PT, col if j % 2 else PALE, bold=bool(j % 2))
        y += nl * step + item_gap

    # The annex divider sits on the annex's own track: the main arc is finished
    # behind it, so lighting main-track segments would misreport where we are.
    n_track = N_ANNEX if annex else N_SLIDES
    pitch = 11.37 / n_track
    for i in range(n_track):
        n = i + 1
        # the section ahead in its own colour, what is already behind us dimmed,
        # the rest dark — so the divider agrees with the chrome on either side.
        c = col if annex or rng[0] <= n <= rng[1] else (
            MUTED if n < rng[0] else RULE)
        rect(s, M + i * pitch, 7.28, pitch * 0.90, 0.045, c)
    return s


# =========================================================== 1 · PSI TITLE
keep_only_slide(prs, 0)
title_slide = prs.slides[0]
by_name = {sh.name: sh for sh in title_slide.shapes}
set_para_texts(by_name["TextShape 1"], ["The CUDA ClusterFinder"])
set_para_texts(by_name["CustomShape 4"],
               ["Kernel design · hardware limits · seven optimization steps"])
set_para_texts(by_name["CustomShape 2"], ["Khalil Daniel Ferjaoui"])
set_para_texts(by_name["CustomShape 3"],
               ["Paul Scherrer Institut · aare", "September 2026"])
notes(title_slide, """SAY THIS -- 52 words, about 20 s.

- Good morning. I am Khalil, and this is about making the cluster finder in aare
  run on a GPU.

- It is a performance talk, so there are numbers, but the story is simple: I
  expected to spend my time on the maths, and I spent it on moving data.""")

# =========================================================== 2 · HERO
s = new_slide()
rect(s, 0, 0, 0.17, H, ACCENT)
tf = tb(s, M + 0.31, 0.85, 11.37, 0.3)
run(para(tf, True), "AARE · HYBRID PIXEL DETECTORS · CUDA CLUSTERFINDER",
    9.5, MUTED, bold=True, spc=1.8)

# TWO PARAGRAPHS, NOT ONE WRAPPED LINE. At 42 pt the sentence is wider than the
# column either way, so the only question is WHERE it breaks -- and left to the
# wrap it breaks mid-clause. Splitting it puts the claim on line 1 and the
# punchline alone on line 2, which is the shape a TLDR wants.
#
# "rarely", not "never": at 3x3 the kernel is 5.5 us against 16.6 of H2D, but at
# 9x9 it is 39.9 us and IS the tallest bar -- slide 19 says so in its title. An
# absolute here would read as a walk-back sixteen slides later.
tf = tb(s, M + 0.31, 1.28, 11.78, 1.45)
run(para(tf, True, line=1.02),
    "The kernel was rarely the bottleneck.", 42, PALE, bold=True)
run(para(tf, line=1.02), "Feeding it was.", 42, PALE, bold=True)
tf = tb(s, M + 0.31, 2.80, 11.78, 0.45)
run(para(tf, True, line=1.05),
    "One kernel, one thread per pixel, and seven steps to keep it fed", 22, ACCENT)

# The body paragraph is gone. Every number in it -- 5.5 µs of kernel against
# 16.6 µs of H2D -- is said out loud in the notes and shown at full size on
# slide 19 and in A7, and a hero that has to be READ is a hero nobody hears.
tf = tb(s, M + 0.31, 4.05, 11.78, 0.26)
run(para(tf, True),
    "WHERE THIS ENDS UP · THE SHIPPED f32 BUILD AFTER ALL SEVEN STEPS · "
    "ENGINE TIMES [f32 · s4] · RECONCILED IN A7",
    8.5, AMBER, bold=True, spc=1.4)

# THE MISMATCH STAT IS AGAINST THE ORDINARY CPU FINDER, NOT THE TWIN. It used
# to read "6 / 23 M vs CPU twin", which is measured and correctly labelled, but
# leans on a word the audience does not have until slide 25 -- and once slide 29
# gives the 19-cluster update-timing residual a title of its own, a hero quoting
# only the 6 reads as the flattering half. The two are different causes on
# different builds, so they add rather than replace:
#     serial vs frozen        19   pedestal update timing   [f64]
#     frozen vs CUDA [f32]     6   float32 ULP ties
# and |serial /\ CUDA_f32| <= 19 + 6 by the triangle inequality on symmetric
# differences. 25 / 23 244 605 = 1.08 per million, so the box says ~1 in 10^6:
# a bound the deck can defend, and one nobody has to hold six digits to read.
#
# THE CLUSTER SIZE LEADS, because every other box is quoted at it and ×9.1 means
# nothing without it. "16.3 µs per frame" is gone: it is on slide 22 and again
# on slide 34, where there is room to qualify it.
stats = [("3×3", "CLUSTER SIZE", PALE),
         ("×9.1", "VS 24-THREAD CPU", ACCENT),
         ("61,312", "FRAMES / SECOND", PALE),
         ("~1 in 10⁶", "CLUSTER MISMATCH VS CPU", AMBER)]
for i, (v, l, c) in enumerate(stats):
    x = M + 0.31 + i * 2.95
    rect(s, x, 4.43, 0.04, 1.05, c)
    tf = tb(s, x + 0.22, 4.43, 2.58, 0.58)
    run(para(tf, True), v, 32, c, bold=True)
    tf = tb(s, x + 0.22, 5.11, 2.58, 0.3)
    run(para(tf, True), l, 8.5, MUTED, spc=1.2)

rect(s, M + 0.31, 5.80, 11.37, 0.012, RULE)
tf = tb(s, M + 0.31, 6.00, 11.78, 0.6)
run(para(tf, True, line=1.35),
    "RTX 4090 (Ada, sm_89) · PCIe 4.0 ×16 · Mönch 400×400 uint16 · 3×3 clusters · "
    "100 000 frames · Cu fluorescence, MAX IV", 10, MUTED)
# The asterisk and its footnote ("this em dash was hand written") went with the
# dash. The joke only worked while there was a dash on the slide to point at.
notes(s, """SAY THIS -- 92 words, about 35 s.

- The title is the finding, and it surprised me. I expected to spend my time on
  the kernel. The kernel was fast almost immediately.

- At the small cluster size the maths takes 5 microseconds per frame, and just
  getting that frame across the PCIe wire takes 16. So six of my seven steps are
  not about computing. They are about feeding, and about getting results back.

- Where we end up: nine times faster than 24 CPU threads, 61 thousand frames per
  second, and the answers are the same to about one cluster in a million.

DETAILS: WHERE THE ~1 IN A MILLION COMES FROM, IF ASKED.

It is a bound, and a provable one, not a measurement. Three comparisons, all at
3x3 over 10 000 frames and 23 244 605 clusters:

    serial CPU  vs  frozen CPU        8 / 11 = 19   pedestal update timing
    frozen CPU  vs  CUDA [f64 ped]     0 /  0 =  0   the port, on its own
    frozen CPU  vs  CUDA [f32 ped]     0 /  6 =  6   float32 ULP ties

The shipped build is CUDA [f32], and the finder a user would otherwise reach for
is the serial CPU one. That pair is never compared directly, but symmetric
difference obeys the triangle inequality, so it is at most 19 + 6 = 25. Two
independent causes, neither of them the GPU port: slide 29 has the first, slide
28 the second, and slide 26 shows the port itself contributes exactly zero.

25 / 23 244 605 = 1.08 per million. Say "about one in a million"; do NOT say
"below one in a million", which the bound does not support.

THE PARAGRAPH THAT USED TO BE ON THIS SLIDE is now only spoken: at 3x3 the
kernel needs 5.5 us per frame while getting that frame across PCIe costs 16.6,
or 13.2 uncontended. Both numbers return at full size on slide 19 and in A7.""")
# U+263A, not an emoji smiley: the colour-emoji planes (U+1F600+) are dropped
# entirely by the LibreOffice PDF export on this box. U+263A lives in DejaVu
# and Segoe UI Symbol, renders monochrome, and inherits the run colour.

# ------------------------------------------------------- divider · context
s = section("Context · what the code does",
        "The algorithm, the hardware, the limits",
        "No optimizations yet, only what the hardware has to do.",
        [("3–4", "The algorithm"),
         ("5–6", "The two machines"),
         ("7–8", "The CUDA kernel"),
         ("9", "What limits it"),
         ("10", "The roadmap")],
        rng=(3, 10))
notes(s, """SAY THIS -- 56 words, about 20 s.

- First, some context. What the algorithm does, what the two machines look
  like, and what we ask the GPU to do.

- No speed numbers in this part. We only set up the problem.

- One thing to keep in mind all the way through: every frame has 160 000
  pixels, and the work on each pixel is independent of the others.""")

# =========================================================== 3 · THE PHYSICS
s = new_slide()
chrome(s, 3, "The algorithm · what it is for",
       "A photon is not a pixel: it is a cluster")
bullets(s, M, 1.82, 8.16, [
    "One photon's charge **spreads over neighbouring pixels**; summing the 3×3 "
    "patch recovers its energy.",
    "The histogram of those cluster energies **is** the measurement: peak and "
    "width give gain and **energy resolution**.",
], size=PT_BODY)
figure(s, "fig_frame", M, 3.20, 7.70)
rail(s, [("label", "Why it matters")], y0=1.95)
h = card_figure(s, "img_spectra", RAIL_X, 2.32, RAIL_W)
caption(s, RAIL_X, 2.32 + h + 0.10, RAIL_W,
        "Cluster-energy spectra from an energy scan.", size=PT_META)
rail(s, [
    ("label", "MÖNCH03 · the detector this feeds"),
    ("row", "Array · pitch · active area", "400 × 400 · 25 µm · 10 × 10 mm²", TEXT2),
    ("row", "Frames per second", "1.3 k standard, 3–6 k optimised", AMBER),
    ("row", "Peak pixels = photons / frame", "~2 330   ·   1.5 %", ACCENT),
], y0=4.92, divider=False)
caption(s, M, 7.02, 8.16,
        "Real MÖNCH data, Cu fluorescence, MAX IV. One cluster per local maximum, "
        "so 2 330 counts photons.")
notes(s, """SAY THIS -- 118 words, about 45 s.

- A photon does not land on one pixel. Its charge spreads over the pixels next
  to it. So to get its energy, you have to add up a small patch -- three by
  three around the brightest pixel. That patch is what we call a cluster.

- The histogram on the right is the real measurement. Where the peak sits gives
  the gain, how wide it is gives the energy resolution. Finding clusters is
  what makes that plot possible.

- The detector is MOENCH: 400 by 400 pixels, 1 300 frames per second today, up
  to 6 000 with better readout boards. About 2 300 photons per frame -- so only
  one and a half percent of the pixels have a photon.

DETAILS: The detector, and why its frame rate is the yardstick.

MOENCH03 is a hybrid silicon pixel detector: charge integration with analog
readout, 25 x 25 um^2 pitch, 400 x 400 pixels over a 10 x 10 mm^2 active area.
Its standard frame rate is 1.3 kHz; with optimised readout boards the design
reaches 3-6 kHz depending on configuration.

That is the number to hold on to. 400 x 400 = 160 000 pixels is 312.5 KiB per
frame at 16 bit, so 1.3 kHz is ~0.4 GB/s off the detector and 6 kHz is ~1.9
GB/s. The 24-thread CPU finder does 6 762 FPS at 3x3, so it already keeps up
with the standard mode and roughly matches the optimised ceiling -- with nothing
left over for anything else on the machine, and nothing left at 9x9, where it
manages 1 503 FPS. Slide 33 closes this loop.

Pixels above 5 sigma are ~5 700, about 3.6 % of the frame and 2.4 per photon;
that is the number that sets how much of the frame the 3x3 windows cover.""")

# =========================================================== 4 · PER FRAME
s = new_slide()
chrome(s, 4, "The algorithm · per frame", "Per pixel: subtract, threshold, update the pedestal")
bullets(s, M, 1.82, COL, [
    "Per pixel: subtract a **running pedestal**, keep what clears **nσ · rms**, "
    "cut a 3×3 cluster around each local maximum.",
    "The pedestal is **updated by ~80 % of pixels every frame**, so arithmetic "
    "and data movement are coupled.",
])
code(s, M, 3.36, COL, [
    "// the whole algorithm, per pixel",
    "v   = frame[i] - pedestal_mean[i]",
    "rms = pedestal rms at i",
    "m   = max(v) over the 3x3 window at i",
    "if (m > «nSigma» * rms)                  // a photon is within reach",
    "    if (v == m) -> emit cluster        //   ... and I am its peak",
    "    else        -> «nothing»             //   ... I am in its shadow",
    "else            -> update pedestal     // I saw nothing",
], title="THE WHOLE ALGORITHM · THREE OUTCOMES, NOT TWO")
callout(s, M, 5.86, COL,
        "**Thesis of this talk:** the compute was fast almost immediately. "
        "Six of the seven steps are about feeding it.")
rail(s, [
    ("label", "The shape of the work"),
    ("gap", 0.20),
    ("stat", "Work items per frame", "160 000", PALE),
    ("gap", 0.10),
    ("row", "Operations on each", "~5, identical", TEXT2),
    ("gap", 0.18),
    ("row", "Communication between them", "none", ACCENT),
    ("gap", 0.18),
    ("row", "Order they may run in", "any, once the pedestal is fixed", ACCENT),
    ("gap", 0.20),
    ("row", "Bytes per frame", "312.5 KiB · ~2 330 clusters", TEXT2),
], y0=2.10)
notes(s, """SAY THIS -- 127 words, about 46 s.

- This is the whole algorithm, for one pixel. First subtract the pedestal --
  that is the dark baseline, what the pixel reads with no light. Then look at
  the three by three window around you.

- Now three outcomes, not two. If the brightest value in that window is above
  the threshold and it is ME, I write out a cluster. If it is a neighbour, I do
  nothing -- I am in its shadow. And if nothing is above the threshold, then I
  saw only noise, so I add my value to the pedestal. That last case is about
  80 % of pixels, every frame.

- So: 160 000 pixels, about five operations each, and no pixel needs to talk to
  another. That shape is why a GPU should work here.

DETAILS: Why order does not matter, and the one place it does.

A pixel reads what its neighbours MEASURED, never what they decided. With the
pedestal held at its frame-start value, order cannot matter, and that is what
licenses one thread per pixel.

The serial CPU finder updates the pedestal mid-scan and so gives that property
up: a pixel scanned late sees a pedestal that earlier pixels have already moved.
Slides 30-31 measure what that costs -- it is the channel behind 11 of the 19
cluster differences between the serial CPU and the frozen-pedestal reference.""")

# ==================================================== 5 · THE CPU
# The two machine slides. They exist because the audience is asked, from slide 7
# on, to accept "one thread per pixel" and "occupancy" without ever having been
# shown what a thread costs on either machine. Both diagrams are the CS149 ones
# (credited in the captions): redrawing them would lose the shared visual
# grammar — orange fetch/decode, yellow ALU, blue execution context — which is
# the entire reason the pair reads at a glance.
s = new_slide()
chrome(s, 5, "The machine we are starting from",
       "CPU: latency-oriented, built to finish one thread fast")
figure(s, "img_cpu_core", M, 1.82, 6.72)
# The die photo mirrors slide 6's: same grammar, compute units boxed, so the two
# machines are compared as objects and not only as block diagrams. It is kept
# WHOLE rather than cropped to the ten cores, because the L3 slab on the right
# and the I/O block on the left are half the die area -- which is the callout's
# point standing next to it in silicon.
figure(s, "img_cpu_die", M, 5.04, 4.70)
callout(s, 5.38, 5.04, 3.29,
        "Count the boxes: **6 fetch/decode** and two levels of private cache, all "
        "to keep **two** instruction streams fed. The ALUs are the small part.", h=1.34, size=PT_LEAD)
caption(s, 5.38, 6.58, 3.29,
        "Both diagrams after Stanford CS149, Fall 2025.", size=PT_META)
rail(s, [
    ("label", "pc-moench-04 · AMD Ryzen 9 7950X"),
    ("gap", 0.10),
    ("stat", "Cores × SMT", "16 × 2", PALE),
    ("gap", 0.04),
    ("stat", "CPU, 1 thread", "1.75 ms", MUTED),
    ("stat", "CPU MT, 24 threads", "148 µs", PALE),
    ("gap", 0.06),
    ("row", "That is the bar", "6 762 frames / s", ACCENT),
    ("gap", 0.16),
    ("note", "Thread count is swept per cluster size, not assumed: "
             "24 at 3×3, 32 at 9×9."),
])
notes(s, """SAY THIS -- 124 words, about 45 s.

- This is the machine we start from. Look at what the chip is spent on: six
  fetch and decode units, an instruction selector, two levels of private cache.
  All of that keeps just two instruction streams running fast.

- The yellow boxes, the parts that actually compute, are the small ones. Most
  of the chip is there to FIND independent work inside one program, and to hide
  slow memory behind a cache. That is the right design when you have a few
  threads and each one must finish quickly.

- Our problem is the opposite. The independent work is handed to us: 160 000
  pixels. The bar to beat is here: 24 threads, 148 microseconds per frame,
  6 762 frames per second.

DETAILS: The CPU slide. The point is not that CPUs are bad -- it is what the
silicon is SPENT on.

On the diagram: Intel Skylake is shown, one core, schematically. The Zen 4 core
in this machine differs in detail -- 4 FP pipes rather than 3, AVX-512
double-pumped on 256-bit datapaths -- but not in kind: ~6-wide front end, 4
scalar ALUs, 2 SMT contexts, private L1 + L2. Load/store units are not drawn.

DO NOT ADD THE BOXES UP INTO A LANE COUNT. Counting the picture gives 3 x 8-wide
vector + 4 scalar = 28 ALU boxes, and 28 is not a number to say next to the GPU's
16 384. Three reasons:

  1. They are not separate hardware. The drawing gives the scalar ALUs and the
     vector ALUs their own boxes, but on Skylake they sit on the SAME issue
     ports (0, 1, 5). Adding 24 + 4 counts the same execution resources twice.
  2. Only two of the three vector units multiply -- read the red labels: "MUL or
     ADD", "MUL or ADD", and then "ADD" alone. FLOPs need FMA-capable width, so
     it is 2 x 8, not 3 x 8.
  3. "8-wide" is 8 x 32 bit = 256 bit = AVX2, so those ARE fp32 lanes. That part
     of the reading is right.

The number to quote is 16 FP32 FMA lanes per core: on Skylake client, 2 FMA units
x 256 bit. And on the Zen 4 in this machine, also 16 -- 4 FP pipes of which 2 are
FMA-capable, 256 bit wide. 16 cores x 16 = 256, which is the figure on slide 6's
rail. Say it comes from the Zen 4, NOT from this Skylake picture: the two
microarchitectures landing on the same 16 is a coincidence, not a derivation.

Same shape of caveat as the V100 diagram on slide 6, and the same answer if the
count is challenged: the picture is there for the RATIO of control silicon to
arithmetic -- 6 fetch/decode, an instruction selector, two private caches, all
feeding 2 execution contexts -- not for its lane count.

On the die photo: Comet Lake, a 10-core Core i9, at the same scale of argument
as the AD102 die on the next slide. Ten cores, and they do not fill the die --
the orange slab on the right is L3 and the block on the left is I/O and the
memory controller. If someone asks why a 16-core Ryzen is not simply 16x a
1-core Ryzen, that picture is the answer.

Six fetch/decode units, an instruction selector, L1 + L2 private cache: all of
that machinery exists to find independent work INSIDE one instruction stream,
and to hide memory latency behind a cache. It is the right design when you have
a few threads that must each finish fast.

WHAT THE DIAGRAM DOES *NOT* SHOW, if asked. It draws exactly four things:
superscalar issue (6 fetch/decode + several ALUs + one "instruction selection"
box), SIMD (3 x 8-wide vector ALUs), 2-way SMT (the two Execution Contexts), and
a private L1/L2. There is NO out-of-order hardware on it -- no reorder buffer,
no reservation stations, no register renaming, no load/store queue, no branch
predictor, no retirement unit. The single "instruction selection" box is the only
nod to issue logic, and its label does not say out-of-order; on a diagram whose
subject is SMT it reads just as naturally as "pick which of the two threads
issues this cycle".

So do not point at this picture and say "out-of-order". The real Zen 4 in this
machine is deeply out-of-order and that fact supports the argument -- but say it
as a fact about the chip, not as a box on the slide, or someone will spend the
next minute looking for the box.

Our problem has the opposite shape: 160 000 work items that are already
independent. We do not need a machine to FIND the parallelism -- it is handed to
us. Every transistor spent looking for it is a transistor not doing arithmetic.

On the thread sweep: 6 762 FPS is the best of a measured sweep (perf/cpu_threads.py),
not a default. 48 threads on 16 cores is 24 % SLOWER than 24 -- worth saying out
loud, because an oversubscribed baseline is the easiest way to inflate a GPU
speedup without lying about anything.""")

# ==================================================== 6 · THE GPU
s = new_slide()
chrome(s, 6, "The machine we are moving to",
       "GPU: throughput-oriented, the whole frame at once")
figure(s, "img_gpu_die", M, 1.82, 2.84)
caption(s, M, 5.16, 2.84,
        "AD102 · 144 blocks, 128 enabled on this card. One SM boxed.", size=PT_META)
figure(s, "img_gpu_sm", M + 3.15, 1.82, 4.91)
# The colour key is a separate crop: in the source it spans the full slide width
# while the diagram spans 60 % of it, so one rectangle cannot hold both.
figure(s, "img_gpu_legend", M + 3.15, 5.16, 4.91)
callout(s, M, 5.86, COL,
        "Same grammar, inverted proportions: **4 fetch/decode** for **64 warp "
        "contexts**.", h=0.72, size=PT_LEAD)
caption(s, M, 6.74, COL,
        "One SM: a V100 is shown; this card's is the same idea: 128 FP32 lanes, "
        "48 warp slots, 100 kB shared memory. After Stanford CS149, Fall 2025.", size=PT_META)
rail(s, [
    ("label", "NVIDIA GeForce RTX 4090"),
    ("gap", 0.12),
    ("stat", "FP32 lanes", "16 384", PALE),
    # NOT "vs 32 CPU streams": that compared GPU SIMD lanes against CPU hardware
    # threads, two different quantities, and inflated the ratio to ×512. SMT adds
    # threads, not lanes — the 7950X's 32 threads share 16 cores × 16 FP32 lanes.
    # Like for like it is ×64 on width, ×32 on peak FP32. The notes carry the grid.
    ("row", "128 SMs × 128 lanes", "vs 16 cores × 16 = 256", TEXT2),
    ("gap", 0.14),
    ("stat", "Resident thread slots", "196 608", ACCENT),
    ("row", "The frame needs", "160 000  →  all at once", AMBER),
    ("gap", 0.20),
    ("note", "Device memory runs at 1 008 GB/s (384-bit, 21 Gbps); PCIe 4.0 ×16 "
             "peaks at 31.5 and measures 24. Two copy engines, so H2D and D2H "
             "run at the same time."),
])
notes(s, """SAY THIS -- 144 words, about 52 s.

- Same drawing, opposite proportions. Four fetch and decode units, and 64 warp
  contexts. Almost all the area is compute and registers.

- The GPU hides slow memory in a different way. Not with a big cache, but by
  having many other threads ready to run. When one group waits, another runs.

- The number to keep: this card can hold 196 000 threads at the same time. Our
  frame needs 160 000. The whole frame fits in the machine at once.

- And compare like with like: 256 float lanes on this CPU, 16 384 here. That is
  64 times the width, about 32 times once you count clock speeds. We measure
  nine. So we are not short of lanes -- we are short of a wire to feed them.
  Device memory runs at a thousand gigabytes per second; PCIe brings us 25.

DETAILS: The GPU slide, and the number to land.

128 SMs x 1 536 threads = 196 608 thread slots that can be RESIDENT at the same
time. Our frame is 160 000 pixels. The entire frame fits in the machine at once,
one thread per pixel, with room left over -- which is why slide 7's "one thread
per pixel" is not a figure of speech, and why slide 9 can talk about occupancy
as a real quantity (3x3 achieves 100 %: 1 536 threads resident per SM).

Latency hiding, in one sentence: the CPU hides memory latency with a cache and
out-of-order execution; the GPU hides it by having 48 other warps ready to run.
That is why there is no reorder buffer on this diagram and no branch predictor.

Count the pink cells in the legend before moving on: FP64 units are a small
minority of the SM, 8 MUL/ADD per clock against 16 for fp32 and 16 for int. On
a GeForce part the ratio is far worse than this V100 diagram suggests -- 1/64.
That is half of why opt7 pays, and it is worth planting here so opt7 lands as a
consequence rather than a surprise.

THE DIAGRAM IS A V100 AND ITS LANE COUNT DOES NOT TRANSFER. Anyone who counts
squares gets 64 FP32 lanes per SM, not 128, and they are right about the picture.
Per sub-core the legend reads 16 yellow "SIMD fp32" and 16 salmon "SIMD int" --
on Volta those are SEPARATE units and the int ones cannot do float. Ada changed
exactly that: per sub-core, 16 FP32-only lanes plus 16 that do FP32 OR INT32, so
4 x 32 = 128 per SM, and 128 SMs x 128 = 16 384. Cross-check: 16 384 x 2 flop x
2.52 GHz = 82.6 TFLOP/s against NVIDIA's quoted 82.58.

Two more Volta-isms in the same picture: it draws 64 warp slots per SM where Ada
has 48 -- slide 9 divides by 48 -- and it draws FP64 at half the FP32 rate,
where a GeForce Ada part is 1/64. The caption on the slide already says a V100
is shown; this is the detail behind that sentence.

THE TRAP TO PRE-EMPT, AS A LIKE-FOR-LIKE RATIO. The old version of this line
compared 16 384 GPU lanes against 32 CPU threads and called it 512x. That is
lanes against hardware threads -- two different quantities -- and it is the kind
of thing this audience catches. The defensible comparisons:

    FP32 lanes     CPU 16 cores x 16 = 256     GPU 16 384         x64
    peak FP32      ~2.6 TFLOP/s at ~5 GHz      82.6 TFLOP/s       x32
    memory BW      ~71 GB/s measured           1 008 GB/s         x14
    MEASURED       6 762 FPS                   61 312 FPS         x9.1

The CPU's 16 FP32 lanes per core are 2 FMA pipes x 256 bit, and SMT adds threads,
not lanes -- 32 threads share the same 256. Quote x64 or x32. Never x512.

The point survives either way, and it is the last row: even 32x is not what we
get, because every frame crosses a 25 GB/s wire once in each direction. At 9x9,
where there is far more work per byte moved, we DO get x26.5 -- the same
neighbourhood as the compute ratio. That contrast between the two cluster sizes
is the talk.""")

# =========================================================== 7 · THE KERNEL
s = new_slide()
chrome(s, 7, "The CUDA kernel · execution model",
       "One thread per pixel, one block per 16×16 tile")
bullets(s, M, 1.82, 7.75, [
    "A **16×16 block = 256 threads**; the grid tiles the frame in **625 "
    "identical blocks**.",
    "The output is **sparse**, but the **decision work is dense**: every pixel "
    # Non-breaking space: the line broke between "160" and "000".
    "is tested, **160 000 times per frame**.",
], size=PT_BODY)
code(s, M, 3.30, 7.75, [
    "«block» = dim3(BLOCK_X, BLOCK_Y);           // 16 x 16 = 256 threads",
    "«grid»  = dim3((ncols + BLOCK_X - 1)/BLOCK_X,  // 25 x 25 = 625 blocks",
    "               (nrows + BLOCK_Y - 1)/BLOCK_Y);",
    "",
    "find_clusters_in_single_frame<<<«grid», «block», shmem, stream>>>(",
    "    d_frame, d_pd_mean, /* ... 9 more ... */ d_clusters);",
], size=PT_CODE, title="LAUNCH CONFIGURATION · ClusterFinderCUDA.hpp")
caption(s, M, 5.10, 7.75,
        "400×400 → a 25×25 grid of 16×16 blocks = 625 blocks per frame, over 128 "
        "SMs. Nothing in the launch depends on the cluster count.")
# The right column used to restate slide 4's three outcomes -- same logic, same
# panel title, same ~80 % figure. What belongs here instead is the thing this
# slide is FOR and the notes already argue: the CPU's loop over pixels becomes
# the launch geometry, and the pixel a thread owns is arithmetic on its own ids.
#
# THE CPU SIDE IS QUOTED, NOT INVENTED. It used to read "find_cluster(iy, ix)",
# which is not an entry point that exists: aare has ONE call, find_clusters(frame)
# at ClusterFinder.hpp:123, and the per-pixel test is INLINE in its double loop
# (:137-139). A reader who knows the code hears a per-pixel API that isn't there;
# a reader who doesn't, learns the wrong shape. The loop itself stays -- it is
# real, and it is the whole subject of the panel.
code(s, 8.56, 1.82, 4.24, [
    # "one thread per FRAME" rather than "one call per frame": it is equally true
    # (ClusterFinderMT::process pops a whole frame per thread, MT.hpp:65-75) and
    # it rhymes with the title, so the panel reads frame-vs-pixel on the same noun.
    "// CPU: one thread per FRAME, loop inside",
    "«find_clusters»(frame) {",
    "  for (int iy = 0; iy < 400; iy++)",
    "    for (int ix = 0; ix < 400; ix++)",
    "      /* test pixel (iy,ix) inline */",
    "}",
    "// GPU: no loop. The launch is the loop:",
    "ix = «blockIdx».x*blockDim.x + «threadIdx».x;",
    "iy = «blockIdx».y*blockDim.y + «threadIdx».y;",
], size=PT_CODE, title="WHERE THE LOOP WENT")
callout(s, 8.56, 4.22, 4.24,
        "**The loop is the launch.** 625 blocks × 256 threads = **160 000**, one "
        "per pixel, decided before any data is seen.", h=0.86, size=PT_LEAD,
        color=RED)
callout(s, 8.56, 5.22, 4.24,
        "A **block** is the 256 threads that share one tile and can wait for each "
        # h=0.72 clipped: at 4.24 in wide this wraps to three lines and "slide."
        # sat on the panel border. 0.88 ends at 6.10, clear of the flow at 6.20.
        "other. That is the next slide.", h=0.88, size=PT_LEAD, color=AMBER)
# The five steps go LAST. __syncthreads is the second one and it cannot mean
# anything until "a block is 256 threads that can wait for each other" has been
# said -- which is the amber callout directly above this line.
flow(s, M, 6.20, 12.30, ["load tile + halo", "__syncthreads", "stencil reduction",
                        "classify", "append or update pedestal"])
notes(s, """SAY THIS -- 89 words, about 30 s. Two differences from the CPU code.

- WHAT ONE WORKER OWNS. CPU: one thread takes one WHOLE frame -- 24 threads on
  24 frames -- and walks it pixel by pixel. GPU: one thread takes ONE PIXEL,
  160 000 at once on the same frame.

- WHERE THE LOOP WENT. On the CPU the loop is in the code: for y, for x. On the
  GPU there is no loop -- the loop IS the launch. A 25 by 25 grid of 16 by 16
  blocks, each thread finding its own pixel. And a block is 256 threads that
  share one tile: the next slide.

DETAILS (if asked):
  - Every thread runs the same function. The kernel is written once, for one
    thread; the launch runs it 160 000 times.
  - "One thread per frame" is why the CPU baseline moves with thread count:
    24 threads at 3x3, 32 at 9x9.
  - Each CPU thread carries its own pedestal, so neighbouring frames train
    separate baselines. Act IV, not here.

Do NOT say "a chunk of the frame" for the CPU: ClusterFinderMT hands each thread
a whole frame, and ClusterFinder walks it in raster order. Slide 30 needs that
-- its whole argument is that a late pixel is judged against a pedestal this
same scan already moved.

DETAILS: what this slide used to carry, and where it went. The right column was a
copy of slide 4's three-outcome logic -- same panel title, same ~80 % figure --
so it has been replaced by the loop/launch contrast the notes above argue. The
two facts that lived only here, if either comes up:

  - The branch split per frame: ~80 % of threads update the pedestal, ~18 % sit
    in a peak's shadow, ~1.5 % are peaks. (Also on slide 30.)
  - A shadow pixel is neither recorded NOR fed back: it is not a cluster, and it
    does not update the pedestal either. That is the third outcome, and it is
    why the algorithm has three and not two.

DETAILS: the launch line is elided on purpose. The real call passes 12 arguments
(ClusterFinderCUDA.hpp:715-719); nine of them are pedestal and geometry pointers
that say nothing about the execution model, and at full length they were the
widest thing on the slide. The bold is on grid and block because those two are
the subject; shmem is deliberately left plain -- it is slide 8's subject, one
slide later.""")

# =========================================================== 6 · TILING
s = new_slide()
# The eyebrow says "tiling and halos", not "shared memory", because the title now
# carries that phrase and the two sit four lines apart.
chrome(s, 8, "The CUDA kernel · tiling and halos",
       "Shared memory: load the tile once, reuse it nine times")
bullets(s, M, 1.82, 12.40, [
    "Neighbouring threads need **overlapping** windows; without shared memory "
    "each pixel is fetched up to **nine times**.",
    "Each block stages one **pedestal-subtracted** tile plus its halo, loaded "
    "by the threads on the block edges.",
], size=PT_BODY)
figure(s, "fig_tile", M, 2.70, 8.36)
code(s, 9.02, 2.70, 3.88, [
    "«extern __shared__» unsigned char smem[];",
    "auto *sh = (COMPUTE_TYPE*)smem;",
    "auto stride = blockDim.x + 2*col_radius;",
    "auto tid = (threadIdx.y + row_radius)*stride",
    "         + (threadIdx.x + col_radius);",
    "// pedestal subtraction fused into the load",
    "«sh»[tid] = d_frame[gid] - d_pd_mean[gid];",
], size=PT_CODE, title="clusterfinder_kernel.cuh")
callout(s, 9.02, 4.72, 3.88,
        "Only **odd** cluster sizes are supported, so the centre pixel is unique.",
        h=0.92, size=PT_LEAD)
notes(s, """SAY THIS -- 123 words, about 45 s.

- Each thread needs its own 3x3 window, so it reads nine values. Which means
  every pixel is read nine times, by nine different threads. At 9x9 it is
  eighty-one times.

- What CUDA gives you here is control over the chip's memory. Shared memory is
  the same on-chip memory as L1, except you decide what goes in it.

- So the block loads its tile once, plus a one-pixel halo, and after that every
  thread reads from shared memory. One trip to VRAM instead of nine -- and the
  pedestal subtraction happens on the way in, for free.

- The cost: 1.3 kilobytes at 3x3, 2.2 at 9x9, out of 100 per SM. One to two
  percent. Shared memory is never what runs out. So what does?

DETAILS: shared memory is never the binding resource here. The tile is stored in
COMPUTE_TYPE (float), not in the pedestal type: 1.3 KB for 3x3 and 2.3 KB for
9x9, against 100 KB of shared memory per SM. Even the old double-precision tile
only reached 4.5 KB. What runs out first is registers -- the next slide.

DETAILS: "nine reads" is per pixel, not per thread. A thread reads 9 values (its
own window); a pixel is read by 9 threads (every window it falls in). Both are
nine at 3x3, which is why the sentence works either way -- but the WASTE is the
second one, so say it that way.

DETAILS: on Ada, shared memory and L1 are one 128 KB block of on-chip storage per
SM, split at launch, up to 100 KB of it usable as shared. So "same silicon as L1,
but user-managed" is literal, not an analogy.""")

# ================================================ 8 · OCCUPANCY (the output)
s = new_slide()
chrome(s, 9, "Hardware · occupancy",
       "Occupancy is a latency-hiding budget")
bullets(s, M, 1.82, 12.30, [
    "When a warp stalls, the SM switches to another **already-resident** warp. "
    "**Occupancy = resident warps / the maximum**: how many alternatives it "
    "has to switch to.",
], size=PT_BODY)
# THE STRIP CARRIES THE BUDGETS, NOT THE ANSWER. "block size 16 × 16" and
# "occupancy 100 % · 33 %" were both already in the figure below -- 100/33 at
# full size on the left panel, the block sizes across the right one -- so the
# strip was repeating the conclusion instead of supplying the three numbers you
# divide to reach it. Slots, registers and shared memory are the three ceilings,
# and with 256 threads/block beside them the whole arithmetic is on the slide.
# The register-pressure slide that used to precede this one is annex A2, and the
# strip did NOT need to absorb its numbers: fig_occupancy already prints "38
# regs/thread, 6 blocks/SM" and "128 regs/thread, 2 blocks/SM" under its own
# bars. Only the pointer was missing, and that is in the callout.
# The three RINGED boxes are the hardware's own ceilings -- the things that can
# run out. The other two are properties of this launch (256 threads a block) or
# arithmetic on a ceiling (48 warps = 1 536 / 32), and reading them as limits is
# exactly the confusion the ring removes.
statstrip(s, M, 2.55, 12.30, [
    ("threads / block", "256"),
    ("thread slots / SM", "1 536", RED),
    ("max warps / SM", "48"),
    ("registers / SM", "65 536", RED),
    ("shared mem / SM", "100 kB", RED),
], h=0.84)
figure(s, "fig_occupancy", 1.19, 3.44, 10.95)
callout(s, M, 6.62, 12.30,
        "**16×16 is the balance point.** 33 % is not a failure to fix: it is what "
        "the register budget allows (**annex A1**), and at 9×9 one kernel nearly "
        "fills the machine on its own.", h=0.62, size=PT_LEAD)
notes(s, """SAY THIS -- 120 words, about 45 s.

- This is the last two slides in one number. Occupancy is the warps resident
  on an SM over the most it could hold. It is a latency-hiding budget: when a
  warp stalls, the SM switches to another one that is already there.

- At 3x3: six blocks, 48 warps, 100 percent. At 9x9: two blocks, 16 warps,
  33 percent.

- And 33 is not a bug to fix. Occupancy only buys latency hiding, and hiding
  only pays if there is latency left unhidden. At 9x9 each thread does a lot of
  work per byte it loads, so 16 warps keep the SM fed.

- High occupancy is a means, not the goal -- and on this workload the kernel was
  never the bottleneck anyway.

CAREFUL: occupancy has nothing to do with arithmetic intensity. Intensity is
FLOP per byte, a property of the ALGORITHM; occupancy cannot change it. The true
statement is "more occupancy does not mean more throughput" -- a kernel with
enough independent work per thread runs full speed at low occupancy. Saying
"highest occupancy is not highest arithmetic intensity" invites a correction
from anyone in the room who knows the roofline.

DETAILS: why "33 % occupancy" is not the alarm it looks like.

Occupancy only buys latency hiding, and latency hiding only matters if there is
latency left unhidden. At 9x9 the kernel is long and each thread does a lot of
work per byte loaded, so 16 resident warps are enough to keep the SM fed.

It is also build-dependent, and that is the honest caveat. On the shipped f32
build the 9x9 kernel overlaps across four streams by only 1.02x -- extra streams
buy nothing in kernel concurrency. On the f64 build the kernel is 3.5x longer
and DOES overlap, 1.32x, and there four streams genuinely lower the floor. So
quote the overlap factor, never "occupancy is fine" on its own.

Shared memory never binds here: 1.3 KB per block at 3x3, 2.3 KB at 9x9, against
100 kB per SM. That is 1-2 % of the budget, which is why the strip lists it as a
ceiling we never approach rather than as a number that decides anything.

ON THE 100 kB, if someone who knows sm_89 asks. The SM has 128 kB of UNIFIED
L1 + shared memory; 100 kB is the most of it that can be carved out as shared,
and 99 kB the most a single block may request. So 100 kB is the shared-memory
ceiling, not the size of the array. The three numbers on the strip are the three
things that can run out -- and on this kernel it is never this one.

Provenance for both hardware slides: shipped f32 build, measured not estimated.
cuobjdump -res-usage for registers and spills;
cudaOccupancyMaxActiveBlocksPerMultiprocessor for blocks per SM; RTX 4090
(sm_89). Reproduce with python/tests/perf/kernel_resources.py.""".rstrip())
# =========================================================== 10 · THE LADDER
s = new_slide()
chrome(s, 10, "Roadmap",
       "Three acts, ordered by which bar is tallest")
rows = [
    ("act", "ACT I  ·  feed the GPU", "the host cannot submit work fast enough", "at 3×3  [f64]", ACCENT),
    ("opt1", "First CUDA port", "1 stream, one launch per frame", "×2.34", ACCENT),
    ("opt2", "Streams + batching", "4 streams, drained every 4 frames", "×3.66", ACCENT),
    ("opt3", "Pipeline rework", "sync barriers removed", "×4.32", ACCENT),
    ("opt4", "Pinned memory", "DMA-speed host transfers", "×5.69", ACCENT),
    ("act", "ACT II  ·  get the results back", "the host copy is now the tallest bar", "at 3×3  [f64]", PALE),
    ("opt5", "Host↔GPU overlap", "chunked submit / collect", "×7.45", PALE),
    ("opt6", "Zero-copy collection", "read in place, never copy", "×8.65", PALE),
    ("act", "ACT III  ·  the kernel", "the kernel is the tallest bar, at 9×9", "at 9×9  [f32]", AMBER),
    ("opt7", "FP32 pedestal + variance rewrite", "the only kernel change in the deck",
     "kernel −41%", AMBER),
]
y = 1.78
for i, (tag, name, sub, gain, col) in enumerate(rows):
    if tag == "act":
        rect(s, M, y + 0.31, 12.30, 0.016, col)
        tf = tb(s, M + 0.02, y, 8.27, 0.28)
        run(para(tf, True), name.upper(), 12, col, bold=True, spc=1.4)
        tf = tb(s, M + 5.37, y + 0.04, 6.92, 0.26)
        run(para(tf, True, align=PP_ALIGN.RIGHT), f"{sub}   {gain}", 10, MUTED)
        y += 0.42
        continue
    rect(s, M, y, 12.30, 0.50, PANEL)
    rect(s, M, y, 0.04, 0.50, col)
    tf = tb(s, M + 0.29, y + 0.10, 1.03, 0.32)
    run(para(tf, True), tag, 13, col, bold=True, font=MONO)
    tf = tb(s, M + 1.50, y + 0.04, 5.17, 0.28)
    run(para(tf, True), name, 12.5, PALE, bold=True)
    tf = tb(s, M + 1.50, y + 0.27, 5.79, 0.26)
    run(para(tf, True), sub, 10, MUTED)
    tf = tb(s, 9.49, y + 0.09, 3.20, 0.35)
    run(para(tf, True, align=PP_ALIGN.RIGHT), gain, 13.5, col, bold=True)
    y += 0.55
caption(s, M, 6.98, 12.30,
        "Speedups are 3×3 against the best CPU configuration, 24 threads.")
notes(s, """SAY THIS -- 104 words, about 40 s.

- Here is the whole talk on one page. Seven steps, grouped in three acts, and
  the acts are ordered by which part was the slowest at the time.

- Act one: the host cannot feed the GPU fast enough. Act two: getting the
  results back is now the slow part. Act three: only now is the kernel itself
  the slow part -- and that is the only step where I touched the maths.

- The numbers on the right are speed-ups against the CPU. We start at two and
  finish at nine, for the small cluster size. Do not read every row now; we
  come back to this page at the end.""")

# --------------------------------------------------------- divider · ACT I
s = section("Act I of III · feed the GPU",
        "Bottleneck: host submission and PCIe",
        "The kernel was fast almost immediately. This act is about the host, and it "
        "is told at 3×3, where PCIe is the floor.",
        [(11, "two memories, one wire"),
         (12, "opt1 · first port"),
         (13, "opt2 · four streams, four deep"),
         (14, "opt3 · the whole batch in flight"),
         (15, "opt4 · pinned memory")],
        rng=(11, 15),
        carry=("Starting from", "6 762 FPS",
               "24-thread CPU · 14.8 s for 100 000 frames"))
notes(s, """SAY THIS -- 62 words, about 22 s.

- Act one. The kernel was already fast. The problem is getting the data to it.

- Five slides, all at the small cluster size, because that is where the PCIe
  wire is the limit.

- We start from the CPU: 6 762 frames per second, about 15 seconds for a
  hundred thousand frames. Everything from now on is compared to that.""")

# ============================================ 12 · WHAT H2D AND D2H MEAN
# Act I and Act II are five slides attacking two arrows, and until this slide
# existed the audience met the arrows already named and already being optimised.
# It buys the whole of both acts a shared vocabulary for one page.
s = new_slide()
chrome(s, 11, "Act I · the two memories",
       "The GPU consumes 32× faster than PCIe delivers")
bullets(s, M, 1.82, 12.30, [
    "A CUDA kernel addresses **the GPU's own memory and nothing else**. So every "
    "frame is copied in (**H2D, host to device**) and every result copied back "
    "out (**D2H**). Those two names are what the next five slides are about.",
], size=PT_BODY)
figure(s, "fig_gpu_model", 0.95, 2.82, 11.20)
callout(s, M, 6.40, 12.30,
        "**The asymmetry is the whole talk.** Host DRAM runs at ~71 GB/s and VRAM "
        "at 1 008, but everything between them crawls through **31.5**.",
        h=0.78, size=PT_LEAD)
notes(s, """SAY THIS -- 119 words, about 43 s.

- One rule of CUDA explains the next five slides. A GPU program can only read
  the GPU's own memory. It cannot touch host memory. So every frame has to be
  copied IN, and every result copied back OUT. Those two copies have names:
  H2D, host to device, and D2H, device to host.

- Now look at the three speeds. Inside the GPU, memory runs at a thousand
  gigabytes per second. Host memory runs at seventy. But the wire between them,
  PCIe, gives us thirty-one.

- So the GPU can eat about thirty times faster than the wire can deliver. Our
  data crosses that wire once per frame, in each direction. That is the whole
  problem of the next two acts.

DETAILS: This slide is for the half of the room that has never written CUDA. If
they are all GPU people, say the two sentences on the left and move on.

The one idea: a kernel cannot dereference a host pointer. Everything it reads
must already be in VRAM, and everything it produces must be copied back. That is
not a performance detail, it is the programming model, and it is why six of the
seven optimisations are about movement rather than arithmetic.

THE THREE BANDWIDTHS, and be precise about which comparison is which. VRAM to SM
is 1 008 GB/s. Host DRAM measured ~71 GB/s on pc-moench-04 (threaded copy over a
1 GB array, counting 24 B per element: 8 read, 8 write, 8 read-for-ownership;
dual-channel DDR5, so this is in the right family). PCIe 4.0 x16 is 31.5.

So the wire is 32x below VRAM but only about 2.3x below host DRAM. The claim to
make is NOT "the wire is 32x slower than memory" -- it is that the wire is the
NARROWEST link in the chain, and that the GPU can consume about thirty times
faster than the wire delivers. Any calculation whose input crosses that wire once
per use is gated by the wire, and this one does exactly that: a frame arrives, is
touched about five times, and leaves.

If asked why not keep the data resident on the GPU: for this workload the frames
come off a detector at 100 kHz and the results go to analysis on the host, so
both ends are genuinely on the host side. Streaming is the problem, not a
choice.

The DMA engines are the reason the story is not simply "the wire is slow, the
end". There are two, one per direction, independent of the SMs, so transfer and
compute CAN overlap. Act I is getting the input arrow overlapped; Act II is the
output arrow.

Numbers on this slide, all from A1: H2D 13.2 us [s1] / 16.6 [s4], kernel 5.5 us
at 3x3 f32, D2H payload 93 kB at 3x3.""")

# =========================================================== 13 · OPT1
s = new_slide()
chrome(s, 12, "Act I · the first CUDA port",
       "The first port runs at 26 % of the GPU's floor", opt=1)
# The tiling/halo bullet is gone: slide 8 is entirely about it, and repeating it
# here made the reader start on kernel internals when the slide's subject is the
# pipeline around the kernel.
bullets(s, M, 1.82, COL, [
    "One copy in, one kernel, one copy out. The calls are **async**, but the "
    "clusters have to be **on the host and copied out** before the next frame "
    "starts, so **nothing overlaps**.",
])
# The calls are cudaMemcpyAsync, NOT cudaMemcpy: what serialises opt1 is the
# cudaStreamSynchronize, not a blocking copy. The 1-4 comments name the three
# engines by the SAME labels the rail, the callout and the timeline picture use.
#
# TWO DELIBERATE ELISIONS, both so this slide does not forward-reference later
# ones. (1) opt1 really stages the frame through a pinned buffer first
# (std::memcpy into sc.h_frame); showing it would mean explaining pinned memory
# four slides before opt4 introduces it, and it is the hop opt4 deletes, so it
# belongs there. h_src is neutral about where the source lives. (2) There are
# really TWO D2H copies -- the count, then the payload sized from it. That is
# annex A4's entire reveal, so this panel shows the payload only.
# SKIM STRUCTURE, not just shorter lines. Three things do the work:
#   - the two cudaMemcpyAsync calls are padded so dst, src and the H2D/D2H token
#     land in the SAME columns -- the pair reads as one call twice, differing
#     only in direction;
#   - every trailing comment starts at column 59, so the 1-2-3-4-5 spine can be
#     read straight down without reading any code at all;
#   - the blank line splits the panel into its two halves: three calls that do
#     not block, then the two lines that do.
code(s, M, 2.82, COL, [
    "cudaMemcpyAsync(d_frame,    h_src,      H2D, sc.stream);   // 1. frame in",
    "find_clusters_in_single_frame<<<..., sc.stream>>>(...);    // 2. kernel",
    "cudaMemcpyAsync(h_clusters, d_clusters, D2H, sc.stream);   // 3. clusters out",
    "",
    "«cudaStreamSynchronize»(sc.stream);                          // 4. host stops",
    "for (c : found) out.push_back(c);                          // 5. copy out",
], title="ClusterFinderCUDAOpt2.hpp · find_clusters()")
# The picture goes BEFORE the floor callout, not after. It is a picture of the
# code panel directly above it -- three engines, one at a time, the host blocking
# between -- so it belongs next to the code it illustrates. The floor is then
# READ OFF that picture, and the 26 % in the rail is the consequence. Stating the
# floor first and showing the evidence afterwards had the argument backwards.
figure(s, "fig_opt1_timeline", M, 4.66, COL)
callout(s, M, 6.16, COL,
        "**PCIe is full-duplex**: H2D, kernel and D2H run on independent engines "
        "and overlap, so the **floor** is **max(H2D, kernel, D2H)**, never the sum. "
        "At 3×3 that is **16.2 µs → 61 859 FPS** [s4].", h=0.95, size=PT_LEAD)
rail(s, [
    ("label", "opt1 · 3×3 · 100 k frames · f64"),
    ("gap", 0.10),
    ("stat", "Throughput", "15 807 FPS", PALE),
    ("stat", "vs 24-thread CPU", "×2.34", ACCENT),
    ("row", "Per frame", "63.3 µs", TEXT2),
    ("gap", 0.10),
    ("row", "H2D · kernel · D2H  [s4]", "16.2 · 15.2 · 7.7 µs", TEXT2),
    ("row", "% of floor · 61 859 FPS", "26 %", AMBER),
    ("note", "[s1] opt1's own, uncontended: 13.1 · 14.7 · 5.3 µs"),
])
notes(s, """SAY THIS -- 144 words, about 53 s.

THIS PICTURE IS THE MAP FOR THE NEXT TEN SLIDES. Both ideas that drive the whole
ladder are visible as empty space in it. Point at the gaps while you say bullet 2.

- The first port is the simple one: copy a frame in, run the kernel, copy the
  clusters out. The calls do not block, but the last line does -- the clusters
  have to be copied into the caller's vector before the next frame starts. So
  the host stops, once per frame.

- Look at the picture: it holds the whole plan. There are two kinds of empty
  space in it. First: the three GPU steps run one after another, even though
  they use separate hardware and could overlap. Closing that gap is Act one.
  Second: the GPU waits for the host, then the host waits for the GPU. Closing
  that one is Act two.

- The first gap also gives the target: if the three overlapped perfectly, a
  frame would cost the longest of them, not the sum. 16 microseconds, 62 000
  frames a second. We are at 26 percent of it.

DETAILS: The two-idea framing, and where Act III fits.

The mapping is exact, so it is safe to promise here and collect later:

  ACT I   opt1-opt4   overlap the GPU with ITSELF.  H2D, kernel and D2H are
                      three independent engines; streams (opt2), removing the
                      barriers (opt3) and pinning (opt4) are all about getting
                      those three to run at the same time.
  ACT II  opt5-opt6   overlap the GPU with the HOST. Slide 18 is titled
                      host-GPU overlap; opt6 removes the last host copy.

Act III (opt7, the f32 kernel) is NOT a third kind of overlap, and do not
promise it as one. It makes one of the three engines cheaper rather than hiding
it behind another. That is also why it comes last: measured end to end it is
worth 1.5 %, and only once the host is off the critical path is the same change
worth 21 %. If someone asks "so what is act three?", the answer is "once nothing
is waiting any more, the only thing left is to make the work itself smaller".

DETAILS: Why the rail quotes [s4] on a slide about a ONE-STREAM build.

The floor is a property of the machine, not of opt1, and it has to be the SAME
denominator on all seven slides or "% of floor" means nothing. So every opt slide
divides by the floor derived from the s4 columns of A1: at 3x3 f64 the engines
are H2D 16.17, kernel 15.17, D2H 7.69 us, the max is H2D at 16.17, and
1/16.17 us = 61 859 FPS. That is the 16.2 in the callout.

The small [s1] line under it is opt1's own measurement -- one stream, nothing
else on the link -- and it is what the timeline figure is drawn from: 13.1 + 14.7
+ 5.3 = 33.2 us of serialized engine time inside a 63.3 us frame. Quote it if
someone asks what a single frame costs on a quiet GPU; do not divide by it. Its
max is the kernel, not H2D, and a floor built from it would move every
percentage on the next six slides.

Both conventions and the reconciliation are in annex A7.

DETAILS: What the code panel leaves out on purpose.

Two simplifications, both to avoid forward-referencing later slides. If someone
who has read the header asks, these are the honest answers:

  - opt1 does not DMA from the caller's array. It std::memcpy's the frame into
    sc.h_frame, a pinned buffer owned by the class, and the H2D reads from
    there. That extra host-side copy of 312 kB per frame is exactly what opt4
    deletes by registering the caller's buffer instead -- slide 15. Explaining
    pinned memory here would spend the idea four slides before it pays.

  - There are TWO D2H copies per frame, not one. The kernel writes a count, the
    host fetches the 4-byte count, waits, and only then can it ask for that many
    clusters -- so there is a second sync as well. That round trip is annex A3's
    whole subject, so this panel shows the payload copy only.

Neither changes the claim on this slide: the calls are async, the host still
stops once per frame, and nothing overlaps.""")

# =========================================================== 12 · OPT2
s = new_slide()
# NOT "2 000-frame batches". opt2 chops the batch into 500 rounds of four with a
# drain after each, so the pipeline is only ever FOUR frames deep -- the batch
# API is the vehicle that gets four streams, not a depth. Claiming 2 000 here
# left opt3 with nothing to add except a sync count. The depth story belongs on
# slide 14, where it becomes true.
chrome(s, 13, "Act I · streams",
       "Four streams, but only four frames deep: ×1.56", opt=2)
# WHAT A STREAM IS, then the picture of it, and only then the struct. The struct
# is an implementation detail of an idea the audience does not have yet; showing
# it before the picture made the slide open on field names. The figure is now the
# second thing on the slide, which is where the idea actually lands.
bullets(s, M, 1.82, COL, [
    "A **stream** is an ordered queue of GPU work. Work in **different** "
    "streams may overlap, so a copy can run while another stream computes.",
])
figure(s, "fig_opt2_timeline", M, 2.58, 7.60)
bullets(s, M, 4.40, COL, [
    "Each stream owns its **own buffers and pedestal**; frames go "
    "**round-robin**.",
])
# The struct is StreamContextOpt2 and it lives in the frozen opt1/opt2 header.
# The shipped ClusterFinderCUDA.hpp struct is a different shape -- no
# d_clusters, one fused uint8_t *d_output, plus d_pd_off -- so naming that file
# here pointed at a struct that does not have the fields on the slide.
code(s, M, 4.86, COL, [
    "struct StreamContextOpt2 {",
    "    cudaStream_t   stream;      FRAME_TYPE *h_frame;  // pinned",
    "    FRAME_TYPE    *d_frame;     ClusterType *d_clusters;",
    "    PEDESTAL_TYPE *d_pd_mean, *d_pd_sum, *d_pd_sum2;",
    "};",
    "auto &sc_k = v_sc[k];  // round-robin: frame = round*«n_streams» + k",
], title="ClusterFinderCUDAOpt2.hpp · per-stream state")
callout(s, M, 6.66, COL,
        "**Scaffolding, not yet the payoff.** The streams exist, but the host still "
        "synchronises after every round: see opt3.", h=0.58)
rail(s, [
    ("label", "opt2 · 3×3 · 4 streams · 4 deep"),
    ("gap", 0.10),
    ("stat", "Throughput", "24 726 FPS", PALE),
    ("stat", "vs 24-thread CPU", "×3.66", ACCENT),
    ("gap", 0.05),
    ("row", "Per frame", "40.4 µs", TEXT2),
    ("row", "Step gain over opt1", "×1.56", ACCENT),
    ("row", "% of floor · 61 859 FPS", "40 %  (was 26)", AMBER),
])
notes(s, """SAY THIS -- 116 words, about 42 s.

- A stream is just a queue of GPU work. Things inside one queue run in order.
  Things in DIFFERENT queues are allowed to overlap. So with four queues, a
  copy in one can run while another one computes.

- Each stream gets its own buffers and its own pedestal, and frames go to the
  four streams in turn. We also stop sending one frame at a time: 2 000 frames
  per call.

- This gives 1.56 times, and takes us from 26 to 40 percent of the floor. But
  be honest about what it is: it is only the scaffolding. The host still stops
  and waits after every round of four. Removing that is the next slide.""")

# =========================================================== 13 · OPT3
s = new_slide()
chrome(s, 14, "Act I · remove the sync barriers",
       "One sync per batch, not one per round: ×1.18", opt=3)
# The diagram is the argument on this slide, so it takes the whole left column
# and the two bullets go underneath it as the caption they effectively are.
figure(s, "fig_streams", M, 1.86, 8.05)
# BOTH SIDES CALL cudaStreamSynchronize. The earlier sketch here showed opt2
# calling cudaDeviceSynchronize, which is wrong -- that function appears nowhere
# in aare. opt2's barrier is the per-round DRAIN LOOP at Opt2.hpp:366-374, one
# cudaStreamSynchronize per stream, repeated every round; opt3 (kept commented
# at ClusterFinderCUDA.hpp:971-973) makes the same four calls ONCE per batch.
# What changed between them is where the loop closes, not which API is used.
code(s, 8.80, 1.86, 4.10, [
    "// opt2:  a drain after EVERY round of 4",
    "for (round) {",
    "   for (k : 4) submit(frame);",
    "   for (k : 4) «cudaStreamSynchronize»(sc[k]);",
    "}",
    "// opt3:  submit the whole batch, drain once",
    "for (frame : batch) {",
    "   cudaMemcpyAsync(..., sc.stream);",
    "   kernel<<<..., sc.stream>>>(...);",
    "   cudaMemcpyAsync(..., sc.stream);",
    "}",
    "for (k : 4) «cudaStreamSynchronize»(sc[k]);",
], title="THE BARRIER, BEFORE AND AFTER")
callout(s, 8.80, 4.86, 4.10,
        "**29 188 FPS · ×4.32**\n34.3 µs/frame · 47 % of floor (was 40)",
        h=0.86, size=PT_LEAD)
bullets(s, M, 5.92, 12.30, [
    "opt2 drained **all four streams after every round**. The host then "
    "copied four frames out, with every engine idle.",
    "opt3 submits every frame's H2D → kernel → D2H **asynchronously** and "
    "synchronises **once per batch**.",
], size=PT_BODY)
caption(s, M, 6.80, 12.30,
        "Each lane is one stream; H2D and D2H are one FIFO engine each. The red "
        "blocks are the host: 48 % of opt1's frame, 46 % of opt2's round. opt3 "
        "removes a second, per-frame barrier as well: annex A3.")
notes(s, """SAY THIS -- 122 words, about 44 s.

- Three pictures, same time axis. Top: version one, one stream, one engine
  working at a time. Middle: four streams, but the host says "wait for
  everybody" after every round of four -- you can see the copy engine go idle
  in the orange gaps. Bottom: the barrier is gone.

- The fix is small in the code and large on the picture. Instead of waiting
  after every round, we push all the work of the whole batch into the queues,
  and wait once at the end.

- 1.18 times faster, and 47 percent of the floor. The point is not the number,
  it is the shape: from here on the copy engine never has to stop and wait for
  the host.

CAREFUL: BOTH VERSIONS CALL THE SAME FUNCTION.

Do not say opt3 "replaced cudaDeviceSynchronize with cudaStreamSynchronize".
cudaDeviceSynchronize is not used anywhere in aare. Both sides call
cudaStreamSynchronize, once per stream. What changed is HOW OFTEN:

  opt2  ClusterFinderCUDAOpt2.hpp:324-388
        per round of 4: a launch loop, then a DRAIN loop that calls
        cudaStreamSynchronize on each of the 4 streams. With a batch of 2 000
        that is 500 rounds, so 2 000 sync calls per batch.

  opt3  ClusterFinderCUDA.hpp:896-990 (kept commented, and the comment there
        says these are the opt3/opt4 numbers)
        one launch loop over the whole batch, then the same 4 calls ONCE.

So the honest sentence is "opt2 drained the queues 500 times per batch, opt3
drains them once" -- an argument about the loop boundary, not about which API
was called. If someone asks why not cudaDeviceSynchronize: it would be worse
here, since it waits on every stream on the device rather than the four we
own.""")

# =========================================================== 14 · OPT4
s = new_slide()
chrome(s, 15, "Act I · pinned (page-locked) memory",
       "Pinning the input buys DMA-speed H2D: ×1.32", opt=4)
bullets(s, M, 1.82, COL, [
    "Normal host memory is **pageable**, so the driver stages every transfer "
    "through a **hidden pinned buffer**: copied twice.",
    "**Pinning** locks the pages in RAM, so the DMA engine reads host memory "
    "**directly**.",
], size=PT_BODY)
figure(s, "fig_pinning", M, 3.36, 8.10)
code(s, 8.80, 3.36, 4.10, [
    "// pin the whole dataset once",
    "«cudaHostRegister»(ptr, bytes,",
    "                 cudaHostRegisterDefault);",
    "",
    "// ... run the whole campaign ...",
    "",
    "«cudaHostUnregister»(ptr);",
], size=PT_CODE, title="ClusterFinderCUDA.hpp")
callout(s, 8.80, 5.30, 4.10,
        "**38 486 FPS · ×5.69**\n26.0 µs/frame · 62 % of floor, the largest step in Act I",
        h=1.00, size=PT_LEAD)
callout(s, M, 6.44, 12.30,
        "**Measured H2D:** one 312.5 KiB frame in **13.2 µs = 24.2 GB/s**, which is "
        "77 % of PCIe 4.0 ×16: true DMA speed. **Pageable staging manages ~15 GB/s** "
        "for the same frame, because it is copied twice.", h=0.78, size=PT_LEAD,
        color=AMBER)
notes(s, """SAY THIS -- 124 words, about 45 s.

- Normal host memory can be moved around by the operating system. The GPU
  cannot copy from memory that might move, so the driver does it in two steps:
  first into a hidden safe buffer, then to the GPU. Every frame is copied twice.

- Pinning tells the operating system: never move these pages. Then the GPU
  reads host memory directly. One line of code, once, for the whole dataset.

- And it is the biggest single step in act one: 1.32 times, 62 percent of the
  floor. The copy now runs at 24 gigabytes per second, which is about 77 percent
  of what the PCIe wire can do -- the two-step version managed 15. Note this
  helps because the copy IN is the tallest bar here. At 9x9 it buys almost
  nothing.

DETAILS: Which H2D number, and why there are two.

13.2 us is uncontended [s1]. In the shipped pipeline H2D reads 16.6 us [s4]; the
extra 26 % is H2D<->D2H contention on the link, and the full engine grid is in
annex A7.

The rule this slide is the first sighting of: pinning attacks H2D, which is the
tallest bar at 3x3 and the shortest at 9x9. So the same change is worth x1.32 at
3x3 and only x1.03 at 9x9.""")

# -------------------------------------------------------- divider · ACT II
s = section("Act II of III · get the results back",
        # NOT "the host-to-host result copy": "host-to-host" is one unbreakable
        # 8.2 em token, so _fit had to drop the whole title to 28 pt while every
        # other divider sits at 36. Splitting it into three words costs nothing
        # and puts this title back on the same size as its siblings.
        "Bottleneck: the result copy, host to host",
        "Both transfers are DMA already. What is slow is what happens AFTER the "
        "D2H: copying clusters out of the finder's pinned buffer into heap the "
        "caller owns. Host to host. Still 3×3, but 9×9 is where this act pays most.",
        [(16, "opt5 · host↔GPU overlap"),
         (17, "9×9 · why overlap runs out"),
         (18, "opt6 · zero-copy")],
        rng=(16, 18), col=PALE,
        carry=("Arriving at", "38 486 FPS",
               "opt4 · 26.0 µs per frame · 62 % of the GPU floor"))
notes(s, """SAY THIS -- 65 words, about 23 s.

- Act two. Both copies now run at full wire speed, so the wire is no longer the
  thing to fix.

- What is slow now happens AFTER the results come back: the host copying
  clusters out of our buffer into memory the caller owns. Host to host, no GPU
  involved.

- Three slides. It pays a little at 3x3, and a lot at 9x9.""")

# =========================================================== 15 · OPT5
s = new_slide()
chrome(s, 16, "Act II · host↔GPU overlap",
       "Overlapping host and GPU hides min(host, GPU): ×1.31", opt=5)
bullets(s, M, 1.82, COL, [
    "opt3 overlapped the GPU's own engines, never the **host** with the GPU: "
    "the batch synchronised, then built ClusterVectors with the GPU idle.",
    "opt5 keeps **one batch in flight while materialising the previous one**.",
])
figure(s, "fig_overlap", 0.42, 2.94, 8.50)
code(s, M, 5.70, COL, [
    "tok = cf.«submit_batch»(data[a0:b0], first_frame=a0)",
    "for a, b in bounds[1:]:",
    "    nxt = cf.«submit_batch»(data[a:b], first_frame=a)   // GPU starts i+1",
    "    results.extend(cf.«collect»(tok))                   // host unpacks i",
    "    tok = nxt",
], size=PT_CODE, title="you never write these: find_clusters_batched() wraps them")
notes(s, """SAY THIS -- 118 words, about 43 s.

- Up to now we overlapped the GPU with itself. We never overlapped the GPU with
  the HOST. The batch finished, and then the host built its result objects while
  the GPU sat and did nothing.

- The fix is to keep one batch running on the GPU while the host is still
  unpacking the previous one. Send batch two, then unpack batch one, and repeat.

- Note what is NOT here: no new CUDA call, no kernel change. Only the order of
  the host loop. That is 1.31 times faster, and it takes us to 81 percent of the
  floor. And the user does not even see it -- it is inside the normal function
  everyone already calls.

- Last line on the right, and then we move on: at the big cluster size the same
  change only reaches 45 percent, not 81. That gap is the whole next slide.

DETAILS: opt5 — host<->GPU overlap. Code and chunk sizing are on annex A4.

    tok = cf.submit_batch(data[a0:b0], first_frame=a0)
    for a, b in bounds[1:]:
        nxt = cf.submit_batch(data[a:b], first_frame=a)   # GPU starts N+1 ...
        results.extend(cf.collect(tok))                   # ... host unpacks N
        tok = nxt
    results.extend(cf.collect(tok))                       # drain

1. This is now INTERNAL to find_clusters_batched(), so every caller gets it for
   free. submit_batch/collect stay public for anyone who wants the token by hand.

2. The chunk size is rounded to a multiple of n_streams, because the device
   pedestal is per-stream: an uneven chunk would advance the four stream
   pedestals by different amounts and the finders would stop being comparable.

Why the gain differs by cluster size: the saving is min(GPU, host) per chunk, so
it is largest when the two terms are comparable. At 3x3 they nearly are (GPU
16.2 us, host ~9.8) and opt5 is worth x1.31. At 9x9 the host term is roughly
twice the GPU term, so overlap hides only the smaller one and opt5 is worth
x1.20 - which is exactly the diagnosis that motivates opt6: you cannot overlap
your way out of a host term that is simply larger. Report SS8.2.

CAREFUL: THE 45 % ON THE RAIL IS NOT A STABLE NUMBER.

Both rail blocks come from the same ladder table (SS8, [f64]), but they are not
equally solid, and the difference is worth knowing before someone asks:

    3x3   19.84 us   spread  1.2 %    clean, 30-96 k minor faults
    9x9   66.39 us   spread 22.4 %    151 601 minor faults, best of five reps

opt5 at 9x9 is THE ONE ROW IN THE LADDER THAT NEVER REACHES A FAULT-FREE STEADY
STATE (SS8.3). 66.39 is the best rep, not a plateau. So quote 45 % as "about 45"
and do not defend the second figure.

This does not weaken the argument -- it strengthens it. The faults are the host
allocating and first-touching a fresh ClusterVector per frame, which is exactly
the cost opt6 deletes. The instability IS the symptom. opt6 at 9x9 lands on
30.01 us, 100 % of the floor, with the fault count back to normal.""")
# The 9x9 row goes LAST, on purpose. This slide's argument is 3x3 -- x1.31 for
# no CUDA work -- and that has to land whole before it is qualified. Putting 9x9
# at the bottom makes it the hand-off: the last thing read on the slide is the
# number the NEXT slide exists to explain. 45 % is amber against 81 %'s accent,
# so the drop is visible before anything is said about it.
rail(s, [
    ("label", "opt5 · f64 · no CUDA work at all"),
    ("gap", 0.10),
    ("stat", "3×3 throughput", "50 410 FPS", PALE),
    ("row", "per frame · step · of floor", "19.8 µs · ×1.31 · 81 %", ACCENT),
    ("gap", 0.16),
    ("label", "the whole change"),
    ("gap", 0.06),
    ("row", "CUDA API calls added", "none", TEXT2),
    ("row", "what moved", "the host loop", PALE),
    ("gap", 0.20),
    ("row", "9×9 · same change, half the ceiling", "66.4 µs · ×1.20 · 45 %", AMBER),
])


# ============================================ 18 · WHY OVERLAP RUNS OUT AT 9x9
# The bridge from opt5 to opt6. Slide 17 is told at 3x3, where the host copy fits
# underneath the GPU floor and two slots are plainly enough. At 9x9 the host is
# the taller bar and the room reliably guesses "add more slots" -- so the picture
# answers that guess directly, by drawing three slots and landing on the same
# finish line. Once buffering is ruled out, opt6 is the only move left.
s = new_slide()
chrome(s, 17, "Act II · at 9×9 · why overlap runs out",
       "Overlap runs out: the host is the taller bar", opt=5)
bullets(s, M, 1.82, 12.30, [
    "At 3×3 the host copy is **shorter than the GPU floor** and hides under it. "
    "At 9×9 it is **roughly twice the floor**, so overlap has less to hide: "
    "opt5 is worth **×1.20** here against ×1.31 at 3×3.",
], size=PT_BODY)
figure(s, "fig_overlap_9x9", 0.66, 2.50, 11.68)
callout(s, M, 6.30, 12.30,
        "A deeper buffer **relocates the GPU's idle, it does not close it**: the "
        "host lane is already back-to-back in both strips, so it alone sets the "
        "pace. **The only way down is to make the host term smaller.**",
        h=0.66, size=PT_LEAD)
caption(s, M, 7.06, 12.30,
        "Measured proportions: GPU 30.01 µs/frame, host ~62 µs steady-state [A10].",
        size=PT_META)
notes(s, """SAY THIS -- 121 words, about 44 s.

- At the small cluster size the host work is shorter than the GPU work, so it
  hides underneath it completely. At 9x9 the host work is about twice the GPU
  work. Now it sticks out, and overlap can only hide the smaller of the two.

- The natural idea is: use more buffers, go deeper. The lower picture is
  exactly that -- three slots instead of two -- and the finish line does not
  move. More buffers move the GPU's idle time somewhere else; they do not
  remove it.

- The reason is simple: the host lane is already busy back to back. It sets the
  pace. The only way to go faster is to make the host work itself smaller. That
  is the next slide.

DETAILS: This slide exists because "add more slots" is the reliable guess here,
and it is worth letting the room say it out loud before the second strip goes up.

The queueing argument, if it is asked for: with producer period G and consumer
period H, an N-buffer pipeline has steady-state period max(G, H) for every N >= 2.
Buffers DECOUPLE two stages; they do not speed up either. Depth beyond 2 only
helps when the periods vary -- it absorbs jitter, at the cost of latency and
pinned memory. Here the work per chunk is near constant, so there is no jitter to
absorb.

The concrete version lands better: a third slot needs somebody to fill it, and
the only host thread is inside collect(). Adding slots without adding a producer
thread is adding storage to a queue that is not storage-bound.

"So would a producer thread help?" It would let submit and collect overlap, but
the host term is dominated by malloc + first touch, which is allocator-serialised
anyway. That was measured and reverted (ClusterFinderCUDA.hpp:130-136): 2.27 M
faults at 8 threads against 9.7 k at 1, a 6 % gain at best and a 33 % LOSS when
results are freed promptly. The comment there ends with the right conclusion --
stop allocating per frame, do not copy faster. That is opt6.

ON THE 62 US, if challenged. It is not measured directly; no one timed the host
loop. Two independent routes agree on it. (1) Fault-correct the five f64 reps at
0.68 us/fault -- the same rate fitted at 3x3 and applied out of sample -- and a
22 % raw spread collapses to 4.6 %, at 61-64 us. (2) f32 rep 3 happened to run
with only 10 128 faults and measured 61.85 us with no correction at all. The
"~40 us" that used to be on the opt6 slide was the memcpy alone, computed at
bandwidth; the loop is allocation-bound, not bandwidth-bound, so it under-counted
the host term by about half.""")

# =========================================================== 19 · OPT6
s = new_slide()
chrome(s, 18, "Act II · zero-copy collection",
       "Read the results in place: ×2.21 at 9×9", opt=6)
bullets(s, M, 1.82, COL, [
    "The D2H lands in a **pinned host buffer**; collect() then allocates one "
    "ClusterVector per frame and memcpys into it: **467 kB per frame at 9×9**.",
    "collect_view() returns a **BatchView** onto that buffer instead: it "
    "withholds **ownership** past the chunk, not access.",
], size=PT_BODY)
figure(s, "fig_resultpath", 0.42, 3.36, 8.50)
callout(s, M, 6.32, COL,
        "The win is **max(0, host copy − GPU floor)**. At 3×3 the copy already hid "
        "under the floor; at 9×9 it is **twice the floor** and cannot hide at any "
        "overlap.", h=0.80, size=PT_LEAD)
notes(s, """SAY THIS -- 126 words, about 46 s.

- So what IS the host doing? The results arrive in our own pinned buffer. Then,
  for every frame, we allocate a new object and copy the clusters into it -- 467
  kilobytes per frame at 9x9. Allocate, copy, free, sixty thousand times.

- The fix is to stop copying. We hand back a view: a small object that points
  INTO the buffer that already holds the data. You can read every cluster. The
  only thing you do not get is ownership past the end of your chunk.

- At 9x9 that is 2.2 times faster on its own, and we land exactly on the GPU
  floor -- 100 percent. At 3x3 it is only 1.16, because there the copy was
  already hidden. Same change, very different value.

DETAILS: Where 93 kB and 467 kB come from: clusters/frame x sizeof(Cluster).

sizeof — Cluster.hpp:28 is two CoordType coords then std::array<T, X*Y> data.
CoordType=uint16, T=int32, alignof 4, so the coords pack with no padding:
    3x3:  2x2 B +  9x4 B =  40 B
    9x9:  2x2 B + 81x4 B = 328 B

clusters/frame — report section 13, the same counts used for correctness:
    3x3:  233 094 390 / 100 000 fr = 2 330.9 /fr
    9x9:   28 447 962 /  20 000 fr = 1 422.4 /fr

    2 330.9 x  40 B =  93.2 kB/frame
    1 422.4 x 328 B = 466.5 kB/frame  ->  x 20 000 fr = 9.33 GB/run

NOT the D2H transfer. D2H is fixed at cap x sizeof regardless of how many clusters
were found: at the 9x9 campaign cap of 1700 that is 1700 x 328 B = 558 kB every
frame, i.e. the 467 kB of real clusters is an 84 % fill (section 4.2). The 467 kB
is the HOST memcpy inside collect(), pinned buffer -> freshly allocated
ClusterVector. That is what collect_view() removes; the wire traffic is unchanged.

Report: section 12 (payload table), section 9.4 (materialize_slot), section 13.""")

rail(s, [
    ("label", "opt6 · 3×3 and 9×9 · collect_view()"),
    ("gap", 0.10),
    ("stat", "3×3 throughput", "58 495 FPS", PALE),
    ("row", "per frame · step · of floor", "17.1 µs · ×1.16 · 95 %", ACCENT),
    ("gap", 0.14),
    ("stat", "9×9 throughput", "33 323 FPS", PALE),
    ("row", "per frame · step · of floor", "30.0 µs · ×2.21 · 100 %", AMBER),
    ("gap", 0.12),
    ("row", "Results", "bit-identical to opt5", TEXT2),
])

# ------------------------------------------------------- divider · ACT III
s = section("Act III of III · the kernel",
        "Bottleneck: the kernel itself",
        "The story moves to 9×9, where the kernel is finally the tallest bar, and "
        "where most of what it does is update the pedestal, not find photons.",
        [(19, "why the kernel, and where its time goes"),
         (20, "opt7 · FP32 pedestal"),
         (21, "why it comes last")],
        rng=(19, 21), col=AMBER,
        carry=("Arriving at", "58 495 FPS",
               "opt6 · 3×3, and 33 323 FPS at 9×9, where this act pays"))
notes(s, """SAY THIS -- 64 words, about 23 s.

- Act three. Only now is the kernel itself the slowest part -- and only at the
  big cluster size, 9 by 9.

- This is the one act where I change the maths instead of the plumbing.

- First, the measurement that says the kernel is now the problem -- and what it
  is spending that time on, which turns out not to be photons.

- Three slides. The trap that the change brings with it -- and it is a real
  one -- is annex ten.""")

# =========================================================== 19 · WHY THE KERNEL
# Act III used to open on a definitions slide. It now opens on the MEASUREMENT
# that motivates the act -- the s1 reading where the kernel is finally the
# tallest bar -- and on what the kernel actually spends that time doing. s1 is
# introduced in one parenthesis; the convention itself is annex A8, which is
# where it belongs once the arc stops needing it to make an argument.
s = new_slide()
chrome(s, 19, "Act III · why the kernel, and where its time goes",
       "The kernel is the tallest bar, and 80 % of it is pedestal")
bullets(s, M, 1.82, 12.30, [
    "At 9×9, **one stream with nothing else running**, the kernel takes "
    "**39.9 µs** against **22.0** for the copy out and **13.2** for the copy in. "
    "For the first time in the talk it is the thing to fix.",
    "**~80 % of pixels see only noise**, and each takes the **pedestal-update** "
    "branch: six running values read and written. Most of the kernel is "
    "**moving pedestal numbers**, not finding photons.",
], size=PT_BODY)
# 3.26, not 3.10: two body bullets of two lines each end at 3.04 once the
# renderer's own line height and the inter-bullet gap are counted, and the
# figure is an opaque PNG -- it does not collide with the second bullet, it
# paints over it, which is why the layout audit could not see this one.
figure(s, "fig_kernel_binds", M, 3.26, 12.30)
callout(s, M, 6.52, 12.30,
        "**So attack the pedestal, not the photon search.** The branch that runs "
        "on four pixels in five is the one worth making cheaper.",
        h=0.62, size=PT_LEAD)
notes(s, """SAY THIS -- 118 words, about 43 s.

- Act three, and for the first time the GPU itself is the problem. Left picture:
  one stream, nothing else running, the big cluster size. The kernel is 40
  microseconds. The two copies are 22 and 13. The kernel is nearly twice the
  next tallest bar.

- So what is the kernel actually doing? Right picture: one frame, finished,
  coloured by what happened to each pixel. Green is a photon we stored. The dark
  cells are pixels that sit inside some photon's window. And amber -- about four
  pixels in five -- is "nothing here", which means: update the pedestal.

- So the kernel is mostly not finding photons. It is bookkeeping noise. That is
  what act three attacks.

DETAILS: the two numbers on this slide, and where they come from.

The engine bars are the UNCONTENDED reading: one stream, nothing else on the
device, so each bar is a true duration. 9x9, cap 1 700, f64 build:
H2D 13.20, kernel 39.86, D2H 21.97 us. In the shipped four-stream pipeline the
numbers move -- copies rise, the kernel falls, because four kernels overlap each
other -- and that distinction is annex A7. It does not change this slide's
claim: the kernel binds in BOTH readings on the f64 arm.

The ~80 % is measured, not assumed: the branch map is read out of the finder
after a frame completes (branch_site_dump.py). Of a 21x21 neighbourhood, ~1.5 %
of pixels are stored photons, ~18 % sit in a photon's shadow -- inside some
window whose max clears 5 sigma -- and ~80 % clear nothing at all and push the
pedestal. The same picture returns on slide 31 to argue about Test3; here only
the amber fraction matters.

WHY THAT MAKES THE KERNEL BANDWIDTH-BOUND. The pedestal update is six values per
pixel, read and written, and almost every pixel does it. There is very little
arithmetic per byte moved. So the kernel's cost is dominated by pedestal
TRAFFIC, which is what the next slide halves.""")


# =========================================================== 20 · OPT7
# Two panels, and only two: the s1 kernel pair, which IS the -40.5 % claim, and
# the s4 grid, which is the only place the handover to D2H can be read. The s1
# copy engines are not on this slide -- they do not move under the typedef and
# they were the third of the picture that carried no argument. Full reading: A9.
s = new_slide()
chrome(s, 20, "Act III · FP32 device pedestal",
       "One typedef: −41 % on the kernel, and the floor moves to D2H", opt=7)
bullets(s, M, 1.82, 7.75, [
    "One typedef halves all four pedestal arrays (**48 bytes per updating "
    "pixel in FP64, 24 in FP32**), and the kernel is **bandwidth-bound**.",
], size=PT_BODY)
figure(s, "fig_f32_kernel_s1s4", M, 2.86, 8.16)
caption(s, M, 6.26, 8.16,
        "Left: one stream, nothing else running, so the bar is a true duration and "
        "carries the −40.5 % claim. Right: the shipped four-stream pipeline, "
        "where the f32 kernel falls under a D2H bar that never moved. Full "
        "engine grid: A8.",
        size=PT_META)
code(s, 8.80, 1.82, 4.10, [
    "// clusterfinder_kernel.cuh",
    "using COMPUTE_TYPE    = float;",
    "using DEVICE_PED_TYPE = «float»;",
    "//            was: double",
], size=PT_CODE, title="ONE TYPEDEF")
callout(s, 8.80, 3.22, 4.10,
        "Kernel, 9×9 **uncontended**\n**39.86 → 23.70 µs  (−40.5 %)**",
        h=0.94, size=PT_LEAD)
callout(s, 8.80, 4.30, 4.10,
        "**The bottleneck moves.** In the shipped pipeline the kernel lands at "
        "**23.94 µs**, under a **25.24 µs** D2H. The limit is now the result "
        "path, not the maths.", h=1.30, size=PT_LEAD, color=AMBER)
callout(s, 8.80, 5.74, 4.10,
        "Naive FP32 is **wrong**. See **annex A9**.", h=0.56, size=PT_LEAD,
        color=RED)
caption(s, 8.80, 6.48, 4.10,
        "Both builds, same git rev, 20 000 frames.", size=PT_META)
notes(s, """SAY THIS -- 120 words, about 44 s.

- Remember that four pixels in five only update the pedestal, and that update
  reads and writes six running values. Those were stored as doubles, 8 bytes
  each.

- Change them to floats. One typedef. 48 bytes per pixel becomes 24. The kernel
  is limited by memory traffic, so half the bytes is a big win: 40 percent off
  the kernel.

- Left bar pair is the honest measurement of the change: one stream, nothing
  else running, 39.9 down to 23.7.

- The right picture is the ending. In the real pipeline the kernel drops to
  23.9, and the copy back out -- 25.2, unchanged -- is now the tallest bar. We
  optimised the kernel until it stopped being the problem.

CAREFUL: TWO DIFFERENT PERCENTAGES, AND THEY MEASURE DIFFERENT THINGS.

-40.5 % is the uncontended kernel: one stream, a true duration, and the honest
headline for what the typedef does to the arithmetic. -26.7 % is the same change
read in the shipped four-stream pipeline, where the kernel's per-frame busy time
is already shortened by overlapping itself, so part of the win was hidden before
the typedef arrived. Both are true. A7 is the rule for which to quote, and A8
has the full grid including the copy engines and the 3x3 arm.

THE 3x3 ARM, if asked: the same typedef is worth -70.6 % in the kernel there and
4.6 % end to end, because at 3x3 the kernel was never the tallest bar -- H2D
was. It also frees 9 registers, taking 3x3 from five blocks per SM to six. Neither
changes anything end to end.

WHY FP64 WAS SLOW FOR TWO REASONS, not one: the bytes are the main story, but on
a GeForce part FP64 arithmetic also runs at 1/64 of FP32. The bytes are why it
is worth 40 %; the arithmetic rate is a bonus.""")


# =========================================================== 21 · OPT7 when
# Back on the main arc. It is the only slide that says WHY the acts are in the
# order they are in, and it is where the talk hands the floor from the kernel
# back to D2H -- an ending, not a detail. The allocator-state caveat behind the
# excluded opt3 arm and the daggered opt5 arm stays in annex A11.
s = new_slide()
chrome(s, 21, "Act III · why this act comes last",
       "opt7 saves the same 4.7 µs at every step: 6 % of opt4, 16 % of opt6",
       opt=7)
figure(s, "fig_f32_absolute", M, 1.82, 12.30)
# The callout carries opt7'S OWN RUNG, which no slide in the deck stated: every
# other step gets "N FPS / x1.NN / % of floor" and opt7 got only kernel numbers.
# 25.14 us/frame IS the 9x9 f32 floor (A8), so the ladder closes at 100 %.
callout(s, M, 5.84, 12.30,
        "**opt7 ships at 25.1 µs/frame, 100 % of the 9×9 floor.** And the act "
        "ends by handing that floor away: at **s4** the f64 kernel bound the "
        "pipeline at **32.66 µs**; the typedef puts it at **23.94**, under a "
        "**25.24 µs** D2H that never moved.",
        h=0.80, size=PT_LEAD, color=AMBER)
caption(s, M, 6.72, 12.30,
        "9×9 · 20 000 frames · 4 streams · cap 1 700 · warm · both arms at the same "
        "git rev. opt3 excluded and opt5 marked †: their arms sat in different "
        "allocator states. Worked through in annex A10.", size=PT_META)
notes(s, """SAY THIS -- 123 words, about 45 s.

- One last question about act three: why is it last? Because its value depends
  on everything before it.

- Left picture: the whole frame time shrinking, step by step, from 80
  microseconds down to 25. Right picture: what the typedef itself saves at each
  of those steps. It is the same 4.7 microseconds every time. The saving never
  grew. What changed is the frame around it: 4.7 out of 80 is six percent; the
  same 4.7 out of 30 is sixteen.

- And that is where the talk ends, because act three hands the limit away. Once
  the kernel drops to 24 microseconds, the copy back out, at 25, is the tallest
  bar again. There is no opt8.

DETAILS: WHY opt3 IS EXCLUDED AND opt5 IS DAGGERED. Both steps' two arms sat in
different allocator states, so the difference between them does not report the
typedef -- it reports the heap. opt3 is dropped entirely; opt5's reading agrees
with opt4 and opt6 but is not independently attributable. Annex A10 has both.

WHAT THE BARS ARE. The f32 bar at opt6 is opt7, the shipped build. The f32 bars
at opt4 and opt5 are the same typedef applied at earlier steps -- configurations
that exist only to make this comparison controlled, and that nobody ships.

WHY opt4 AND opt6 ARE THE PAIR TO QUOTE. They are the two readings whose fault
counts match, so the difference between their arms is the typedef and nothing
else. That is the whole reason the saving can be called constant.

THE HANDOVER, if asked for the numbers: [engines, slide 20 of the 33-slide
draft] f64 s4 kernel 32.66 vs D2H 25.25 -> kernel-bound; f32 s4 kernel 23.94 vs
the same 25.25 -> D2H-bound. The constraint moved from the arithmetic to the
result path, which is why Act III is last and why there is no opt8.""")

# ------------------------------------------------------- divider · results
s = section("Results · what came out of it",
        "The full ladder, both cluster sizes",
        "Both cluster sizes end to end, and the audit behind the numbers.",
        [("22–23", "Results, both cluster sizes"),
         ("24", "Where the time went")],
        rng=(22, 24), col=PALE,
        carry=("Arriving at", "58 495 FPS",
               "opt6 · everything after this is the ladder seen whole"))
notes(s, """SAY THIS -- 58 words, about 21 s.

- That is all seven steps. The rest of the talk is the ladder seen as a whole:
  both cluster sizes end to end, and where the time actually went.

- No new optimisation from here. Only results -- and then whether they are the
  RIGHT results, which is the last section.""")

# =========================================================== 22 · RESULTS
s = new_slide()
chrome(s, 22, "Results · 3×3", "×9.1 at 3×3, sitting on the H2D floor")
figure(s, "fig_arc", 1.79, 1.76, 9.71)
callout(s, M, 5.94, 6.04,
        "**×9.1 over 24 CPU threads**, 14.8 s → 1.63 s for 100 000 frames, "
        "and **at the H2D floor**.", h=0.8)
callout(s, 6.75, 5.94, 6.04,
        "Every step is **monotonic**, and correctness held constant throughout: "
        "**exact on the f64 pedestal**, **6 clusters in 23 M** on the shipped f32.", h=0.8, color=AMBER)
caption(s, M, 6.96, 12.30,
        "3×3 · nσ = 5 · 100 000 frames · batch 2 000 · 4 streams · warm = best of "
        "reps 1–4 · each step in its own process.")
notes(s, """SAY THIS -- 112 words, about 41 s.

- The whole ladder at the small cluster size. Every bar is a step, and every
  step is an improvement -- nothing had to be undone later.

- Left to right: 6 762 frames per second on 24 CPU threads, and 58 495 at the
  end. Nine times faster. For a hundred thousand frames that is 15 seconds down
  to 1.6.

- The dashed line is the hardware floor, and the last bar is sitting on it. That
  is the honest way to read this: we are not fast, we are FINISHED. What is left
  is the wire. And the answers stayed correct the whole way -- six clusters
  different out of 23 million.""")

# =========================================================== 23 · RESULTS 9x9
s = new_slide()
chrome(s, 23, "Results · 9×9", "×26.5 at 9×9, and opt7 hands the floor to D2H")
figure(s, "fig_arc_9x9", 1.79, 1.76, 9.71)
callout(s, M, 5.94, 12.30,
        "**×26.5 over 32 CPU threads**; the kernel is the tallest bar for the whole "
        "f64 arm, and opt7's −40 % drops it **below D2H**.", h=0.72)
caption(s, M, 6.86, 12.30,
        "9×9 · cap 1 700 · 20 000 frames · CPU baseline at 32 threads. opt1 and "
        "opt2 are 3×3-only. Why the two arms differ is in the notes.")
notes(s, """SAY THIS -- 109 words, about 40 s.

- Same ladder, big cluster size. Here the numbers are much larger: 26 times
  faster than 32 CPU threads.

- Why bigger? Because the CPU suffers much more when clusters are 9 by 9 -- it
  does 81 pixels of work where the GPU just uses more threads. The GPU work per
  frame grows far less.

- And this is the arm where the story ends differently. For the whole f64 part
  the kernel is the tallest bar. The typedef change cuts it by 40 percent, and
  then it drops BELOW the copy back out. From that point the result path is the
  limit, not the maths.""")

# =========================================================== 24 · WHERE TIME GOES
s = new_slide()
chrome(s, 24, "Where the time actually went",
       "The host bar dies first, then the floor itself drops")
figure(s, "fig_overhead", 1.79, 1.82, 9.71)
callout(s, M, 5.09, 6.04,
        "**Acts I and II never touch the arithmetic.** The GPU floor is a flat "
        "16.2 µs at 3×3 / 30.0 µs at 9×9; what collapses is everything stacked on it.",
        h=0.86, size=PT_LEAD)
callout(s, 6.75, 5.09, 6.04,
        "**Act III is the only step that lowers the floor itself**, and it could not "
        "have been seen until the stack above it was gone.", h=0.86, size=PT_LEAD,
        color=AMBER)
caption(s, M, 6.24, 12.30,
        "Coloured = the GPU floor for that act's build. Grey = everything the host "
        "adds on top: +50 µs at opt3, nothing at opt6.")
notes(s, """SAY THIS -- 116 words, about 42 s.

- This is the same story told as two colours. The coloured part is the hardware
  floor -- what the GPU cannot go below. The grey part is everything the host
  piles on top of it.

- Look at what actually happened. In acts one and two the coloured part never
  moves. It is flat: 16 microseconds at 3x3. Every step is eating grey. By the
  end of act two there is no grey left.

- Only act three moves the coloured part itself, and only there. That is the
  reason for the order of the whole talk: you cannot see the floor, and you
  certainly cannot lower it, while fifty microseconds of host work sit on top
  of it.""")

# --------------------------------------------------- divider · VALIDATION
s = section("Validation · does it find the same photons",
        # Kept short on purpose: _fit scales the title to the column, and a long
        # one drops to 24 pt while every other divider gets 32-36. A divider that
        # whispers reads as a lesser beat than it is.
        "Does it find the same clusters?",
        "Whether the CUDA finder returns the same clusters as the CPU.",
        [("25–26", "The fair comparison"),
         ("27–28", "The f32 residual, dissected"),
         ("29–31", "The CPU-vs-CPU gap, named"),
         ("32–33", "For users"),
         ("34", "What is next")],
        rng=(25, 34), col=PALE,
        carry=("Established", "×9.1 and ×26.5",
               "on the hardware floor at both cluster sizes, if the physics holds"))
notes(s, """SAY THIS -- 60 words, about 22 s.

- Everything so far assumed one thing: that the GPU finds the same photons as
  the CPU. If it does not, none of the numbers mean anything.

- So the last part is correctness. How do you compare two finders fairly, what
  is left over when you do, and what a user should actually call.""")

# ============================================ 26 · PEDESTAL UPDATE TIMING
s = new_slide()
chrome(s, 25, "Validation · why a CPU twin was needed",
       "CPU and CUDA update the pedestal at different moments")
figure(s, "fig_pedtiming", M - 0.15, 1.82, 12.61)
callout(s, M, 5.98, 6.04,
        "The CPU finder updates the pedestal **as the raster reaches each pixel**. "
        "160 000 CUDA threads read it at once, so the update is **applied at the "
        "frame boundary**.", h=0.86, size=PT_LEAD)
callout(s, 6.75, 5.98, 6.04,
        "A straight CPU↔CUDA comparison therefore moves **two** things at once. "
        "**ClusterFinderFrozen** moves only the update: same arithmetic, same "
        "gates, same scan.", h=0.86, size=PT_LEAD, color=AMBER)
caption(s, M, 7.02, 12.30,
        "Frozen is a diagnostic twin, not a product. Every finder trains on the "
        "same 1 000 pedestal frames and runs the same 10 000 data frames.")
notes(s, """SAY THIS -- 124 words, about 45 s.

- To compare two finders you have to change ONE thing at a time. Here is why
  that is hard.

- The CPU finder updates the pedestal as it walks the frame, pixel by pixel. So
  a pixel near the end of the frame is judged against a pedestal that the
  earlier pixels of the SAME frame have already moved.

- The GPU cannot do that. All 160 000 threads read the pedestal at the same
  moment, so the update can only be applied at the end of the frame.

- So a straight CPU-against-GPU comparison changes two things at once: the
  maths, and WHEN the pedestal moves. That is why I wrote a third finder --
  same CPU code, but with the update moved to the frame end.

DETAILS: Why this slide is here.

The obvious experiment - run the CPU finder and the CUDA finder over the same
frames and count the differences - cannot answer the question anyone actually
asks, which is "is the port correct?" It moves two variables at once:

  1. WHEN the pedestal is updated. The serial CPU finder updates per pixel,
     during the raster scan, so a pixel late in the frame is judged against a
     pedestal that already contains this frame's earlier pixels. That is
     scan-order dependent by construction.
  2. WHAT arithmetic runs. float vs double, a different local-max expression,
     a different rounding point.

ClusterFinderFrozen holds (2) fixed and changes only (1): it is the serial
finder, same decision logic line for line, with the update deferred to the frame
boundary. That is the CUDA update model, on the CPU.

So the three-way comparison factorises:
    cpu    vs frozen  -> update timing alone      (19 of 23.2 M)
    frozen vs cuda    -> the port alone           (6 of 23.2 M)

and the second is the number that answers the question. It is also why the
headline mismatch figure in this deck is 6 / 23 M and not 25 / 23 M: the larger
number is dominated by an effect that has nothing to do with CUDA.

Frozen ships in the library as a diagnostic, not as the recommended finder.""")

# =========================================================== 26 · CORRECTNESS
s = new_slide()
chrome(s, 26, "Validation · isolating one variable at a time",
       "CUDA and its CPU twin agree exactly: 0 in 23 million")
bullets(s, M, 1.82, 12.40, [
    "**ClusterFinderFrozen** makes byte-for-byte the same decisions as ClusterFinder "
    "and differs in exactly one thing: **when** the pedestal is updated. Frozen per "
    "frame, pushed at frame end. That is the CUDA model, so comparing against it "
    "isolates everything else the port changes.",
], size=PT_BODY)
# ONLY THE TWO ROWS THAT MEASURE THE PORT. The table used to carry four, and the
# two serial-CPU rows were read as the weak half of it -- when in fact they are a
# separate and equally strong claim, which now has slide 29 to itself.
table(s, M, 3.42, 12.40,
      ["comparison", "the one thing that differs", "A-only / B-only",
       "% of clusters"],
      [["frozen CPU  vs  CUDA  [f64 ped]", "**nothing**", "**0 / 0**", "**0 %**"],
       ["frozen CPU  vs  CUDA  [f32 ped]", "the float32 pedestal EMA drifting: "
        "see next slide", "0 / 6", "0.000026 %"]],
      colw=[0.28, 0.40, 0.17, 0.15], size=PT_TABLE, rowh=0.60)
# The row that measures the PORT with the timing held fixed. Ringed rather than
# recoloured, and drawn after the table so the outline sits over the zebra.
frame_rect(s, M - 0.04, 3.76, 12.48, 0.68)
callout(s, M, 5.40, 12.30,
        "**The ringed row is the port on its own.** With the update timing held "
        "identical, the two finders agree on all **23 244 605** clusters. Not "
        "small: zero. The six at f32 are the next slide.",
        h=0.76, size=PT_LEAD, color=RED)
caption(s, M, 6.38, 12.30,
        "3×3 · 10 000 frames · 23.2 M clusters · exact centre-set difference at "
        "tol = 0. [f64 ped] and [f32 ped] differ only in DEVICE_PED_TYPE. What is "
        "left once the timing is allowed to differ is slide 29.")

notes(s, """SAY THIS -- 121 words, about 44 s.

- Now we can change one thing at a time, and this is the headline of the whole
  validation section.

- The ringed row is the port on its own -- the CPU twin against CUDA, with the
  update timing held identical. Zero differences. Not "small": zero, out of 23
  million clusters. The port is exact.

- The second row, six clusters, is the float pedestal, and that is the next
  slide. Six out of 23 million, and I can show you all six.

- What happens when the timing is allowed to differ is slide 29, and it is a
  CPU-against-CPU story.

DETAILS: WHY THE SERIAL ROWS ARE NOT ON THIS SLIDE. They carry the same 8/11 as
each other, which is a separate claim of its own -- everything CUDA changes is
worth zero -- and reading it off the top of a four-row table made it look like
the weaker half. It has slide 29 to itself.

WHICH TEST the timing moves is measured on slides 30-31, and the answer is
specific:

    frozen-only 11 clusters  ->  ALL from Test3 (total > c3*nSigma*rms)
    cpu-only     8 clusters  ->  ALL from the local-max gate (value == max)
    Test1 threshold          ->  contributes ZERO clusters

Why Test3 is the exposed one: it sums all nine values, so it collects the shift
of EVERY already-scanned neighbour -- three above, one left, exactly the pixels
eligible for update -- where Test1 feels at most the one that happens to be the
argmax.

Confirmed two ways: instrumented (branch codes diffed per frame) and ablated
(Test3 compiled out, at which point the 11 go to zero exactly). Slides 30-31
have both.

Test1 flips MORE often than Test3 -- 619 of 974 divergent pixels -- and costs
nothing, because it only moves a pixel between QUIET_UPDATE and SHADOW and
neither of those stores. It is not merely downstream of Test3 either: with Test3
ablated the first divergence is still a Test1 flip, on a frame where the two
pedestals were provably identical. Test1 initiates on its own, just ~7x slower
(frame 30 vs frame 4), and can never create a cluster by itself.

The f32 row is a different mechanism entirely -- EMA drift, not timing, since
both finders freeze the pedestal there. That is the next slide.""")

# ================================================ 28 · THE MISMATCH, SEEN
s = new_slide()
chrome(s, 27, "Validation · the disagreement, seen",
       "The whole disagreement is one duplicate centre")
figure(s, "fig_mismatch147", 1.19, 1.72, 10.95)
callout(s, M, 6.34, 6.04,
        "Only each finder's **own** 3×3 footprints are drawn, so the panels differ "
        "**exactly** where the finders do.", h=0.78, size=PT_LEAD)
callout(s, 6.75, 6.34, 6.04,
        "cuda's patch is **one row taller**: a second centre directly below the one "
        "both found. **The charge is already counted**: a duplicate, not a new "
        "photon.", h=0.78, size=PT_LEAD, color=AMBER)
notes(s, """SAY THIS -- 116 words, about 42 s.

- This is one of those six differences, drawn. Left is the CPU twin, right is
  CUDA, same frame, same region.

- Each picture only shows the pixels that THAT finder claimed. So where they
  look the same, the two agreed pixel by pixel. And they mostly do.

- The one difference is ringed: CUDA keeps a second centre directly below one
  that both of them keep. Its patch is one row taller.

- And this is the important part for a physicist: nothing was missed, and
  nothing was invented out of noise. The charge is already counted in the
  cluster next to it. The two finders are arguing about which of two touching
  pixels owns the same photon.

DETAILS: Frame 147, the strongest of the six residuals, shipping f32 build.
Every other centre in the patch agrees, including the ordinary photon at bottom
right.

Each panel masks every pixel that is not inside the 3x3 footprint of one of THAT
finder's own cluster centres. So a pixel is visible on the left only if frozen
claimed it, and visible on the right only if cuda did. Where the panels look
identical, the finders agreed pixel for pixel.

The red dots are cluster centres. The amber ring is the one centre cuda keeps
that frozen does not: (202, 8), directly below (202, 7) which both finders keep.
Because the two 3x3 windows overlap, cuda's patch is four rows tall where
frozen's is three - that visible extra row IS the disagreement.

Note what is NOT happening: there is no cluster in one panel that is absent from
the other's neighbourhood entirely. Nothing was missed, and nothing was invented
out of noise. The two finders are arguing about which of two adjacent pixels owns
a photon they both found.

Same construction as helper.plot_masked_mismatch() in the notebook, so this is
reproducible from ClusterFinderFrozen_vs_CUDA.ipynb directly.""")

# ========================================================== 27 · RESIDUALS
s = new_slide()
chrome(s, 28, "Validation · the six residuals, dissected",
       "float32 cannot tell these two pixels apart")
# The walkthrough IS the slide: the whole argument is that two printed numbers
# are the same number in float32, and that only lands if the room can read both.
code(s, M, 1.82, 6.56, [
    "frame 147    centre (x=202, y=8)      3×3 window",
    "raw window (ADU)     pedestal-subtracted, 1 decimal",
    "[[4646 5282 4703]    [[ 45.3 «638.4» -12.1]   frozen and cuda",
    " [4857 5318 4950]     [ 43.7 «638.4»  -7.5]   print the SAME",
    " [4763 4640 4858]]    [  1.1  70.7  136.4]]   window",
    "",
    "the two contenders, at full precision:",
    "                        rival (dy=-1)     centre",
    # Blue on BOTH cuda values, because the claim is that they are one number --
    # the same blue as the two «638.4» cells in the window above. Red on the one
    # frozen value that is bigger, because that single number is the whole reason
    # f64 rejects where f32 accepts. Marking frozen's centre too would say "look
    # at these" when the argument is "look at THIS one".
    "frozen  [f64 ped]     ⟪638.383019956⟫   638.382773664",
    "cuda    [f32 ped]     «638.382812500»   «638.382812500»",
    "",
    "gate:  accept if  centre >= max(window)",
    "   frozen   638.382773664 >= 638.383019956  ->  reject",
    "   cuda     638.382812500 >= 638.382812500  ->  «accept»",
], size=PT_CODE)
callout(s, M, 4.86, 6.56,
        # "**ULP***" -> bold ULP followed by a plain asterisk: split("**") puts
        # the stray * at the head of the following non-bold run.
        "Under the f64 pedestal the two pixels are **0.000246 ADU** apart; one "
        "float32 **ULP*** here is **0.000488**. **Half a ULP**: in float32, the "
        "same number.", h=1.0, size=PT_LEAD)
figure(s, "fig_spectra_valid", 7.32, 1.82, 5.58)
callout(s, M, 6.10, 12.30,
        "**At f64 the residual is zero.**", h=0.52, size=PT_LEAD, color=AMBER)
# The footnote has to say WHICH ulp, or a careful listener computes 2^-14 for a
# number near 638 and concludes the slide is wrong by a factor of 8. The grid is
# set by the ~5 000 ADU operands of the subtraction: 638.3828125 is exactly
# 1 307 408 x 2^-11, which is the tell.
caption(s, M, 6.76, 12.30,
        "*  ULP = unit in the last place, the gap between one float32 and the "
        "next. Here it is 2⁻¹¹ = 0.000488 ADU, set by the ~5 000 ADU operands of "
        "the subtraction, not by the 638 ADU answer, whose own ULP is 8× finer.")
notes(s, """SAY THIS -- 118 words, about 43 s.

- And here is why, in numbers you can read. This is the test that decides
  whether a pixel is the local maximum: is my value greater than my neighbour's?

- Look at the two sums printed there. They differ in the eleventh digit. In
  float, with 32 bits, those two numbers are THE SAME number -- there is no bit
  left to hold the difference.

- So the comparison goes one way on the CPU, which uses doubles, and the other
  way on the GPU, which uses floats. Both answers are defensible. Neither is a
  bug.

- Six clusters out of 23 million, all of this shape. If you need exactly the CPU
  answer, use the double pedestal build; it is exact.

DETAILS: Frame 147, centre (x=202, y=8). The full numbers.

                              centre               rival      centre - rival
  frozen (f64 host ped)  638.382773664       638.383019956        -0.000246291
  cuda   (f32 dev  ped)  638.382812500       638.382812500         0.000000000

  pedestal mean, centre:  f64 4679.617226336   f32 4679.617187500
  pedestal mean, rival :  f64 4643.616980044   f32 4643.617187500

  1 float32 ULP at 4679.6 ADU = 0.000488 ADU
  the f64 separation between the two pixels = 0.50 ULP

Read that last line slowly: the two pixels differ by HALF a float32 ULP. There is
no float32 number between them, so once the pedestal is stored as float32 they
round to the identical bit pattern. The local-max gate is `value >= max`, which
accepts on equality - so CUDA keeps the centre and frozen, which can still see
the rival is 2.5e-4 ADU higher, rejects it.

This is not drift in the usual sense. It is not that the f32 EMA wandered away
from the f64 one over thousands of updates - look at the pedestal columns, they
agree to 4e-5 and 2e-4 ADU. It is that the QUESTION being asked ("which of these
two pixels is larger?") has an answer that float32 cannot represent.

Consequences worth stating out loud if asked:
  - All six exist only in the shipping f32 build.
  - Each of the six sits ONE PIXEL from a cluster both finders found: a duplicate
    neighbour, never a spurious photon and never a missed one.
  - It is one-directional. cuda-only 6, frozen-only 0. A tie can only ever ADD a
    centre, never remove one, because >= accepts.
  - All six are at Chebyshev distance 1 from a cluster both finders found: a
    duplicate neighbour, never a spurious photon.
  - With DEVICE_PED_TYPE = double the residual is 0 / 0 at both 3x3 and 9x9.
  - Fixing it in f32 would mean a strict > in the gate, which changes the CPU's
    documented behaviour on genuine ties. Not worth 6 clusters in 23 million.""")

# ============================================ 29 · THE CPU-vs-CPU RESIDUAL
# The top half of what used to be one four-row table. Splitting it is the point:
# those two rows carry the SAME 8/11, which is the strongest statement in the
# deck about the port -- everything CUDA changes is worth exactly zero -- and it
# was being read as the weak half of a table whose ringed pair said "0 / 0".
s = new_slide()
# THE TITLE NAMES THE CAUSE, NOT THE VENUE. It used to read "the gap is CPU
# against CPU", which is where the difference can be reproduced, not what makes
# it -- and it read as a contradiction over a table whose second row is serial
# against CUDA. Slide 26 established frozen == CUDA exactly at f64, so the only
# variable left between serial and CUDA is the one frozen was built to isolate:
# when the pedestal is updated. That is a choice in the algorithm, not a defect
# in the port, and it is what the slide should say.
chrome(s, 29, "Validation · the residual that is left",
       "All 19 clusters come from when the pedestal is updated")
bullets(s, M, 1.82, 12.40, [
    "The serial finder **updates the pedestal as it scans**, so a pixel can be "
    "tested against a pedestal its own neighbours already moved. Frozen and "
    "CUDA both **hold the pedestal still for the whole frame**.",
    "That is the only thing separating these two rows from the exact pair. And "
    "the rows carry **the same 8 and 11 clusters**, so **nothing CUDA does "
    "contributes anything at all**.",
], size=PT_BODY)
table(s, M, 3.42, 12.40,
      ["comparison", "the one thing that differs", "A-only / B-only",
       "% of clusters"],
      [["serial CPU  vs  frozen CPU", "update timing alone: **a CPU-only effect**",
        "**8 / 11**", "0.000082 %"],
       ["serial CPU  vs  CUDA  [f64 ped]", "the same thing, **and nothing else**",
        "**8 / 11**", "0.000082 %"]],
      colw=[0.28, 0.40, 0.17, 0.15], size=PT_TABLE, rowh=0.60)
callout(s, M, 5.40, 12.30,
        "**19 in 23 244 605, and the port is in none of them.** frozen and CUDA "
        "agree exactly at f64, so the one variable left between serial and CUDA "
        "is **when the pedestal moves**. Which test it moves is the next two "
        "slides.", h=0.76, size=PT_LEAD, color=RED)
caption(s, M, 6.38, 12.30,
        "3×3 · 10 000 frames · 23.2 M clusters · exact centre-set difference at "
        "tol = 0. A-only means found by the left finder and not by the right.")
notes(s, """SAY THIS -- 112 words, about 41 s.

- So what IS left? Nineteen clusters, and they come from the one thing the CPU
  twin was built to hold fixed: when the pedestal is updated.

- The serial finder updates as it scans. So by the time it reaches a pixel, some
  of that pixel's neighbours have already pushed their values into the pedestal.
  Frozen, and CUDA, keep the pedestal still for the whole frame.

- Now read the two rows. Serial against frozen: 8 and 11. Serial against CUDA:
  the same 8 and 11. Identical.

- And remember the previous slide: frozen and CUDA agree EXACTLY at f64. So the
  only variable left between the serial finder and CUDA is the one frozen was
  built to isolate -- when the pedestal is updated. All nineteen clusters come
  from that, and none of them from the port.

- The next two slides say which test does it.

DETAILS: the decomposition, in case it is asked before the next slide gets there.

    frozen-only 11 clusters  ->  ALL from Test3 (total > c3*nSigma*rms)
    cpu-only     8 clusters  ->  ALL from the local-max gate (value == max)
    Test1 threshold          ->  contributes ZERO clusters

Why Test3 is the exposed one: it sums all nine values, so it collects the shift
of EVERY already-scanned neighbour -- three above, one left, exactly the pixels
eligible for update -- where Test1 feels at most the one that happens to be the
argmax.

WHY THIS IS NOT A DEFECT IN EITHER FINDER. Neither update order is more correct;
they are two defensible readings of "the pedestal at this pixel". The claim the
deck makes is narrower and provable: the CUDA port reproduces the frozen model
exactly, and the frozen model differs from the serial one by 19 clusters in 23
million, for a reason that has nothing to do with the GPU.""")


# ================================================ 30 · THE THREE TESTS
# Promoted out of the annex. Once the arc says out loud that the residual is
# a CPU-vs-CPU effect (slide 29), the next question is WHICH test moves, and
# the answer is measured rather than argued -- so it belongs in the arc.
# Two slides. Slide 4's code panel is titled "THE WHOLE ALGORITHM" but shows a
# two-test simplification -- Test3 is not in it. So part 1 has to put the third
# branch back before part 2 can blame it for anything.
s = new_slide()
chrome(s, 30, "Validation · the third test", "There is a third test")
bullets(s, M, 1.80, 12.30, [
    "Slide 4 showed the decision with **two** tests, which is the right "
    "simplification for the arc. The finder has **three**, and the missing one "
    "is what the next slide is about.",
], size=PT_BODY)
code(s, M, 2.48, 7.54, [
    "v = frame[i] - ped_mean[i];   rms = ped_rms[i]",
    "if (v < -nSigma*rms)           -> skip, no update",
    "m     = «max» over the 3x3 window",
    "total = «sum» over the 3x3 window",
    "if (m > «nSigma»*rms)              // TEST 1",
    "    if (v == m) -> emit cluster  //   local-max gate",
    "    else        -> shadow, no update",
    "else if (total > «c3*nSigma»*rms)   // TEST 3",
    "    if (v == m) -> emit cluster",
    "else            -> update pedestal",
# 146-186, not 104-142: find_clusters() opens at 123 and the branch chain runs
# 149-186. A listener who opens the file at 104 lands in steal_clusters().
], size=PT_CODE, title="ClusterFinder.hpp:146-186 · all three branches")
rail(s, [
    ("label", "what each test asks"),
    ("gap", 0.12),
    ("row", "Test 1  ·  nSigma · rms", "is any ONE pixel bright?", ACCENT),
    ("gap", 0.12),
    ("row", "Test 3  ·  c3 · nSigma · rms", "is the WINDOW bright?", AMBER),
    ("gap", 0.12),
    ("row", "local-max gate", "is this pixel the peak?", PALE),
    ("gap", 0.18),
    ("note", "Test 1 reads the window's MAX; Test 3 reads its SUM. That one "
             "difference is the next slide."),
], y0=2.15)
callout(s, M, 5.86, 12.30,
        "**c3 = √(3×3) = 3 falls out of variance addition.** Variances add, "
        "standard deviations do not: nine pixels of noise σ give **Var(Σ v) = "
        "9σ²**, so the sum's noise is **3σ**. **c3·nSigma·rms** is therefore the "
        "**same 5σ criterion as Test 1**, asked of the window instead of the "
        "pixel.", h=1.1, size=PT_LEAD)
caption(s, M, 7.02, 12.30,
        "Roughly 80 % of pixels reach the last branch and push the pedestal; "
        "~1.5 % are peaks and ~18 % sit in a peak's shadow.", size=PT_META)
notes(s, """SAY THIS -- 84 words, about 31 s.

- Earlier I showed the decision with two tests. The finder actually has three,
  and I left one out to keep that slide readable.

- The third one looks at the total of the whole cluster, not just the brightest
  pixel or the brightest four. So it can accept a photon that is spread out and
  has no single strong pixel.

- At the small cluster size it is a small correction. It matters here because it
  is the test behind the CPU-versus-CPU difference.

DETAILS: Why slide 4 leaves Test3 out: at 3x3 it is a small correction to the
cluster count, and the arc's point there is the THREE OUTCOMES shape -- store,
shadow, update -- not completeness. This slide is where completeness belongs.

Read c3 out loud, it is the part people find satisfying. Test1 asks whether one
pixel is 5 sigma above ITS OWN noise. Test3 asks the same question of the SUM of
nine pixels -- so the only thing needed is the noise on that sum, and that is
just variance addition:

    Var(v1 + ... + v9) = Var(v1) + ... + Var(v9) = 9 * sigma^2     (independent)
    sigma_sum = sqrt(9 * sigma^2) = 3 * sigma

Standard deviations do NOT add; variances do. That is the whole content of the
sqrt: nine pixels give nine times the variance but only three times the rms, so
a threshold on the sum has to be 3x larger to carry the same 5-sigma meaning.
It is also why summing helps at all -- the signal adds linearly (x9) while the
noise adds in quadrature (x3), a net sqrt(9) = 3x gain in signal-to-noise for a
photon that is genuinely spread over the window.

c3 = sqrt(ClusterSizeX * ClusterSizeY) in the constructor, so it generalises: at
9x9 it is sqrt(81) = 9.

If someone challenges the independence assumption: it is the pedestal noise that
is being added, which is per-pixel readout noise and is uncorrelated between
pixels to a good approximation. Correlated noise would need covariance terms and
would make c3 larger than sqrt(N).

The negative-value skip at the top is a fourth branch but not a test in the same
sense -- it drops pixels far BELOW pedestal, which are detector artefacts rather
than photons, and it does not update either.

What matters for the next slide: Test1 reads the MAX of the window, Test3 reads
the SUM. That is the whole difference in sensitivity.""")


# ============================================ 31 · WHICH TEST, MEASURED
s = new_slide()
# 25, not the 27 the arc otherwise uses: at 27 this title takes two lines and
# pushes the first bullet into the figure. 25 keeps it to one line with 8 % of
# margin.
chrome(s, 31, "Validation · which test causes the gap",
       "Test3 makes the clusters; Test1 only moves the pedestal",
       title_size=25)
bullets(s, M, 1.76, 12.30, [
    "**push_fast** touches only the pixel's **own** accumulators and never reads "
    "the stencil, so the update is **order-independent**: same set of updated "
    "pixels, bit-identical pedestal. Only a differing **decision** can diverge.",
], size=PT_BODY)
figure(s, "fig_test3", M + 0.72, 2.49, 10.80)
callout(s, M, 6.13, 12.30,
        "Measured two ways. **Instrumented**: of the 19 disagreeing clusters, the "
        "**11 that only frozen finds are all Test 3**; the 8 only serial finds are "
        "all the local-max gate. **Ablated**: compile Test 3 out and the 11 go to "
        "**zero**.", h=0.72, size=PT_LEAD)
caption(s, M, 7.05, 12.30,
        "3×3 · 10 000 frames · 23 244 605 clusters · per-pixel branch codes "
        "diffed frame by frame, then re-run with Test 3 compiled out.", size=PT_META)
notes(s, """SAY THIS -- 92 words, about 34 s.

- This is the proof that the CPU-versus-CPU difference is the third test, and it
  was done two independent ways.

- First, instrumented: record which branch every pixel takes, in both finders,
  and compare. All eleven clusters that only the frozen twin finds come from
  that one test.

- Second, ablated: compile the test out of BOTH finders and re-run. The eleven
  go to exactly zero.

- The picture shows one such pixel. The two sums differ in the eleventh digit,
  and the maximum is identical in both.

DETAILS: TWO EXPERIMENTS, INDEPENDENT ROUTES.

Instrumented (both finders shipped logic, branch codes recorded): 974 pixels out
of 1.6e9 take a different branch. Only some change the cluster set, and those
decompose onto slide 29's 8/11 with no remainder:

    frozen-only 11  =  QUIET_UPDATE -> TEST3_STORE      (all Test 3)
    cpu-only     8  =  TEST3_STORE  -> TEST3_SKIP  (5)  (local-max gate)
                    +  TEST1_STORE  -> SHADOW      (3)

Ablated (Test 3 compiled out of BOTH finders):

                        Test3 ON     Test3 OFF
    clusters (cpu)    23 244 602    23 241 342
    divergent pixels         974           584
    first divergence   frame   4     frame  30
    frozen-only               11             0
    cpu-only                   8             4

The 11 go to zero exactly. The cpu-only 4 are all TEST1_STORE -> SHADOW, the
local-max gate, which does not need Test 3 -- the count is not preserved because
ablating changes which pixels push, so the pedestal follows a different path and
the downstream ties are a different realisation.

WHY TEST1 FLIPS MOST AND COSTS NOTHING. A Test1 flip moves a pixel between
QUIET_UPDATE and SHADOW. Neither stores, so no cluster appears or disappears; it
only changes whether that pixel pushed, which feeds forward.

TEST1 IS NOT MERELY DOWNSTREAM. With Test3 off, the first divergence is still at
frame 30, and nothing diverged before it -- so the pedestals were identical and
that Test1 flip came straight from the within-frame asymmetry. Test1 initiates
independently; it is just ~7x slower to do so (frame 30 vs frame 4) and cannot
create a cluster on its own.

THE SITE. Frame 203, pixel (125,245). Threshold c3*5*rms = 256.511. Serial sums
256.489 and calls it a pedestal sample; frozen sums 256.581 and stores. The gap
is 0.09 ADU, from four neighbours whose pedestals differ by 0.12 ADU in total.
max is 56.473 in BOTH -- printed because it is the argument: Test1 could not
have caused this.

READING THE LEFT PANEL. It is the FINISHED frame, not a scan in progress -- the
branch map is read once find_clusters() returns, so every one of the 441 cells
carries a final decision. The arrows are raster ORDER, drawn outside the grid so
they cannot be mistaken for a cursor.

If asked why a photon looks 3x4 and not 3x3: it does not. The green cells are
the stored clusters and they are single pixels, nine of them in this patch. The
dark region around them is SHADOW, and shadow is not the photon's 3x3 -- it is
every pixel whose OWN 3x3 window contains something above 5 sigma. Note the
shadow pixel need not be bright itself. Three tiers, all visible here:

    810.3 ADU   IS the window max          -> stored
    345.0 ADU   above the 85.5 bar, not the peak  -> shadow
     17.7 ADU   nowhere near the bar, but its 3x3 reaches the 345 -> shadow

Charge shared across two adjacent pixels therefore lights up the union of every
window that can see either one, which here is 3x4. In the whole patch 27 pixels
clear 5 sigma, 9 are local maxima, and 98 end up shadowed.

The marked pixel is amber with a green ring on purpose: the panel is coloured by
the SERIAL finder, and serial sampled the pedestal there. Only frozen stored. If
it were painted plain green the picture would contradict the sigma lines beside
it. Amber means one thing in BOTH panels -- a pixel that pushed the pedestal --
which is why the frozen readings on the right are red rather than amber.

WORTH POINTING AT IF THE ROOM IS ENGINEERS, four rows below the marked pixel.
(129,245) is an isolated shadow cell with pedestal samples on both sides of it.
The cause is the 85.150 ADU pixel at (130,244), one row down and one column
left, which sits in the window of BOTH (129,245) and (129,244). Their window
maxima are IDENTICAL. They take different branches anyway, because the bar is
per-pixel -- the test is m > nSigma * rms[centre], the noise of the pixel being
tested. The branch codes alone bound the two:

    rms(129,245)  <  17.030  <=  rms(129,244)

The left neighbour is noisier, so its bar is higher, so the same 85.150 fails to
clear it. And (130,244) is itself only a pedestal sample: 85.2 does not clear its
OWN bar and its window sums to 119.5 against a ~256 Test3 threshold. So a pixel
can silence a neighbour's pedestal update without being bright enough to be
anything itself -- "bright" is not a property of a pixel, it is a relation
between a value and whose noise you measure it against.

Deliberately NOT marked on the figure. It is a statement about the algorithm,
not about serial-vs-frozen, and a second leader line would compete with the one
that carries the slide's actual argument. Full numbers in docs/deck/QA.md.

The branch_map property on both finders is diagnostic scaffolding. The ablation
is AARE_TEST3_ENABLED in the two headers, 1 by default.""")


# =========================================================== 28 · API 1
s = new_slide()
chrome(s, 32, "For users · Python API", "The fast path in eight lines")
code(s, M, 1.82, 7.85, [
    "from aare import File, ClusterFinderCUDA",
    "cf = ClusterFinderCUDA(image_size=(400, 400), cluster_size=(3, 3),",
    "                       n_sigma=5, «n_streams»=4,",
    "                       «max_clusters_per_frame»=3000)",
    "for _ in range(1000):                    «# 1. train the pedestal»",
    "    cf.push_pedestal_frame(pd.read_frame())",
    "data = f.read_n(100_000)                 «# 2. one contiguous array»",
    "cf.‹register_input_buffer›(data)           «# 3. pin it once»",
    "for s in range(0, len(data), 2000):      «# 4. batch through it»",
    "    clusters = cf.‹find_clusters_batched›(data[s:s+2000], first_frame=s)",
    "cf.‹unregister_input_buffer›()             «# 5. release the pages»",
], size=PT_API, title="THE RECOMMENDED PATTERN")
# The zero-copy variant sits directly under the pattern it replaces, so the two
# can be read against each other instead of across the slide.
code(s, M, 5.42, 7.85, [
    "for v in cf.‹find_cluster_views_batched_iter›(data, 2000):  «# 4, zero-copy»",
    "    hist.fill(v.sums())        # consume inside the loop:",
    "                               # the view dies with its chunk",
], size=PT_API, title="ZERO-COPY · SAME PATTERN, STEP 4 SWAPPED")
# Right column: one entry per numbered step, so the mapping needs no arrows.
bullets(s, 8.66, 1.82, 4.24, [
    "**3 · register_input_buffer()**: opt4 in one call. Pin **once**, outside "
    "the loop.",
    "**4 · find_clusters_batched()**: one ClusterVector per frame, with opt5's "
    "overlap built in.",
], size=PT_BODY)
callout(s, 8.66, 3.72, 4.24,
        "**Zero-copy: ×1.16 at 3×3, ×2.21 at 9×9.** The views expose every "
        "cluster; they only withhold **ownership** past the chunk.",
        h=1.15, size=PT_LEAD)
callout(s, 8.66, 5.42, 4.24,
        "Reduce as you go and you never allocate. Hold a view past its chunk and "
        "you stall the pipeline.", h=1.15, size=PT_LEAD, color=AMBER)
notes(s, """SAY THIS -- 117 words, about 42 s.

- If you take one slide away for your own code, take this one. This is the fast
  path, and it is eight lines.

- Two of them matter. Register the input buffer once, outside the loop -- that
  is the pinning step, and you do it on the whole dataset, not per frame. And
  call find_clusters_batched, which has the host-GPU overlap built in. You do
  not have to think about streams at all.

- The lower panel is the zero-copy version: same code, one call swapped. You get
  a view instead of your own copy. Read it inside the loop, and reduce as you
  go. If you keep it past the chunk, it is gone.""")

# =========================================================== 27 · API 2
s = new_slide()
chrome(s, 33, "For users · choosing the knobs",
       "Five knobs, and the one that silently truncates")
hdr = [("Parameter", 1.05), ("What it does", 3.72), ("Guidance", 5.37)]
# The ring and the warning are the same point, so they carry the same colour.
callout(s, M, 1.78, 12.30,
        "The single most common mistake: leaving **max_clusters_per_frame** too low. "
        "It does not error; it truncates, and every frame quietly returns the same count.",
        h=0.62, size=PT_LEAD, color=RED)
y = 2.56
rect(s, M, y, 12.30, 0.4, PANEL)
for lab, dx in hdr:
    tf = tb(s, M + dx - 0.85 if dx > 1.05 else M + 0.28, y + 0.09, 5.17, 0.3)
    run(para(tf, True), lab.upper(), PT_LABEL, MUTED, bold=True, spc=1.3)
y += 0.44
# Guidance is what a user will act on, so it says the number and the consequence
# and stops. Everything that explains WHY is in the notes.
params = [
    ("n_streams", "Upper bound on frames in flight.",
     "4 at both cluster sizes. 8 buys no kernel concurrency at 9×9."),
    ("max_clusters_per_frame", "Fixed size of the per-frame D2H transfer.",
     "Must exceed the real maximum or clusters are silently dropped. At 9×9 the "
     "maximum is 1 633, and a cap of 1 700 already makes D2H the bottleneck."),
    ("batch size", "Frames per find_clusters_batched call.",
     "2 000 amortises launch overhead without a large pinned footprint."),
    ("cluster_size", "Compile-time stencil geometry.",
     "3×3 and 9×9 are registered. 9×9 moves the bottleneck off H2D."),
    ("register_input_buffer", "Page-locks the host array for DMA.",
     "Always, if the data is already in RAM."),
]
flagged = None
for i, (p_, what, guide) in enumerate(params):
    if i % 2 == 0:
        rect(s, M, y, 12.30, 0.82, PANEL)
    tf = tb(s, M + 0.29, y + 0.16, 2.69, 0.5)
    run(para(tf, True, line=1.1), p_, 11, ACCENT, font=MONO, bold=True)
    tf = tb(s, M + 3.10, y + 0.16, 3.00, 0.6)
    run(para(tf, True, line=1.2), what, 11, PALE)
    tf = tb(s, M + 6.36, y + 0.16, 5.58, 0.6)
    run(para(tf, True, line=1.2), guide, 11, TEXT2)
    if p_ == "max_clusters_per_frame":
        flagged = y
    y += 0.86
# The one row that loses data if you get it wrong, ringed so it is found without
# reading the table. Drawn last so the outline sits over the zebra fill.
frame_rect(s, M - 0.04, flagged - 0.05, 12.38, 0.96)
notes(s, """SAY THIS -- 113 words, about 41 s.

- Five settings, and only one of them can hurt you. That is the ringed one, the
  maximum clusters per frame.

- It is the size of the envelope we copy back for every frame. If you set it too
  low, you do not get an error. You quietly lose clusters, and every frame comes
  back with the same count -- which looks like a plateau in your data, not like a
  bug.

- The others are easy. Four streams is right at both cluster sizes; eight buys
  nothing. Batches of two thousand. Register the input buffer whenever your data
  is already in memory. Those are the defaults, and they are the ones we
  measured.""")

# =========================================================== 28 · NEXT
s = new_slide()
chrome(s, 34, "Where this leaves us",
       "The bottleneck has walked from the host, to the GPU, to PCIe")
# Four conclusions, not four paragraphs: each card states the fact and the
# number that carries it. The reasoning behind each is in the notes.
cards = [
    ("DONE", ACCENT, "×9.1 at 3×3, ×26.5 at 9×9",
     "16.3 and 25.1 µs/frame end to end, both on their hardware floor. At 3×3 "
     "that is 47× MÖNCH03's standard frame rate."),
    ("DONE", ACCENT, "FP32 pedestal, safely",
     "−40.5 % kernel time, and correct: the variance accumulates on a frozen "
     "per-pixel offset. 6 duplicate clusters in 23 million."),
    ("NEXT", AMBER, "3×3: transfer granularity",
     "The pipeline sustains 16.31 µs against a 13.15 µs uncontended H2D: 2 000 "
     "separate 320 kB descriptors, plus H2D↔D2H contention."),
    ("NEXT", AMBER, "9×9: the result path, not the kernel",
     "At a lossless cap D2H already binds: 25.24 µs against a 23.94 µs kernel. "
     "More kernel work buys nothing."),
]
callout(s, M, 1.78, 12.30,
        "**Every number in this deck is measured on one machine, one dataset, one "
        "git revision**, and the annexes carry the reconciliation for each of them.",
        h=0.62, size=PT_LEAD)
for i, (tag, col, title, body) in enumerate(cards):
    cx = M + (i % 2) * 6.25
    cy = 2.62 + (i // 2) * 2.30
    rect(s, cx, cy, 6.04, 2.00, PANEL)
    rect(s, cx, cy, 6.04, 0.035, col)
    tf = tb(s, cx + 0.3, cy + 0.26, 1.45, 0.26)
    run(para(tf, True), tag, 9.5, col, bold=True, spc=1.5)
    tf = tb(s, cx + 0.3, cy + 0.58, 5.37, 0.4)
    run(para(tf, True, line=1.1), title, 15, PALE, bold=True)
    tf = tb(s, cx + 0.3, cy + 1.08, 5.37, 0.85)
    run(para(tf, True, line=1.3), body, 12, TEXT2)
notes(s, """SAY THIS -- 126 words, about 46 s.

- Where this leaves us. Nine times at 3x3, twenty-six at 9x9, and both are
  sitting on their hardware floor. At the small size that is 45 times the
  detector's standard frame rate, so the finder is no longer the limit.

- And the pedestal is now float, safely -- 40 percent off the kernel, with six
  duplicate clusters in 23 million, all of them the same tie between two
  touching pixels.

- What is left is not the maths. At 3x3 it is how we hand the wire its work --
  two thousand separate small transfers. At 9x9 it is the copy back out, which
  is now the tallest bar. Both of those are transfer problems. The bottleneck
  has walked from the host, to the GPU, to the wire.

DETAILS: The headroom claim, stated carefully.

WHERE THE MATERIAL IS, if anyone asks for it after the talk rather than from the
floor: docs/ClusterFinderCUDA_benchmark_results.md for every number and its
provenance; python/tests/ClusterFinderCUDA_perf.ipynb for the timing protocol;
python/tests/perf/ for the campaign; python/tests/validation_tiers.py and
ClusterFinderFrozen_vs_CUDA.ipynb for the validation; python/tests/branch_trace.py
and branch_site_dump.py for slides 30-31. Paths are deliberately off the slides.

MOENCH03 runs at 1.3 kHz as standard and 3-6 kHz with optimised readout boards.
The finder does 58 495 FPS at 3x3 and 39 775 at 9x9, so it is 45x and 31x the
standard rate, and about 10x and 7x the optimised ceiling. The point is not the
multiple: it is that cluster finding has stopped being the thing that decides
how fast you can take data, on one GPU, at either cluster size.

Two honest qualifications. First, this is throughput on frames already in host
RAM -- getting them there from the detector is a separate problem and is not
measured here. Second, at 9x9 the D2H slot is cap-sized, so the margin shrinks
if max_clusters_per_frame has to grow for a busier beam.

What the 24-thread CPU finder does, for contrast: 6 762 FPS at 3x3, which
matches the standard mode with nothing to spare, and 1 503 FPS at 9x9, which
does not.""")

# --------------------------------------------------------- divider · ANNEX
s = section("",
        "Annexes",
        "",
        [("A1", "Register pressure"),
         ("A2", "Rejected route · CUDA graphs"),
         ("A3", "opt3 · the second D2H"),
         ("A4", "opt5 · the overlap code"),
         ("A5", "Rejected routes · the result copy"),
         ("A6", "Why not one big pinned buffer?"),
         ("A7", "s1, s4, and every engine number"),
         ("A8", "opt7 · the full engine reading"),
         ("A9", "Catastrophic cancellation, in full"),
         ("A10", "The fault model, tested"),
         ("A11", "Behind the numbers · page faults"),
         ("A12", "The three benchmark artefacts")],
        rng=(1, N_ANNEX), col=AMBER, annex=True, item_pt=13, item_gap=0.115,
        carry=("Everything so far", "33 slides",
               "the arc is finished; what follows answers questions"))
notes(s, """SAY THIS -- 45 words, about 18 s. Only if you reach it.

- That is the talk. What follows is backup: every engine number reconciled, the
  three things I tried that did not work, and the measurement traps.

- I will jump to one of these if a question needs it.""")

# ===========================================================================
# ANNEX — the measurement detail behind slide 24, and the rejected routes
# ===========================================================================

# ---- A1 · REGISTER PRESSURE ----------------------------------------------
# Was main-arc slide 9. It is the CAUSE behind the occupancy percentages, and
# the arc can state the two register counts in one bullet on the occupancy
# slide; the cuobjdump evidence and the per-build caveat live here.
s = new_slide()
annex_chrome(s, 1, "what runs out first · expands slide 9",
             "Register pressure: 38 per thread at 3×3, 128 at 9×9")
bullets(s, M, 1.82, COL, [
    "An SM has **65 536 registers** and **1 536 thread slots**; whichever runs "
    "out first decides how many blocks fit at once.",
    "Every thread keeps a private **clusterData[CSX × CSY]**, so register "
    "demand grows with the **square** of the cluster size. Neither build "
    "spills.",
], size=PT_BODY)
figure(s, "fig_regpressure", M, 3.20, 8.16)
callout(s, M, 5.04, COL,
        "**9×9:** the register file is exactly full at two blocks "
        "(2 × 128 × 256 = 65 536).  **3×3:** the slots run out first, and the "
        "registers still have room.", h=0.80, size=PT_LEAD)
code(s, M, 6.20, COL, [
    "«cuobjdump -res-usage» build/aare/_aare_cuda*.so | c++filt",
    "  3x3: REG:«38»  STACK:0 LOCAL:0     # STACK/LOCAL 0 = no spills",
    "  9x9: REG:«128» STACK:0 LOCAL:0",
], size=PT_CODE, title="MEASURED, NOT ESTIMATED · READ FROM THE BUILT .SO")
rail(s, [
    ("label", "Per SM · sm_89 · f32 build"),
    ("gap", 0.15),
    ("stat", "3×3 · blocks resident", "6", ACCENT),
    ("stat", "9×9 · blocks resident", "2", AMBER),
    ("gap", 0.05),
    ("row", "Spills, either case", "0 bytes", TEXT2),
    ("row", "3×3 on the f64 build", "47 regs → 5 blocks", AMBER),
])
notes(s, """SAY THIS -- 117 words, about 45 s. Annex: only on a question.

- Shared memory is not the limit. Two other things are: thread slots, and
  registers. An SM has 1 536 thread slots and 65 536 registers, and whichever
  runs out first decides how many blocks live on that SM at once.

- Our block is 256 threads. At 3x3, 38 registers per thread: six blocks fit, and
  it is the SLOTS that run out first -- six times 256 is exactly 1 536. The
  registers still have room.

- At 9x9 every thread keeps a private 9-by-9 staging array, so 128 registers.
  Two blocks fill the register file exactly, and two thirds of the thread slots
  sit empty.

- Both numbers are read out of the built binary, and neither build spills.

CAREFUL: do not say "six versus two, because of register pressure". Registers
bind at 9x9 ONLY. At 3x3 the thread slots bind and the registers are at 89 %.
The callout on the slide says so, so the room can check you.

DETAILS: the arithmetic, and the one build-dependent caveat.

One SM holds 65 536 registers and 1 536 thread slots. A 16x16 block is 256
threads, so a block costs regs_per_thread x 256 registers. At 3x3 that is 9 728,
six blocks fit, and the THREAD SLOTS run out first. At 9x9 it is 32 768, so two
blocks exactly fill the REGISTER FILE and strand two thirds of the slots. Same
kernel, same block size, opposite binding resource.

Build dependence of the register count itself: the f64 pedestal costs 3x3 nine
extra registers (47), which loses a block per SM. 9x9 is unmoved at 128 -- there
the limiter is the clusterData[9][9] staging array, not the pedestal, which is
why opt7 helps 3x3's occupancy and not 9x9's.

Reproduce blocks/SM at runtime with python/tests/perf/kernel_resources.py, which
calls cudaOccupancyMaxActiveBlocksPerMultiprocessor on the built kernel.""")

# ---- A2 · ROUTE A · IDEA -------------------------------------------------
s = new_slide()
annex_chrome(s, 2, "rejected route · CUDA graphs · expands slide 10",
             "CUDA Graphs, a sound idea that the next act overtook", part=1, nparts=2)
bullets(s, M, 1.95, 12.40, [
    "Every cudaMemcpyAsync / kernel launch costs the **CPU** a few microseconds of "
    "driver work, per frame and per operation. After opt4 that looked like the budget.",
    # "captures" was the wrong verb twice over: it names the stream-capture API
    # we do not use, and it promised a one-call replay the code never delivers.
    "A **CUDA Graph** builds the whole dependency DAG once, so replay should be "
    "**one** cudaGraphLaunch. **Ours needs three per frame**, because every frame has "
    "new source and destination addresses.",
], size=PT_BODY)
# 3.42, not 3.15: two four-line bullets at 15 pt end at 3.32 once the renderer's
# own line height is counted, and the figure is an opaque PNG -- it does not
# collide with the text, it covers it, which is why the layout audit could not
# see this one.
figure(s, "fig_graphs", M, 3.42, 7.85)
# NOT STREAM CAPTURE. The earlier panel here showed cudaStreamBeginCapture /
# cudaStreamEndCapture; neither exists in the file. build_graphs() (graph.hpp
# 491-576) adds every node by hand. And the launch is NOT one call per batch:
# find_clusters_batched (graph.hpp 355-373) repoints the H2D source and the D2H
# destination per frame before launching, so it is three driver calls per frame,
# not one. That is the honest reason the route did not pay, and it was missing.
code(s, 8.56, 3.42, 4.24, [
    "// setup: nodes added by hand, not captured",
    "cudaGraphCreate(&sc.graph, 0);",
    "cudaGraphAddMemset/Memcpy/KernelNode(...);",
    "«cudaGraphInstantiate»(&sc.graphExec, ...);",
    "",
    "// per FRAME: 3 driver calls, not 1",
    "cudaGraphExecMemcpyNodeSetParams(...);  «x2»",
    "«cudaGraphLaunch»(sc.graphExec, sc.stream);",
], size=PT_CODE, title="ClusterFinderCUDA_graph.hpp")
# 5.62: the panel above is eight lines and ends at 5.48. At 5.06 the callout was
# starting inside it, on top of the cudaGraphLaunch line.
callout(s, 8.56, 5.62, 4.24,
        "**REJECTED**\n3×3: 39 752 FPS, inside noise of opt4.\n9×9: **11 072 FPS, 12 % slower**.",
        h=1.10, size=PT_LEAD, color=AMBER)
caption(s, M, 6.86, 12.40,
        "The 3×3 edge was never established: the instrumentation tax (2.8 µs) is "
        "larger than the gap (0.8 µs). Detail in the notes.")
notes(s, """SAY THIS -- 88 words, about 32 s. Annex: only on a question.

- Rejected route one: CUDA graphs. The idea is sound. Every copy and every
  launch costs the CPU a few microseconds of driver work. A graph is that whole
  sequence built once, and then replayed.

- After step four that looked like the right target. But it did not pay: even at
  3x3 the gain was inside the noise, and at 9x9 it was 12 percent SLOWER.

- And then the next act made the question go away: once host work hides under
  the GPU, saving host work is worth nothing.

CAREFUL: OUR GRAPH IS NOT "ONE CALL PER BATCH".

The textbook pitch is "four driver calls per frame become one". Our
implementation does not get there, and the code on the slide now says so:

  - It is built with the EXPLICIT node API, not stream capture --
    cudaGraphCreate + cudaGraphAddMemsetNode / AddMemcpyNode / AddKernelNode +
    cudaGraphInstantiate, all inside build_graphs(), graph.hpp 491-576.
  - Every frame reads from a different place in the input array and writes to a
    different output slot, so the H2D source pointer and the D2H destination
    pointer have to be repointed BEFORE each launch. That is two
    cudaGraphExecMemcpyNodeSetParams calls, then cudaGraphLaunch: three driver
    calls per frame (graph.hpp 355-373).

So the fair sentence is "four calls per frame became three, not one". If asked
why it could not be one: a graph is a fixed set of nodes with fixed parameters,
and a streaming workload gives every frame new addresses. Only a workload that
reuses the same buffers every time gets the textbook win.""")

# ---- A2 · ROUTE A · BUDGET (part 2) -------------------------------------
s = new_slide()
annex_chrome(s, 2, "rejected route · CUDA graphs",
             "What a CUDA Graph actually saves, in microseconds", part=2, nparts=2)
bullets(s, M, 1.90, COL, [
    "The stream path issues **four runtime calls per frame** (memset, H2D, "
    "launch, D2H) and a graph replaces all four with **one** cudaGraphLaunch. "
    "This is the **ceiling**: our own graph needs **three** (see part 1).",
], size=PT_BODY)
table(s, M, 2.80, COL,
      ["call", "per frame", "host cost", "µs/frame"],
      [["cudaMemcpyAsync", "2", "1.98 µs", "3.97"],
       ["cudaLaunchKernel", "1", "2.13 µs", "2.13"],
       ["cudaMemsetAsync", "1", "1.57 µs", "1.57"],
       ["**submission total**", "**4**", "", "**7.67**"]],
      colw=[0.40, 0.18, 0.22, 0.20], size=PT_TABLE, rowh=0.44)
code(s, M, 5.40, COL, [
    "7.67 us/frame x 3/4  = 5.75 us/frame   eliminated, AS MEASURED (under nsys)",
    "5.75 / 4 (see A12)   ~ 1.4 us/frame    eliminated, unprofiled estimate",
], size=PT_CODE, title="THE ARITHMETIC")
callout(s, M, 6.38, COL,
        "**~1.4 µs against a 16.17 µs floor = 8.7 %**, real while the host is the "
        "critical path, and **worth nothing after opt5**, which hides host work under "
        "the GPU entirely.", h=0.80, size=PT_LEAD, color=AMBER)
rail(s, [
    ("label", "route A · measured verdict"),
    ("gap", 0.10),
    ("stat", "3×3 vs opt4", "×1.03", TEXT2),
    ("stat", "9×9 vs opt4", "×0.88", AMBER),
    ("gap", 0.05),
    ("row", "vs opt5 at 9×9", "36 % behind", AMBER),
    ("gap", 0.20),
    ("note", "Read from CUPTI_..._RUNTIME, the host table, so this is the one "
             "order-of-magnitude number in the deck, ±50 %. It is enough to show "
             "graphs cannot pay against a 16 µs floor, and not enough to quote to "
             "two figures."),
])
notes(s, """SAY THIS -- 84 words, about 31 s. Annex: only on a question.

- This is the arithmetic behind that rejection, so it is not a matter of taste.

- Four driver calls per frame, about 7.7 microseconds of host time in total. A
  graph removes three of the four, so about 5.7 -- but that is measured under
  the profiler, which inflates host calls about four times. Unprofiled, call it
  1.4 microseconds.

- Against a 16 microsecond floor that is under nine percent, and only while the
  host is the critical path. After step five it is zero.

CAREFUL: THE 3/4 IN THIS TABLE IS THE BEST CASE, NOT WHAT WE BUILT.

"A graph removes three of the four calls" is the textbook ceiling. Our graph
removes ONE of the four: it still makes two cudaGraphExecMemcpyNodeSetParams
calls plus the launch, three per frame, because every frame needs new H2D and
D2H addresses (see part 1 and its note).

The 5.75 and 1.4 microsecond figures on this slide are therefore an UPPER BOUND
on what route A could ever have saved, not an estimate of what it did save. Say
"at best" when you quote them. The conclusion does not move -- it gets stronger,
since the real saving is smaller than the bound that was already too small to
matter -- but do not present the bound as the measurement. The 7.67 us and the
per-call host costs in the table ARE measured; only the 3/4 is modelled.""")


# ---- A3 · OPT3’S OTHER BARRIER -------------------------------------------
# opt3's title has always said "barriers", plural, but the deck only ever told
# one of them: the per-round drain loop. The count-then-fetch round
# trip went at the same step and was never shown, which made opt2 -> opt3 look
# like a refactor instead of the change of contract it was.
s = new_slide()
annex_chrome(s, 3, "opt3 · the other barrier · expands slide 14",
             "One D2H per frame, not two")
bullets(s, M, 1.82, 12.30, [
    "opt2 asked the device **how many clusters**, blocked until the answer came "
    "back, then asked for **that many**. The size of the second copy was a "
    "function of data that had not arrived yet.",
    "opt3 gives every frame a **fixed envelope** (count, then room for **cap** "
    "clusters), so the copy's size is known at construction and can be queued "
    "with the kernel. The count is still read, but **afterwards**, on the host.",
], size=PT_BODY)

flow(s, M, 3.20, 12.30,
     ["kernel", "copy 4 B", "BLOCK", "read count", "copy N B", "BLOCK"], h=0.62)
caption(s, M, 3.97, 12.30,
        "opt2 · two transfers and two stalls per frame: the host must learn the "
        "count before it knows how much to fetch.", size=PT_META)

flow(s, M, 4.31, 12.30,
     ["kernel", "copy the whole envelope", "→ next frame, host not involved"], h=0.62)
caption(s, M, 5.08, 12.30,
        "opt3 · one transfer, no stall. Nothing in the loop waits on a value.", size=PT_META)

callout(s, M, 5.50, 12.30,
        "**You can only stream a transfer whose size is known in advance.** The "
        "price is bytes: the envelope is sized by the **cap**, not by the clusters "
        "actually found.  **3×3: 120 kB shipped for 93 kB of clusters.  9×9: 558 "
        "for 467.**", h=0.86, size=PT_LEAD, color=RED)
notes(s, """SAY THIS -- 128 words, about 46 s. Annex: only on a question.

- There was a second barrier, and it is less obvious. In the old version the
  host asked the GPU "how many clusters did you find?", waited for the answer,
  and only then asked for that many bytes. The size of the second copy depended
  on data that had not arrived yet.

- You cannot queue up a copy whose size you do not know. So the host had to
  stop, twice, on every single frame.

- The fix is to give every frame a fixed envelope: room for a maximum number of
  clusters, decided once at the start. Now the size is known, so the copy can be
  queued right behind the kernel and nobody waits. It costs about twenty percent
  more bytes on an engine that had time to spare.

DETAILS: The point to say out loud: this is the one step in the ladder that is
not a setting. It changed the kernel signature, the buffer ownership and the
collection loop, and it is why the opt1/opt2 class is frozen in a separate header
(ClusterFinderCUDAOpt2.hpp) rather than being a flag on the current one.

The dependency edge is the whole argument. In opt2 the second memcpy's SIZE
argument is *sc.h_cluster_count -- host memory that only becomes valid after a
cudaStreamSynchronize (Opt2.hpp:277-294, then :442-446). So the sequence is
forced: kernel, copy 4 bytes, BLOCK, read, copy N bytes, BLOCK. Two of those six
steps are the host doing nothing, every frame.

opt3 fixes the size once in the constructor (m_output_bytes_per_frame =
m_clusters_offset + cap * sizeof(ClusterType), ClusterFinderCUDA.hpp:440-445), so
the copy is enqueued in the same loop iteration as the kernel launch (:725-728).
The kernel gets two pointers into ONE allocation (:701 and :709).

What it costs: 120 kB instead of 93 kB per frame at 3x3, 558 instead of 467 at
9x9. Roughly 20 % more bytes on an engine that had spare time, to buy back a
barrier that was stalling everything. That is also why the cap becomes a
throughput knob only from opt3 onward: under opt2 it bounded an allocation, under
opt3 it sets the D2H bar directly (report section 4.2).

If asked why opt2 did not simply copy a cap-sized buffer and skip the sync: that
IS opt3. It could not be done in opt2 because the clusters lived in their own
device allocation with no count field and no fixed per-frame stride -- there was
no single object to copy. Merging the two buffers is what created one.""")


# ---- A4 · THE OPT5 CODE --------------------------------------------------
s = new_slide()
annex_chrome(s, 4, "opt5 · the overlap code · expands slide 16",
             "The overlap, in six lines, and why you never write them")
code(s, M, 1.82, 7.39, [
    "tok = cf.«submit_batch»(data[a0:b0], first_frame=a0)",
    "for a, b in bounds[1:]:",
    "    nxt = cf.«submit_batch»(data[a:b], first_frame=a)   # GPU starts N+1 …",
    "    results.extend(cf.«collect»(tok))                   # … host unpacks N",
    "    tok = nxt",
    "results.extend(cf.«collect»(tok))                       # drain the last one",
], size=PT_CODE, title="THE WHOLE OF OPT5")
bullets(s, M, 3.64, 7.39, [
    "**You do not write this.** It is inside find_clusters_batched(), so opt5 "
    "arrived as a **speedup, not an API change**.",
    "submit_batch() and collect() stay public for streaming from a detector, or "
    "interleaving other work between chunks.",
], size=PT_BODY)
rail(s, [
    ("label", "chunk sizing · the two constraints"),
    ("gap", 0.12),
    ("row", "multiple of", "n_streams", ACCENT),
    ("row", "capped at", "MAX_SLOT_BYTES", ACCENT),
    ("gap", 0.16),
    ("note", "A multiple of n_streams, because the device pedestal is per-stream: "
             "an uneven chunk leaves the four pedestals at different ages and the "
             "finder stops being reproducible."),
    ("gap", 0.14),
    ("note", "Capped so the two pinned slots stay bounded: 544.5 kB per frame at "
             "9×9, cap 1 700. chunk_size_for(n) applies both rules."),
])
callout(s, M, 5.02, 7.39,
        "Two chunks in flight is enough. A third adds pinned memory and no overlap: "
        "the host is already busy for the whole time the GPU is.", h=0.72, size=PT_LEAD)
caption(s, M, 5.98, 7.39,
        "The saving is min(GPU, host) per chunk, so opt5 pays most when the two terms "
        "are comparable, ×1.31 at 3×3, and least when one dominates, ×1.20 at 9×9, "
        "where the host term is roughly twice the GPU term. That gap is the diagnosis "
        "that motivates opt6 (report §8.2).")
notes(s, """SAY THIS -- 86 words, about 31 s. Annex: only on a question.

- This is the whole of step five, six lines. Send the next chunk, unpack the
  previous one, repeat. Send, unpack, swap.

- The point is that you never write it. It is inside the normal batched call, so
  everyone got the speed-up without changing a line. The two halves stay public
  for people streaming from a detector.

- Two rules on chunk size: a multiple of the stream count, because each stream
  owns its own pedestal, and capped so the pinned buffers stay bounded.""")

# ---- A5 · ACT II REJECTED ROUTES ------------------------------------------
# Was a main-arc slide. It answers "did you try just making the copy faster?",
# which is a question, not a step in the argument, so it belongs here.
s = new_slide()
annex_chrome(s, 5, "rejected routes · the result copy · expands slide 18",
             "The copy is allocation-bound, not bandwidth-bound")
rows = [
    ("B\u2032", "One allocation per chunk", "collect_packed()",
     "Removes the per-frame malloc but keeps the copy, and the copy is ~80% of the "
     "cost. Worse, the replacement allocation is 1.17 GB, far above glibc's mmap "
     "threshold, so it is mmap'd and munmap'd every chunk: 606 566 faults, ~21 \u00b5s/frame.",
     "69.3 \u00b5s  \u00b7  deleted from the API"),
    ("B\u2033", "Parallel materialisation", "8-thread copy pool",
     "Each worker gets its own glibc arena, which destroys the cross-run heap reuse "
     "that makes the single-threaded path cheap. Faults went 9 700 \u2192 2 270 000; "
     "MALLOC_ARENA_MAX=1 collapsed them back to 138 k, which is the proof.",
     "+6% at best, \u221233% when results are freed promptly"),
]
y = 2.80
for tag, name, how, body, verdict in rows:
    rect(s, M, y, 12.30, 2.05, PANEL)
    rect(s, M, y, 0.04, 2.05, AMBER)
    tf = tb(s, M + 0.31, y + 0.22, 1.03, 0.4)
    run(para(tf, True), tag, 20, AMBER, bold=True, font=MONO)
    tf = tb(s, M + 1.34, y + 0.20, 6.20, 0.3)
    run(para(tf, True), name, 14, PALE, bold=True)
    tf = tb(s, M + 1.34, y + 0.52, 6.20, 0.3)
    run(para(tf, True), how, 10, MUTED, font=MONO)
    tf = tb(s, M + 1.34, y + 0.90, 10.54, 1.0)
    run(para(tf, True, line=1.3), body, 10, TEXT2)
    tf = tb(s, M + 1.34, y + 1.68, 10.54, 0.3)
    run(para(tf, True), verdict, 10.5, AMBER, bold=True)
    y += 2.25
callout(s, M, 1.82, 12.30,
        "**Copying faster does not help when the cost is the OS populating pages.** "
        "The only winning move is not to allocate, which is exactly what opt6 does. "
        "materialize_slot() is deliberately single-threaded and carries a comment "
        "saying so, to stop the experiment being repeated.", h=0.80, size=PT_LEAD)
notes(s, """SAY THIS -- 92 words, about 34 s. Annex: only on a question.

- The obvious answer to a slow copy is "copy faster". I tried both versions of
  that, and both are here because both failed.

- One allocation per chunk instead of per frame: the copy itself is 80 percent
  of the cost, so it barely helped -- and the big allocation went straight to
  mmap, which the OS then has to zero every chunk.

- Eight threads copying in parallel: worse. Each thread gets its own heap arena,
  which destroys the reuse that made the single-threaded version cheap. Faults
  went from ten thousand to two million.""")


# ---- A6 · THE BIG PINNED BUFFER ------------------------------------------
# This annex exists because the question is asked every time opt6's views are
# shown: "so give the finder one big pinned buffer at construction, let D2H land
# straight in it, and the clusters outlive the finder". It is a GOOD question --
# opt6 genuinely does not solve "I want every cluster at the end" -- so it must
# not be waved away. The answer is arithmetic, and it is a per-BYTE comparison,
# which is why it does not depend on the run length.
s = new_slide()
annex_chrome(s, 6, "the big pinned buffer · expands slide 18",
             "Why not pin one buffer and keep every cluster?")
bullets(s, M, 1.76, COL, [
    "**The proposal.** Give the finder a big **pinned** buffer; every frame's "
    "D2H lands in it at its own offset. No host-to-host copy, and the results "
    "outlive the finder. **It works, and it is faster than opt5.**",
    "**What binds is not time, it is RAM**, and the input has to stay pinned "
    "**at the same time**, since every frame's H2D still reads from it.",
], size=PT_BODY)
table(s, M, 3.58, COL,
      ["9×9 · pinned at once", "per frame", "N = 20 000", "N = 100 000"],
      [["input frames · opt4 already does this", "320 000 B", "6.4 GB",
        "32.0 GB"],
       ["output envelope · what this adds", "557 568 B", "11.2 GB", "55.8 GB"],
       ["**both, for the whole run**", "877 568 B", "**17.6 GB**",
        "**87.8 GB**"]],
      colw=[0.40, 0.20, 0.19, 0.21], size=PT_TABLE, rowh=0.44)
callout(s, M, 5.44, COL,
        "**98 GB free, no swap.** At the measured size (N = 20 000) it fits "
        "easily. It is **linear in N**: 877 568 B per frame caps a 9×9 run at "
        "**≈112 000 frames** with nothing left over.", h=0.72, size=PT_LEAD)
callout(s, M, 6.28, COL,
        "**Faster than opt5, and it costs 20 % more memory than opt5's own "
        "results.** The envelope stores the **cap** (1 700), not the **1 422** "
        "clusters actually found, and all of it is page-locked.",
        h=0.72, size=PT_LEAD, color=AMBER)
rail(s, [
    ("label", "9×9 · per frame"),
    ("gap", 0.10),
    ("stat", "opt5 · measured", "66.39 µs", TEXT2),
    ("row", "opt6 · no host copy either", "30.01 µs · 100 % of floor", ACCENT),
    ("gap", 0.16),
    ("label", "results kept, N = 100 000"),
    ("gap", 0.06),
    ("row", "opt5 · pageable, compacted", "46.6 GB", PALE),
    ("row", "this · pinned, padded to cap", "55.8 GB", AMBER),
    ("gap", 0.10),
    ("note", "Registration is in no timed number here: at 3×3 the whole "
             "cold–warm gap is 0.097 s, and covering the 32 GB input inside it "
             "would need 320 GB/s. Measured DRAM is 71 GB/s."),
])
notes(s, """SAY THIS -- 121 words, about 44 s. Annex: only on a question.

- Fair question, and opt6 really does not answer it. The views die with their
  chunk, so if you want every cluster at the end you still have to copy during
  the run.

- So: one big pinned buffer, the GPU writes straight into it. It works, and it
  IS faster than opt5 -- opt6 removes the same host copy and lands on thirty
  microseconds against opt5's sixty-six.

- What stops it is memory, not speed. The input has to stay pinned as well,
  because every frame is still copied from it. Together that is 878 kilobytes
  of locked memory per frame, so a hundred thousand frames needs 88 gigabytes
  of the 98 we have, and nothing may be swapped out.

DETAILS: The arithmetic. Everything below is measured or arithmetic on measured
numbers -- nothing here is an estimate.

MEASURED INPUTS
    frame                          400 x 400 x 2 B = 320 000 B      SS1
    9x9 D2H envelope, cap 1700     544.5 KiB = 557 568 B            A8
    9x9 clusters actually found    1422 x 328 B = 466 kB / frame    SS5
    host RAM                       125 GiB, no swap, 98 GB free     SS1/SS5
    host DRAM bandwidth            71 GB/s threaded                 SS1
    opt5 9x9                       66.39 us/frame                   SS8
    opt6 9x9                       30.01 us/frame, 100 % of floor   SS9
    opt5 3x3 cold / warm           26.95 / 25.98 us/frame, N=100k   SS5

1. SPEED. opt6 already removes exactly this host-to-host copy and measures
   30.01 against opt5's 66.39, so the proposal is x2.21 at 9x9. It is faster.
   The deck should say so rather than defend opt5.

2. WHY REGISTRATION DOES NOT SHOW UP IN ANY OF OUR NUMBERS. It is called once,
   outside the timed region, and that is provable rather than assumed: the
   entire 3x3 cold-warm gap is (26.95 - 25.98) x 100 000 = 0.097 s. Covering
   the 32 GB input inside it would need 320 GB/s. Measured DRAM is 71 GB/s, so
   it cannot be in there. NOTE THE CONSEQUENCE: opt4's x1.32 also excludes its
   own registration. Whatever accounting we accept for opt4 we must accept
   here.

3. WHAT ACTUALLY BINDS: BOTH BUFFERS ARE PINNED AT ONCE. The input cannot be
   released until the last frame, because every H2D reads from it.
       320 000 + 557 568 = 877 568 B per frame, page-locked
       N =  20 000  ->  17.6 GB      the size we actually measure: fine
       N = 100 000  ->  87.8 GB      of 98 GB free, no swap
   98e9 / 877 568 = 111 700 frames is the ceiling with zero headroom. This is
   why 9x9 is benchmarked at 20 000 frames and not 100 000 (SS5).

4. AND IT IS 20 % MORE MEMORY THAN opt5 KEEPS. opt5's results are compacted to
   what was found (466 kB/frame, 46.6 GB at 100k) and live in pageable heap.
   The envelope stores the cap: 557 568 / 466 000 = 1.20.

VERDICT, IF ASKED FOR ONE. For a bounded run it is the better design: faster
than opt5, and it removes the lifetime limit that opt6 has. It is not a general
replacement, because its footprint is O(N) in PAGE-LOCKED memory while opt5's
output is O(N) in ordinary heap. Pick it when N is known and bounded; pick
opt5/opt6 when it is not.""")


# ---- A7·1 · THE CONVENTION, DEFINED --------------------------------------
# Back in the annex, with its own grid behind it. Act III no longer opens on a
# definitions slide: the arc introduces s1 in one parenthesis on slide 19 and
# spends the stage time on the measurement that motivates the act instead. This
# is where the three words are actually defined, for the question that asks.
s = new_slide()
annex_chrome(s, 7, "measurement convention · expands slide 19",
             "Four kernels overlap, so the kernel floor is busy time, "
             "not run time", part=1, nparts=2)
# THE SLIDE HAS TO EARN ITS PLACE THIS LATE, so it now opens by saying why it is
# here rather than launching into definitions. The honest answer is specific:
# there is exactly ONE copy engine per direction, so an H2D or D2H bar is just
# how long the copy runs -- nothing to define. The kernel is the only engine that
# runs four at a time, and it becomes the tallest bar exactly here, in Act III.
# That is what makes "what is the floor" a real question rather than a caption.
bullets(s, M, 1.82, 12.30, [
    "**Why it needs saying.** In Act III the **kernel** is the tallest bar, "
    "and it is the one engine that runs **four at a time**. A copy engine's "
    "busy time is simply how long it runs; the kernel's is not, so what its "
    "floor means stops being obvious.",
], size=PT_BODY)
# Three definitions, one line of gloss each. Anything longer competes with the
# picture underneath, which is the thing that actually makes s4 make sense.
# CARD COLOURS MUST AGREE WITH THE FIGURE UNDERNEATH. They did not: s1 was blue
# while the blue bar in the picture is the KERNEL, whose 32.66 is an s4 number --
# so the card said "blue = s1" and the picture said "blue = an s4 value" six
# inches apart. s1 is now red (used nowhere in the figures as data), s4 is blue
# so it claims the bars it actually labels, and the floor is green to match the
# green floor line the figure already draws. A6's rail follows the same scheme.
cards = [
    ("s1", RED, "One stream, nothing else running",
     "How long an operation actually takes."),
    ("s4", ACCENT, "The shipped pipeline, four streams",
     "How busy each engine is: a union, not a sum."),
    ("floor", GREEN, "Set by the busiest engine",
     "1 / the busiest engine, in µs/frame or FPS."),
]
x = M
for tag, col, title, body in cards:
    rect(s, x, 2.50, 3.96, 1.24, PANEL)
    rect(s, x, 2.50, 3.96, 0.035, col)
    tf = tb(s, x + 0.26, 2.66, 3.41, 0.30)
    run(para(tf, True), tag, 15, col, bold=True, font=MONO)
    tf = tb(s, x + 0.26, 3.00, 3.50, 0.30)
    run(para(tf, True, line=1.1), title, 12, PALE, bold=True)
    tf = tb(s, x + 0.26, 3.30, 3.50, 0.40)
    run(para(tf, True, line=1.22), body, 11.5, TEXT2)
    x += 4.16
figure(s, "fig_measure", 0.87, 3.96, 11.60)
notes(s, """SAY THIS -- 120 words, about 44 s. Annex: only on a question.

- One question first: why measure engines only now? Because from here on the
  KERNEL is the slowest part. There is one copy engine in each direction, so a
  copy bar is just how long the copy runs. But four kernels run at the same
  time, so "how long the kernel takes" and "how much kernel time a frame costs"
  are two different numbers, and only one of them is the limit.

- From here on, every engine time has a label, and there are only two. "s1"
  means: one stream, nothing else running. That tells you how long an operation
  really takes. "s4" means: the real pipeline, four streams at once. That tells
  you how BUSY each engine is per frame.

- Why two? Because under load they are different numbers. The copies get slower,
  because they share one wire. The kernel gets cheaper per frame, because four
  kernels overlap each other.

- And the floor is the busiest engine in the real pipeline. Not the sum of the
  three -- the biggest one. If you only remember one line: s1 is how long it
  takes, s4 is how busy it is, and the floor is the busiest engine.

DETAILS: Say the floor out loud, in this order.

1. There are three engines and they are independent: the H2D copy engine, the
   SMs, and the D2H copy engine. PCIe is full duplex, so a frame's cost is
   never the sum of the three -- it is the slowest of them.

2. There is exactly ONE copy engine per direction. So two H2D copies can never
   run at the same time; they queue. That is not a modelling assumption, it is
   measured: H2D_overlap and D2H_overlap are 1.000 in every row of probes.csv.
   Kernels are the only row that ever overlaps.

3. That is why s4 is a UNION, not a sum. Four kernels, each 43.2 us long, but
   the SMs are busy only 32.66 us per frame: overlap factor 1.32x.

4. On units, because the card says "1 / the busiest engine" and the axis says
   microseconds: they are the same number. The busiest engine's busy time per
   frame is in us/frame; one divided by it is frames per second. 30.01 us/frame
   IS 33 323 FPS. The deck quotes whichever reads better in context -- us/frame
   when comparing engines, FPS when comparing against the CPU or the detector --
   and "% of floor" is the same ratio either way.

5. The floor is max(H2D, kernel, D2H) at s4. At 9x9 f64 that max is 32.66 us
   = 30 618 FPS. But the unprofiled pipeline actually sustained 30.01 us/frame,
   which is FASTER than the probe says is possible -- because under nsys
   submission is sparser, kernels overlap less, and a less-overlapped interval
   set has a LARGER union. A sustained rate is an existence proof; a probe is an
   estimate. So the floor is the lower of the two, and the 1.32x is a lower
   bound on the real overlap.

If time is short, say only: "s1 is how long it takes, s4 is how busy it is, and
the floor is the busiest engine." The rest is in A7.""")


# ---- A7·2 · THE FULL GRID ------------------------------------------------
s = new_slide()
annex_chrome(s, 7, "measurement convention · every engine number",
             "Uncontended, or as the pipeline runs it", part=2, nparts=2)
bullets(s, M, 1.82, 12.40, [
    "Every engine time in this deck is tagged **[build · s1|s4]**. **s1** is one "
    "stream with nothing else running: what an engine does **on its own**, which is "
    "the right number for a capability claim and for the headroom that remains. "
    "**s4** is the shipped four-stream pipeline: each engine's **busy time per "
    "frame**, the union of its intervals, the only number that can set a floor.",
], size=PT_BODY)
table(s, M, 3.07, 5.99,
      ["3×3 · µs/frame", "s1 f64", "s1 f32", "s4 f64", "s4 f32"],
      [["H2D", "13.14", "13.15", "16.17", "16.63"],
       ["kernel", "14.72", "4.32", "15.17", "5.53"],
       ["D2H", "5.31", "5.27", "7.69", "7.57"],
       ["engine max [s4]", "—", "—", "16.17", "16.63"],
       ["FLOOR = lower of max, sustained", "—", "—", "**16.17**", "**16.31**"]],
      colw=[0.28, 0.18, 0.18, 0.18, 0.18], size=PT_TABLE, rowh=0.44)
table(s, 7.01, 3.07, 5.99,
      ["9×9 · cap 1700", "s1 f64", "s1 f32", "s4 f64", "s4 f32"],
      [["H2D", "13.20", "13.22", "20.77", "20.54"],
       ["kernel", "39.86", "23.70", "32.66", "23.94"],
       ["D2H", "21.97", "21.95", "25.25", "25.24"],
       ["engine max [s4]", "—", "—", "32.66", "25.24"],
       ["FLOOR = lower of max, sustained", "—", "—", "**30.01**", "**25.14**"]],
      colw=[0.28, 0.18, 0.18, 0.18, 0.18], size=PT_TABLE, rowh=0.44)
callout(s, M, 6.00, 12.40,
        "**Two traps.** (1) s4 is engine *occupancy*, not duration: the 9×9 kernel "
        "row **falls** under load. (2) the **floor is not the engine max**: the max "
        "is profiled and runs 2–8 % high.", h=0.92, size=PT_LEAD)
caption(s, M, 7.00, 12.40,
        "The D2H shift is an s4 phenomenon: at s1 the kernel binds in both arms. "
        "Which engine binds must be read from the s4 columns.")
notes(s, """SAY THIS -- 96 words, about 35 s. Annex: only on a question.

- This is every engine number in the deck, in one grid, both cluster sizes and
  both builds.

- Left columns are s1, one stream, nothing else running: how long an operation
  takes. Right columns are s4, the real four-stream pipeline: how busy each
  engine is per frame. Under load the copies get slower and the kernel gets
  cheaper.

- Two traps at the bottom. The s4 kernel number is occupancy, not duration. And
  the floor is not simply the biggest s4 number: the profiler runs a few percent
  high, so we take the lower of the probe and the best real sustained rate.

DETAILS: Neither direction is faster than the other, and the bar heights say so.

At s1, H2D moves 320 000 B in 13.15 us = 24.3 GB/s; D2H moves 120 004 B in
5.27 us = 22.8 GB/s; at 9x9, D2H moves 557 604 B in 21.95 us = 25.4 GB/s. All
three sit at 72-81 % of PCIe 4.0 x16's 31.5 GB/s, the smallest transfer paying
the most fixed cost per byte. So a taller bar in this deck always means MORE
BYTES, never a slower wire -- which is exactly why raising the cap to 1 700
hands the 9x9 floor to D2H: it is 544.5 kB per frame either way the arm is
built, and it does not care about the typedef.

On the last row: the floor is the LOWER of the s4 estimate and the best
unprofiled sustained rate. 16.31 vs 16.63 at 3x3 f32; 25.14 vs 25.24 at 9x9 f32;
30.01 vs 32.66 at 9x9 f64, which is the widest gap and the reason the rule
exists. A probe roofline is an estimate, never a hard denominator.""")

# ---- A8 · OPT7, THE FULL ENGINE READING ----------------------------------
# Was main-arc slide 20. The arc keeps the two panels that carry the argument
# -- the s1 kernel pair that IS the -40.5 %, and the s4 grid that shows the
# handover -- and this slide keeps the rest: the s1 copy engines, the 3x3
# reading, and the s1-vs-s4 rule for which number to quote.
s = new_slide()
annex_chrome(s, 8, "opt7 · the full engine reading · expands slide 20",
             "FP32 halves pedestal traffic: −41 % kernel time")
bullets(s, M, 1.82, 7.75, [
    "**~80 % of pixels** take the **pedestal-update** branch: six accumulator "
    "values read and written.",
    "One typedef halves all four pedestal arrays (**48 bytes per updating "
    "pixel in FP64, 24 in FP32**), and the kernel is **bandwidth-bound**.",
], size=PT_BODY)
figure(s, "fig_f32_kernel", M, 3.20, 8.16)
# 6.60, not 5.72: the caption used to sit inside the chart's own bottom margin,
# which is the margin the taller bars now use.
caption(s, M, 6.60, 8.16,
        "Two more effects, neither the reason it works: FP64 arithmetic runs at "
        "1/64 of FP32 on a GeForce part, and the narrower accumulators free 9 "
        "registers at 3×3.", size=PT_META)
code(s, 8.80, 1.82, 4.10, [
    "// clusterfinder_kernel.cuh",
    "using COMPUTE_TYPE    = float;",
    "using DEVICE_PED_TYPE = «float»;",
    "//            was: double",
], size=PT_CODE, title="ONE TYPEDEF")
callout(s, 8.80, 3.22, 4.10,
        "Kernel, 9×9 **[s1 · cap 1700]**\n**39.86 → 23.70 µs  (−40.5 %)**",
        h=0.94, size=PT_LEAD)
callout(s, 8.80, 4.30, 4.10,
        "At 3×3 the same typedef is worth **−70.6 %** in the kernel and "
        "**4.6 %** end to end: there the kernel was never the tallest bar.",
        h=1.06, size=PT_LEAD)
callout(s, 8.80, 5.50, 4.10,
        "Naive FP32 is **wrong**. See **annex A9**.", h=0.56, size=PT_LEAD, color=AMBER)
caption(s, 8.80, 6.24, 4.10,
        "Both builds, same git rev, 20 000 frames. Full engine grid: A7.", size=PT_META)
notes(s, """SAY THIS -- 122 words, about 44 s. Annex: only on a question.

- Remember that about 80 percent of pixels see only noise, and those all update
  the pedestal. That update reads and writes six running values per pixel. So
  most of what the kernel does is move pedestal numbers, not compute photons.

- Those numbers were stored as doubles, 8 bytes each. Changing them to floats,
  4 bytes, halves that traffic: 48 bytes per pixel becomes 24. One typedef.
  The kernel is limited by memory traffic, so half the bytes is a big win.

- 40 percent off the kernel at 9x9. Two smaller effects come for free: on a
  gaming GPU double-precision maths runs at one sixty-fourth of single, and the
  smaller values also free up registers. Neither is the main reason.

DETAILS: Reading the two panels, and the number to quote.

Left is s1: one stream, nothing else running, so those are true durations. The
kernel binds in BOTH arms there -- 39.86 and 23.70 against a D2H of ~21.95 -- so
s1 alone would say the kernel is the thing to optimise and stop there.

Right is s4, the shipped four-stream pipeline, and it says something different:
the f64 kernel's busy time falls to 32.66 (self-overlap, 1.32x) while every
transfer RISES under contention, and once opt7 puts the kernel at 23.94 the D2H
bar at 25.24 is above it. That is the handover slide 21 is about.

So: -40.5 % is the s1 kernel claim and is the honest headline for what the
typedef does to the arithmetic. -26.7 % is the same change measured at s4, where
self-overlap had already hidden part of the win. Both are true, they measure
different things, and A7 is the rule for which to quote.

The 3x3 occupancy effect is real but not worth stage time: 47 registers -> 38
takes 3x3 from five blocks per SM to six, 83 % -> 100 % occupancy. It changes
nothing end to end, because at 3x3 H2D is the floor.""")


# ---- A9 · CATASTROPHIC CANCELLATION --------------------------------------
# Was main-arc slide 23. It is the trap that comes WITH opt7, not a step of the
# ladder, so the arc keeps the one-line warning on slide 20 and the trap, the
# physics it broke and the fix live here. Part 2 carries the ULP arithmetic of
# the error floor and the rewrite in full.
s = new_slide()
annex_chrome(s, 9, "opt7 · catastrophic cancellation · expands slide 20",
             "Accumulate what is small, not what is large",
             part=1, nparts=2)
bullets(s, M, 1.82, COL, [
    "**The trap.** **var = E[X²] − mean²** takes two numbers near 2 × 10⁷ to "
    "get one near 2 000.",
    "**What it cost.** The rms **clamps to zero**, so a 5σ gate becomes a "
    "**0σ gate**.",
    "**The fix.** Accumulate the **centred Y = X − X₀**: both operands become "
    "O(rms).",
], size=PT_BODY)
figure(s, "fig_cancellation", M, 3.80, COL)
caption(s, M, 6.34, COL,
        "Left: the operands and the answer against the ±3 ADU² error floor. "
        "Right: the f32 curve is reconstructed: measured area, modelled shape. "
        "The two-line patch is on the next slide.", size=PT_META)
rail(s, [
    ("label", "naive f32 · what it did"),
    ("gap", 0.12),
    ("stat", "Extra clusters", "+28.06 %", AMBER),
    ("row", "Pixels affected", "~1–2 % of the sensor", AMBER),
    ("row", "Error floor", "±3 ADU², absolute", AMBER),
    ("gap", 0.26),
    ("label", "after the rewrite"),
    ("gap", 0.12),
    ("row", "f32 vs f64 counts", "3 × 10⁻⁷", ACCENT),
    ("row", "vs the CPU baseline", "0.0039 %", ACCENT),
])
notes(s, """SAY THIS -- 128 words, about 47 s. Annex: only on a question.

- But the naive version of that change is WRONG, and this is the most useful
  slide of the talk if you ever write numerical code.

- The variance was computed as the average of X squared, minus the mean
  squared. Both of those are around twenty million. The answer is around two
  thousand. In float, twenty million has a step size of about two, so you
  subtract two huge numbers, and the small answer drowns in the rounding.

- What it did: the noise estimate collapsed to zero for one or two percent of
  pixels. A five-sigma threshold became a zero-sigma threshold, and we got 28
  percent more clusters -- all fake.

- The fix is two lines: subtract a fixed offset first, so we accumulate small
  numbers. Then float is fine.

DETAILS: The one sentence to leave the room with.

"Precision is relative." A float resolves small numbers finely and large numbers
coarsely, so a small answer must never be computed as the difference of two
large numbers. That is the transferable lesson; everything else on this slide is
this detector's instance of it.

Why the tail matters more than the count. +28 % is a number you might argue
about. A spectrum with a population sitting below the 5 sigma cut is not
arguable: those clusters cannot physically be there, and any gain or resolution
fit done on that spectrum is wrong.

Honesty about the right-hand panel: nobody kept the broken build around to
re-run, so the f32 curve is reconstructed. What is measured is the f64 curve and
the +28.06 % excess (CPU 116 010 113 vs CUDA 148 559 598, SS1 of the write-up).
The placement follows SS9: a correct 5 sigma gate cuts at ~225 ADU, a collapsed
gate admits the whole positive side of the pixel's distribution, so the excess
is smeared upward from zero. The write-up's TL;DR calls it a high-energy tail;
SS9 derives it from ~0 upward. Both describe the same corrupted pixels.

Welford's online variance is the other correct answer and is mentioned in the
write-up. The frozen-offset form was chosen because it is two lines and does not
change the update's arithmetic cost.""")


# ---- A9 · THE VARIANCE REWRITE IN FULL (part 2) --------------------------
# The quantitative backing for the third bullet of part 1: the rewrite itself
# and the error-floor panel that used to share fig_cancellation.
s = new_slide()
annex_chrome(s, 9, "the variance rewrite",
             "The rewrite in full, and which pixels the error reached",
             part=2, nparts=2)
bullets(s, M, 1.82, 7.65, [
    "Freeze **X₀ = round(mean)** once, at the end of pedestal training, and "
    "never move it.",
    "Accumulate the **centred** Y = X − X₀; report **X₀ + sum/n**. Both "
    "operands are now **O(rms)-sized**.",
], size=PT_BODY)
code(s, M, 3.20, 7.65, [
    "// before: both terms ~2.17e7, answer ~2000",
    "var = sum2/n - mean*mean;",
    "",
    "// after: centred on a frozen per-pixel offset X0",
    "DEVICE_PED_TYPE resid  = mean - «d_pd_off»[i];      // ~O(1)",
    "DEVICE_PED_TYPE var_px = sum2[i]/n - resid*resid;  // no cancellation",
], size=PT_CODE, title="clusterfinder_kernel.cuh")
callout(s, M, 5.02, 7.65,
        "Result: the 100 % FP32 build matches the FP64 build to **3 × 10⁻⁷**, "
        "70 clusters out of 233 million.", h=0.72, size=PT_LEAD)
callout(s, M, 5.88, 7.65,
        "**X₀ must never be updated**: the accumulators are defined relative to "
        "it. Welford's is the other correct answer; this one is two lines and "
        "costs nothing.", h=0.86, size=PT_LEAD, color=AMBER)
h = figure(s, "fig_varfloor", 8.35, 1.82, 4.46)
caption(s, 8.35, 1.82 + h + 0.18, 4.46,
        "Which pixels the ±3 ADU² floor actually reaches: the quiet ones, whose "
        "true variance is smallest.", size=PT_META)
notes(s, """SAY THIS -- 95 words, about 35 s. Annex: only on a question.

- The fix in full. At the end of pedestal training we freeze one number per
  pixel: the rounded mean. Call it X-zero. It never moves again.

- Then we accumulate the difference from X-zero instead of the raw value. Both
  sides of the subtraction are now small, about the size of the noise, so there
  is nothing left to cancel. Two lines.

- The result: the all-float build matches the double build to three parts in ten
  million. And the plot shows who was hurt before -- the quietest pixels, the
  ones with the smallest true noise.""")

# ---- A10 · THE FAULT MODEL -----------------------------------------------
s = new_slide()
annex_chrome(s, 10, "the fault model · expands slide 21",
             "The fault model, tested against every step")
bullets(s, M, 1.82, 12.40, [
    "The **0.68 µs per fault** rate is fitted on the **3×3 f32** ladder. Every "
    "row below is **9×9, f64 vs f32**, so the rate is applied **out of sample**, "
    "never refitted.",
], size=PT_BODY)
table(s, M, 2.55, 12.40,
      ["step", "f64 warm (faults)", "f32 warm (faults)", "Δ wall",
       "Δ faults", "predicted", "verdict"],
      [["opt3", "82.44 µs  (128 k)", "95.66 µs  (521 k)", "**+13.22**", "+393 k",
        "**+13.37**", "**the allocator**"],
       ["opt4", "79.83 µs  (128 k)", "75.20 µs  (127 k)", "−4.63", "−0.5 k",
        "−0.02", "clean"],
       ["opt5", "66.39 µs  (152 k)", "61.85 µs  (10 k)", "−4.54", "−141 k",
        "**−4.81**", "**not separable**"],
       ["opt6", "30.01 µs  (0)", "25.14 µs  (0)", "−4.87", "0", "0.00", "clean"]],
      colw=[0.08, 0.18, 0.18, 0.10, 0.10, 0.13, 0.23], size=PT_TABLE, rowh=0.50)
callout(s, M, 5.24, 12.40,
        "**opt3 is the whole argument in one row.** A 393 k fault gap predicts "
        "**+13.37 µs**; **+13.22** was observed. The −40 % kernel is in there, "
        "invisible under 13 µs of the OS zeroing pages.", h=0.9, size=PT_LEAD)
caption(s, M, 6.29, 12.40,
        "9×9 · cap 1 700 · warm = best of reps 1–4. The 0.68 µs/fault rate is "
        "carried in unchanged from the 3×3 fit: nothing here is tuned to make the "
        "columns agree.")
notes(s, """SAY THIS -- 94 words, about 34 s. Annex: only on a question.

- This is the page-fault model tested, not asserted. The cost per fault, 0.68
  microseconds, was fitted on the small cluster size. Every row here is the big
  one, so the number is applied out of sample and never refitted.

- Read the opt3 row. Its two arms differ by 393 thousand faults, the model
  predicts 13.4 microseconds, and 13.2 was measured. So that row tells you about
  the allocator, not about the typedef.

- That is why opt3 is excluded from the earlier chart and opt5 is marked with a
  dagger.""")

# ---- A11·1 · THE RESULT HEAP ---------------------------------------------
# Was main-arc slide 25. It is a warning about how to MEASURE, not a result,
# and the arc reads without it. The code panel that used to sit here is gone:
# A12·1, two slides later, carries the same four lines under the same title,
# and this slide's callout already states the protocol in words.
s = new_slide()
annex_chrome(s, 11, "behind the numbers · expands slide 24",
             "Materialising clusters costs 2.6 M page faults",
             part=1, nparts=2)
bullets(s, M, 1.82, COL, [
    "Every run that keeps its results **allocates a fresh heap and touches it "
    "once**. Linux only finds physical memory on **first touch**, so that pass "
    "is where the OS works, **inside the timer**.",
], size=PT_BODY)
figure(s, "fig_pagefault", M, 2.94, 8.16)
callout(s, M, 5.50, 8.16,
        "**Re-run until minor faults plateau (< 200 k), and quote that run.**",
        h=0.62, size=PT_LEAD)
rail(s, [
    ("label", "one 100 000-frame pass · 3×3"),
    ("gap", 0.12),
    ("stat", "Pages faulted in", "~2.6 M", AMBER),
    ("stat", "Cost inside the timer", "up to 4 s", AMBER),
    ("gap", 0.08),
    ("row", "Clusters materialised", "~10 GB", TEXT2),
    ("row", "Cost per fault, fitted", "0.68 µs", TEXT2),
    ("gap", 0.14),
    ("note", "Zero major faults all campaign. Copying faster cannot avoid this; "
             "not allocating can. That is opt6."),
])
notes(s, """SAY THIS -- 125 words, about 45 s. Annex: only on a question.

- A warning about measuring, and it cost me real time to find. When a run keeps
  its results, it allocates new memory and writes to it for the first time.
  Linux only gives you real physical memory when you first touch it, and that
  work lands INSIDE your timer.

- For one pass here that is 2.6 million page faults, and up to 4 seconds. That
  is not the GPU. That is not even our code. It is the operating system finding
  pages.

- So the rule we used: count the page faults around every timed block, run it
  again until that count settles, and quote that run. If you have never looked
  at this, your first benchmark is probably measuring your allocator.

DETAILS: THE NUMBERS ARE 3x3, opt2, 100 000 frames, one process, results
retained: 2 625 948 minor faults on run 1 at 6.110 s / 16 366 FPS, falling to
185 885 by run 3 at 4.452 s / 22 459 FPS. 2.6 M pages x 4 kB = ~10 GB, which is
100 000 x the 93 kB 3x3 payload. Section 12 of the report has the table.

THE FIT. Correlating delta-wall against delta-faults across those three runs
gives 0.65, 0.77 and 0.68 us/fault. Reconstructing run 1 from run 3:
4.452 s + 2 439 648 x 0.68 us = 6.111 s against 6.110 measured. Kernel time was
constant (0.022 ms) across all three, so the GPU is not involved at all.

WHY IT IS ZEROING, NOT LOOKUP. The kernel must hand out a frame that last
belonged to some other process, so it has to clear it first. That is mandatory
and it is the whole cost. At 4 kB a page, 1 GB of freshly-touched memory is
262 144 faults.

TWO INSTRUMENT CAVEATS, named and set aside. (1) Under nsys, host-side CUDA API
calls read ~4x high and wall time inflates with them, so every wall time in this
deck comes from an UNPROFILED run and only GPU-side timestamps come from nsys.
(2) CUDA events measure a stream, not a kernel. Both worked through in A12.

THE FAIRNESS TRAP that follows from this: freeing a ~10 GB result heap hands it
back to the allocator and the next loop reuses it, so on a cold heap the first
loop pays the entire tax and any printed ratio is meaningless. Both loops must
be at plateau; drop stale result bindings before timing.

Synthetic confirmation, same report section: allocate ~8 GB in ClusterVector-sized
chunks, touch every page, free, repeat in one process -- 2.05 s / 2 046 594
faults cold, then 0.07 s / ~2 000 warm. 30x faster, 1000x fewer faults, same
allocations. glibc retains the arenas.""")


# ---- A11·2 · WHAT A USER ACTUALLY GETS -----------------------------------
s = new_slide()
annex_chrome(s, 11, "behind the numbers · what a first run costs",
             "A first run loses a third of its throughput to page faults",
             part=2, nparts=2)
figure(s, "fig_first_run", 1.57, 1.72, 10.20)
callout(s, M, 5.96, 6.04,
        "**Everything that materialises clusters loses a third of its throughput on "
        "the first run**: +7 to +20 µs per frame.", h=1.00, size=PT_LEAD)
callout(s, 6.75, 5.96, 6.04,
        "**Only the two ends escape, for opposite reasons.** opt1 discards each "
        "frame; opt6 never grows a heap, and reaches **98 % of its peak on a cold "
        "process**.", h=1.00, size=PT_LEAD, color=AMBER)
caption(s, M, 7.06, 12.30,
        "Single pass, one process, every ClusterVector retained · f32 · 100 000 "
        "frames. Repeat the run and the amber bars climb onto the blue ones.", size=PT_META)
notes(s, """SAY THIS -- 114 words, about 42 s. Annex: only on a question.

- The previous slide said "run it twice". But a real user runs it ONCE, in a
  fresh process. So here is that number honestly: blue is the warm run, amber is
  the very first run.

- Almost every version loses about a third of its speed on the first run. That
  is the same page-fault cost, and it is real: your user pays it.

- Two versions escape, for opposite reasons. The first version throws every
  frame away, so it never grows a heap. And the last one never copies the
  results at all, so it does not either -- it reaches 98 percent of its warm
  speed on a cold start. Zero-copy paid twice.""")

# ---- A12·1 · THE THREE ARTEFACTS -----------------------------------------
# A11·1 draws only the first of these three at full size, because it is the only
# one that moves a number the audience is shown. This is the three-card version.
s = new_slide()
annex_chrome(s, 12, "benchmark artefacts · expands A11",
             "Three ways a GPU benchmark lies", part=1, nparts=4)
items = [
    ("First-touch page faults", AMBER,
     "Each run materialises ~10 GB of clusters. The first pass faults in ~2.6 M "
     "pages at 0.7 µs each, up to 4 s of pure OS work inside the timer.",
     "Fix: re-run until getrusage() minor faults plateau (< 200 k)."),
    ("CUDA-event kernel timing", AMBER,
     "avg_kernel_time_ms() measures elapsed time on a stream, including waiting "
     "for other streams. Under 8-stream load it over-reads by up to 3.5×.",
     "Fix: Nsight Systems per-instance times; 1 stream for exclusive numbers."),
    ("The profiler itself", AMBER,
     "Under nsys, wall time per frame inflates ~4× from API tracing.",
     "Fix: GPU op times from nsys, wall times from unprofiled runs."),
]
x = M
for title, col, body, fix in items:
    rect(s, x, 2.0, 3.96, 3.15, PANEL)
    rect(s, x, 2.0, 3.96, 0.035, col)
    tf = tb(s, x + 0.26, 2.28, 3.41, 0.6)
    run(para(tf, True, line=1.15), title, 13, PALE, bold=True)
    tf = tb(s, x + 0.26, 3.02, 3.41, 1.5)
    run(para(tf, True, line=1.3), body, 10, TEXT2)
    tf = tb(s, x + 0.26, 4.42, 3.41, 0.65)
    run(para(tf, True, line=1.25), fix, 10, ACCENT)
    x += 4.16
code(s, M, 5.32, 12.30, [
    "# bracket every timed cell:",
    "mf0 = resource.getrusage(resource.RUSAGE_SELF).ru_minflt",
    "...   t = time.perf_counter() - t0   ...",
    "print(f'minor faults: {mf1-mf0:,}')   # quote the run where this plateaus",
], size=PT_CODE, title="THE FAULT PROTOCOL · BRACKET EVERY TIMED CELL")
callout(s, M, 6.64, 12.30,
        "Validated: **wall = steady-state + faults × 0.68 µs** reproduced a 6.110 s "
        "run to within **1 ms**. Kernel time stayed constant throughout; the GPU was "
        "never the variable.", h=0.62, size=PT_LEAD)
notes(s, """SAY THIS -- 91 words, about 33 s. Annex: only on a question.

- Three ways a GPU benchmark lies, and I hit all three.

- One: the operating system finding pages, inside your timer. Two: CUDA events
  do not time a kernel, they time a stream -- so under load they include waiting
  behind other streams, and over-read by up to three and a half times. Three:
  the profiler itself makes host calls look about four times more expensive.

- The rule that comes out of it: GPU times from the profiler, wall times from
  runs with no profiler attached. Never mix them.""")

# ---- A12 · FAULTS --------------------------------------------------------
s = new_slide()
annex_chrome(s, 12, "benchmark artefacts · expands A11",
             "First-touch page faults: two sources, one counter", part=2, nparts=4)
bullets(s, M, 1.9, 12.40, [
    "A page exists in the address space but has no physical frame yet. First "
    "touch makes the kernel find one, **zero it** (mandatory), and map it. No "
    "disk: ru_majflt stays 0 all campaign. At 4 kB/page, **1 GB = 262 144 faults**.",
], size=PT_BODY)
table(s, M, 2.62, 12.40,
      ["", "(a) result heap", "(b) pinned D2H slots"],
      [["allocator", "malloc → mmap, one ClusterVector per frame",
        "cudaMallocHost, in submit_batch"],
       ["cost / page", "**0.7 µs**", "**1.0 µs**: same fault + pin + DMA map"],
       ["recurs?", "**yes**, every alloc/free cycle above the mmap threshold",
        "**no**: once per buffer, for its lifetime"],
       ["removed by", "**collect_view()**: it allocates nothing",
        "nothing; reserve_output_slots() only moves it out of the timer"]],
      colw=[0.14, 0.43, 0.43], size=PT_TABLE, rowh=0.50)
callout(s, M, 5.30, 6.04,
        "**The two are exactly additive.** Reserving subtracts precisely the pre-pin "
        "count from run 0 and changes nothing else.", h=0.72, size=PT_LEAD)
code(s, 6.75, 5.30, 6.04, [
    "3x3:   572 292 - 455 129 = 117 163   vs 117 192 pre-pin",
    "9x9: 2 759 037 - 2 278 567 = 480 470   vs 480 474",
    "closed form: 2 slots x 2000 x 120 004 B / 4 kB = 117 191",
], size=PT_CODE, title="NOT A CORRELATION, AN IDENTITY")
callout(s, M, 6.5, 12.30,
        "**At 9×9 the heap never plateaus.** ~9.3 GB per pass exceeds glibc's mmap "
        "threshold, so it is re-faulted every pass: **~10 µs/frame, permanently**.", h=0.72, size=PT_LEAD, color=AMBER)
notes(s, """SAY THIS -- 90 words, about 33 s. Annex: only on a question.

- The faults come from two places, and one counter sees both. The result heap,
  which we allocate and free every run, and the pinned buffers, which are
  faulted once and then stay.

- They are exactly additive -- reserving the pinned slots removes precisely the
  expected count and changes nothing else. That is an identity, not a
  correlation.

- One warning for the big cluster size: nine gigabytes per pass is above the
  threshold where the allocator gives memory back to the OS, so it never warms
  up. That cost is permanent.""")

# ---- A12 · EVENTS --------------------------------------------------------
s = new_slide()
annex_chrome(s, 12, "benchmark artefacts",
             "CUDA events measure the stream, not the kernel", part=3, nparts=4)
bullets(s, M, 1.9, COL, [
    "avg_kernel_time_ms() brackets the launch with **cudaEventRecord on the "
    "kernel's own stream**, so it returns time on **that stream's timeline**, "
    "including time queued behind other streams.",
    "Honest at 1 stream, inflated up to **3.5×** at 8. The tell: *PCIe + "
    "overhead* = wall/N − kernel_ms **goes negative**.",
    "It is not free either: **~3.6 µs/frame**, 10–15 % of throughput at 3×3.",
], size=PT_BODY)
code(s, M, 4.75, COL, [
    "if (m_time_kernels)                       // OFF by default, all 3 finders",
    "    cudaEventRecord(start[slot][i], sc.stream);",
    "find_clusters_in_single_frame<<<grid, block, shmem, «sc.stream»>>>(...);",
    "if (m_time_kernels)",
    "    cudaEventRecord(stop[slot][i], sc.stream);   // <- queue-wait lands here",
], size=PT_CODE, title="ClusterFinderCUDA.hpp")
callout(s, M, 6.18, COL,
        "The flag exists for **comparability**, not preference: with events on for one "
        "finder and off for another, the step between them absorbs the tax.",
        h=0.72, size=PT_LEAD)
rail(s, [
    ("label", "9×9 kernel · f64 · nsys"),
    ("gap", 0.10),
    # Same scheme as annex A8's cards: s1 red, s4 blue.
    ("stat", "s1 · duration", "39.86 µs", RED),
    ("stat", "s4 · occupancy", "32.66 µs", ACCENT),
    ("gap", 0.05),
    ("row", "overlap factor", "1.32×", TEXT2),
    ("gap", 0.20),
    ("note", "Same kernel. The s4 column is the union of kernel intervals per frame, "
             "not how long one kernel takes. Quote s1 for duration; s4 feeds the floor, "
             "subject to the sustained-rate rule in A7, which has the full grid."),
])
notes(s, """SAY THIS -- 88 words, about 32 s. Annex: only on a question.

- This is the trap that cost me the most. We time kernels with CUDA events. But
  an event is recorded ON A STREAM, so what you measure is when the stream got
  there -- including all the time the kernel sat waiting behind other streams.

- With one stream it is honest. With eight it reads up to three and a half times
  too high, and you can spot it: the leftover time goes negative.

- It is also not free. Turning it on costs about 3.6 microseconds per frame.""")

# ---- A12 · NSYS ----------------------------------------------------------
s = new_slide()
annex_chrome(s, 12, "benchmark artefacts",
             "Where nsys is sound, and where it is not", part=4, nparts=4)
bullets(s, M, 1.82, 12.40, [
    "Tracing a **host** call runs a callback on entry and exit, **inside the "
    "interval being measured**. GPU work is stamped by the hardware and read "
    "back **afterwards**, so nothing is injected into the execution path.",
], size=PT_BODY)
table(s, M, 2.55, 12.40,
      ["measurement", "sqlite table", "9×9 [f64 · s1 · cap 1500]", "verdict"],
      [["cudaLaunchKernel: the host call", "CUPTI_..._RUNTIME", "1.85 µs",
        "**inflated ~4×**"],
       ["the kernel executing", "CUPTI_..._KERNEL", "39.93 µs", "sound to ~2 %"],
       ["cudaMemcpyAsync: the host call", "CUPTI_..._RUNTIME", "1.65 µs",
        "**inflated ~4×**"],
       ["the H2D / D2H transfer", "CUPTI_..._MEMCPY", "13.25 / 19.44 µs",
        "sound to ~2 %"]],
      colw=[0.36, 0.24, 0.22, 0.18], size=PT_TABLE, rowh=0.52)
callout(s, M, 5.32, 12.30,
        "**Same cudaMemcpyAsync, two numbers: 1.65 µs to ask for the copy, 13.25 µs "
        "for the copy to happen.** A profiler that must be present to measure,"
        "distorts; one that reads a record later does not.", h=0.86, size=PT_LEAD)
caption(s, M, 6.33, 12.30,
        "Every engine number in this deck comes from GPU-side timestamps. Proof "
        "they are sound: opt7 sustains 25.14 µs unprofiled against a 25.24 µs"
        "profiled estimate, 0.4 % apart.")
notes(s, """SAY THIS -- 86 words, about 31 s. Annex: only on a question.

- The profiler is not wrong everywhere -- it is wrong in a specific place, and
  the table says which.

- To trace a HOST call, the tool runs code when you enter and when you leave.
  That code is inside the thing being measured, so host calls read about four
  times too high. GPU work is different: the hardware writes a timestamp, and
  the tool reads it later. Nothing is injected.

- Same call, two numbers: 1.6 microseconds to ASK for a copy, 13 for the copy
  to happen.

DETAILS: Why this one slide is still cap 1500, and what not to read off it.

The RUNTIME column exists only inside the trace: it is the cost of the host call
itself, which does not depend on how large the output slot is. So there is
nothing to re-run -- a cap-1700 trace would report the same ~1.65 and ~1.85 us
for cudaMemcpyAsync and cudaLaunchKernel.

What you must NOT read off this slide is the 19.44 us D2H. That is the cap-1500
figure. At the shipped cap of 1 700 the same engine reads 21.95 us [s1] and
25.24 us [s4], and it is the s4 value that binds the f32 build. A8 has the
shipped grid.

The one number in the whole deck taken from _RUNTIME is the CUDA Graph launch
budget in A2, and it is quoted to a single significant figure for exactly this
reason: a ~4x inflated host-call time can support "launch cost is about 2 us,
against a 24 us floor", and nothing finer than that.""")

prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
